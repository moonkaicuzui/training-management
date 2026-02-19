#!/usr/bin/env python3
"""
Q-TRAIN Training Guide — Korean (관리 모듈 업데이트 가이드)
8 Scenarios with detailed steps and data examples
"""

import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from pptx_helpers import (
    NAVY, BLUE, LIGHT_BLUE, WHITE, DARK, GRAY, LIGHT_GRAY, GREEN, ORANGE, RED, ACCENT,
    add_bg, add_shape, add_text, add_para, add_card, slide_header, add_table_slide,
    create_presentation, title_slide, thank_you_slide,
)

# ── Additional Colors ──
AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)
RED_BG = RGBColor(0xFE, 0xF2, 0xF2)
GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)
TEAL = RGBColor(0x06, 0x69, 0xB4)
PINK = RGBColor(0xDB, 0x27, 0x77)

prs, W, H = create_presentation()


# ════════════════════════════════════════════════════════════
# Helper: Scenario overview slide
# ════════════════════════════════════════════════════════════
def scenario_overview(prs, num, title, character_name, character_role, character_badge, badge_color,
                      situation, steps, step_colors=None):
    """Create a scenario overview slide with character avatar & numbered step flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, f"시나리오 {num}: {title}")

    # Character card (left)
    add_shape(slide, Inches(0.5), Inches(2.3), Inches(3.2), Inches(2.6), WHITE, radius=0.04)
    add_shape(slide, Inches(0.5), Inches(2.3), Inches(3.2), Inches(0.45), badge_color, radius=0.04)
    add_shape(slide, Inches(0.5), Inches(2.55), Inches(3.2), Inches(0.2), badge_color)

    # Badge
    b = add_shape(slide, Inches(0.65), Inches(2.38), Inches(1.1), Inches(0.28), WHITE, radius=0.1)
    tf = b.text_frame
    tf.paragraphs[0].text = character_badge
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = badge_color
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, Inches(0.65), Inches(2.9), Inches(2.9), Inches(0.3),
             character_name, size=15, color=NAVY, bold=True)
    add_text(slide, Inches(0.65), Inches(3.2), Inches(2.9), Inches(0.25),
             character_role, size=10, color=GRAY)

    add_shape(slide, Inches(0.65), Inches(3.55), Inches(2.9), Inches(0.02), LIGHT_BLUE)

    # Situation box
    txBox = add_text(slide, Inches(0.65), Inches(3.7), Inches(2.9), Inches(1.0),
                     "", size=10, color=DARK)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    add_para(tf, "상황:", size=10, color=badge_color, bold=True, space_before=Pt(2))
    add_para(tf, situation, size=10, color=DARK, space_before=Pt(4))

    # Step flow (right area)
    add_shape(slide, Inches(4.0), Inches(2.3), Inches(8.8), Inches(2.6), LIGHT_BLUE, radius=0.04)
    add_text(slide, Inches(4.2), Inches(2.4), Inches(8.4), Inches(0.3),
             "진행 단계", size=14, color=NAVY, bold=True)

    cols = min(len(steps), 4)
    rows_needed = (len(steps) + cols - 1) // cols

    for i, step_text in enumerate(steps):
        col = i % cols
        row = i // cols
        sx = Inches(4.3) + col * Inches(2.15)
        sy = Inches(2.85) + row * Inches(0.9)

        sc = step_colors[i] if step_colors and i < len(step_colors) else BLUE

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, sx, sy, Inches(0.3), Inches(0.3))
        circle.fill.solid()
        circle.fill.fore_color.rgb = sc
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.paragraphs[0].text = str(i + 1)
        ctf.paragraphs[0].font.size = Pt(10)
        ctf.paragraphs[0].font.color.rgb = WHITE
        ctf.paragraphs[0].font.bold = True
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_text(slide, sx + Inches(0.38), sy, Inches(1.7), Inches(0.35),
                 step_text, size=10, color=DARK)

        # Arrow to next step in same row
        if col < cols - 1 and i < len(steps) - 1:
            add_text(slide, sx + Inches(1.95), sy, Inches(0.2), Inches(0.3),
                     "→", size=14, color=GRAY, bold=True)

    return slide


def scenario_detail_table(prs, num, title, subtitle, rows_data, col_widths, note=""):
    """Create a scenario detail slide with a data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, f"시나리오 {num}: {title} — 상세", subtitle)

    add_table_slide(slide, Inches(0.5), Inches(2.3), Inches(12.3), rows_data, col_widths)

    if note:
        y_note = Inches(2.3) + Inches(0.4 * len(rows_data)) + Inches(0.2)
        add_text(slide, Inches(0.5), y_note, Inches(12), Inches(0.4),
                 note, size=10, color=GRAY)

    return slide


# ════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════
title_slide(prs,
            "관리 모듈 업데이트 가이드",
            "8가지 시나리오로 배우는 Q-TRAIN 관리 기능",
            "HWK Vietnam (화승비나) | QIP 교육 관리",
            "https://q-train-web.web.app",
            "2026.02 | Version 1.0")


# ════════════════════════════════════════════════════════════
# SLIDE 2: Table of Contents
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "목차", "8가지 시나리오 기반 가이드")

scenarios = [
    ("01", "신규 교육 프로그램 등록", "새로운 품질 검사 교육 프로그램 생성"),
    ("02", "교육 세션 스케줄링 & 결과 입력", "세션 생성, 참석자 등록, 결과 기록"),
    ("03", "불합격자 재교육 관리", "자동 감지 → 재교육 → 재평가 → 합격"),
    ("04", "TQC 신입교육 프로세스", "4단계 교육 → 면담 → 현장 평가"),
    ("05", "5PRS 교육 추천 시스템 활용", "불량 분석 → 추천 → 교육 → 결과 추적"),
    ("06", "CAPA 프로세스 연동", "문제 발생 → 5-Why → 시정조치 → 교육 연계"),
    ("07", "역량 매트릭스 관리", "역량 정의 → 갭 분석 → 학습 경로"),
    ("08", "대시보드 & 리포트 활용", "KPI 모니터링 → 분석 → 리포트 생성"),
]

y = Inches(2.3)
for i, (num, title, desc) in enumerate(scenarios):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + col * Inches(6.2)
    cy = y + row * Inches(0.85)

    colors = [BLUE, GREEN, RED, ORANGE, ACCENT, RGBColor(0xEF, 0x44, 0x44), TEAL, NAVY]
    sc = colors[i]

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, cy, Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = sc
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, x + Inches(0.55), cy - Inches(0.02), Inches(5), Inches(0.25),
             title, size=14, color=NAVY, bold=True)
    add_text(slide, x + Inches(0.55), cy + Inches(0.25), Inches(5), Inches(0.2),
             desc, size=10, color=GRAY)


# ════════════════════════════════════════════════════════════
# SCENARIO 1: 신규 교육 프로그램 등록
# ════════════════════════════════════════════════════════════
scenario_overview(prs, "1", "신규 교육 프로그램 등록",
    "문경서 차장", "QIP 관리자 (ADMIN)", "ADMIN", BLUE,
    "새로운 품질 검사 교육이 필요합니다.\n외관검사 기초 교육 프로그램을 시스템에 등록해야 합니다.",
    ["프로그램 코드 생성", "평가기준 설정", "유효기간 설정", "강사 배정"],
    [NAVY, BLUE, GREEN, ORANGE])

# Scenario 1 Detail Table
scenario_detail_table(prs, "1", "신규 교육 프로그램 등록", "프로그램 등록 예시 데이터",
    [
        ["항목", "설정 값", "설명"],
        ["프로그램 코드", "QIP-001", "고유 식별 코드 (자동 채번 가능)"],
        ["프로그램명 (EN)", "Visual Inspection Basic", "영문 교육명"],
        ["프로그램명 (VI)", "Kiem tra truc quan co ban", "베트남어 교육명"],
        ["프로그램명 (KR)", "외관검사 기초", "한국어 교육명"],
        ["카테고리", "QIP", "QIP / Production / Retraining / Newcomer"],
        ["대상 직급", "Worker, Line Leader", "교육 대상 직급"],
        ["평가 방식", "SCORE (점수제)", "점수제 또는 합격/불합격"],
        ["합격 점수", "70점", "이 점수 이상이면 PASS"],
        ["등급 기준", "AA>=95 / A>=85 / B>=70", "점수 구간별 자동 등급"],
        ["교육 시간", "4시간", "1회 교육 소요 시간"],
        ["유효기간", "12개월", "자격 만료까지 기간"],
    ],
    [Inches(2.2), Inches(3.5), Inches(6.6)],
    "* 프로그램 생성/수정/삭제 시 자동으로 program_change_logs에 변경 이력이 기록됩니다.")


# ════════════════════════════════════════════════════════════
# SCENARIO 2: 교육 세션 스케줄링 & 결과 입력
# ════════════════════════════════════════════════════════════
scenario_overview(prs, "2", "교육 세션 스케줄링 & 결과 입력",
    "Nguyen Thi Mai", "QA 팀장 / 사내 강사 (TRAINER)", "TRAINER", GREEN,
    "QIP-001 외관검사 기초 교육을 Building A 작업자 대상으로\n진행해야 합니다.",
    ["세션 생성", "참석자 등록", "출석 체크", "점수/등급 입력"],
    [BLUE, GREEN, ORANGE, ACCENT])

# Scenario 2 Detail: Attendance & Results
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "시나리오 2: 출석 & 결과 입력 — 상세", "교육 당일 데이터 기록 예시")

# Attendance table
add_text(slide, Inches(0.5), Inches(2.2), Inches(5.5), Inches(0.3),
         "출석 기록", size=14, color=NAVY, bold=True)

att_data = [
    ["사번", "이름", "출석", "체크인", "체크아웃"],
    ["HV-2024-0847", "Le Thi Hoa", "PRESENT", "08:50", "13:05"],
    ["HV-2023-0312", "Pham Van Thanh", "PRESENT", "08:55", "13:00"],
    ["HV-2022-0156", "Vo Thi Lan", "LATE", "09:20", "13:00"],
    ["HV-2024-0901", "Bui Quang Minh", "ABSENT", "-", "-"],
]
add_table_slide(slide, Inches(0.5), Inches(2.55), Inches(6), att_data,
                [Inches(1.4), Inches(1.5), Inches(1.1), Inches(1), Inches(1)])

# Result table
add_text(slide, Inches(6.8), Inches(2.2), Inches(6), Inches(0.3),
         "결과 입력", size=14, color=NAVY, bold=True)

res_data = [
    ["이름", "점수", "등급", "결과", "재교육"],
    ["Le Thi Hoa", "92", "A", "PASS", "NO"],
    ["Pham Van Thanh", "78", "B", "PASS", "NO"],
    ["Vo Thi Lan", "55", "C", "FAIL", "YES"],
    ["Bui Quang Minh", "-", "-", "ABSENT", "YES"],
]
add_table_slide(slide, Inches(6.8), Inches(2.55), Inches(6), res_data,
                [Inches(1.5), Inches(0.8), Inches(0.8), Inches(1.1), Inches(0.8)])

# Key rules
add_shape(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.2), AMBER_BG, radius=0.04)
add_text(slide, Inches(0.7), Inches(5.1), Inches(12), Inches(0.3),
         "핵심 규칙", size=14, color=RGBColor(0x92, 0x40, 0x0E), bold=True)

rules = [
    "점수 입력 시 등급 자동 계산:  AA (>=95) -> A (>=85) -> B (>=70) -> C (<70)",
    "합격 점수(70) 미만 = FAIL -> needs_retraining 자동 활성화",
    "ABSENT(결석)도 자동으로 재교육 대상으로 분류",
    "결과 수정 시 반드시 '수정 사유' 입력 필수 -> result_edit_logs에 영구 보존",
    "결과 삭제는 절대 불가 (No Delete Policy)",
]

txBox = add_text(slide, Inches(0.7), Inches(5.45), Inches(11.5), Inches(1.6), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for r in rules:
    add_para(tf, "  " + r, size=11, color=RGBColor(0x78, 0x35, 0x0F), space_before=Pt(5))


# ════════════════════════════════════════════════════════════
# SCENARIO 3: 불합격자 재교육 관리
# ════════════════════════════════════════════════════════════
scenario_overview(prs, "3", "불합격자 재교육 관리",
    "Nguyen Thi Mai", "QA 팀장 / 사내 강사 (TRAINER)", "TRAINER", GREEN,
    "시나리오 2에서 Vo Thi Lan이 불합격(55점),\nBui Quang Minh이 결석 처리되었습니다.\n이 두 명을 재교육해야 합니다.",
    ["불합격 자동 감지", "재교육 세션 생성", "재평가 실시", "합격 처리"],
    [RED, ORANGE, BLUE, GREEN])

# Scenario 3 Detail
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "시나리오 3: 재교육 관리 — 상세", "재교육 대상 자동 감지 및 추적 흐름")

# Left: Retraining targets
add_shape(slide, Inches(0.5), Inches(2.3), Inches(5.8), Inches(4.5), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(2.45), Inches(5), Inches(0.3),
         "재교육 대상 자동 식별 (/retraining)", size=14, color=NAVY, bold=True)

retrain_items = [
    ("FAILED — 시험 불합격자", RED, True),
    ("  Vo Thi Lan: 55점 (합격기준 70점 미달)", DARK, False),
    ("", DARK, False),
    ("ABSENT — 교육 결석자", RED, True),
    ("  Bui Quang Minh: 미출석", DARK, False),
    ("", DARK, False),
    ("EXPIRED — 자격 유효기간 만료", ORANGE, True),
    ("  (해당 없음 — 첫 교육)", GRAY, False),
    ("", DARK, False),
    ("EXPIRING — 30일 내 만료 예정", ORANGE, True),
    ("  (해당 없음)", GRAY, False),
]

txBox = add_text(slide, Inches(0.7), Inches(2.9), Inches(5.2), Inches(3.5), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item, color, bold in retrain_items:
    add_para(tf, "  " + item, size=11, color=color, bold=bold, space_before=Pt(3))

# Right: Retraining flow
add_shape(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(4.5), LIGHT_BLUE, radius=0.04)
add_text(slide, Inches(7), Inches(2.45), Inches(5.5), Inches(0.3),
         "재교육 프로세스 흐름", size=14, color=NAVY, bold=True)

flow_items = [
    "1. 시스템이 FAIL/ABSENT 자동 식별",
    "2. 관리자가 재교육 세션 생성 (같은 프로그램)",
    "3. 대상자에게 알림 발송 (RETRAINING_REQUIRED)",
    "4. 재교육 실시 + 출석 관리",
    "5. 재시험 + 결과 입력",
    "6. 합격 시 새 자격증 발급",
    "7. 전체 재교육 이력 추적 가능",
    "",
    "재교육 결과 예시:",
    "  Vo Thi Lan: 82점 (B등급) -> PASS",
    "  Bui Quang Minh: 75점 (B등급) -> PASS",
]

txBox = add_text(slide, Inches(7), Inches(2.9), Inches(5.5), Inches(3.5), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in flow_items:
    c = GREEN if "PASS" in item else DARK
    add_para(tf, "  " + item, size=11, color=c, space_before=Pt(4),
             bold=("PASS" in item))


# ════════════════════════════════════════════════════════════
# SCENARIO 4: TQC 신입교육 프로세스
# ════════════════════════════════════════════════════════════
scenario_overview(prs, "4", "TQC 신입교육 프로세스",
    "Le Thi Hoa", "신입 작업자 (입사 3일차)", "교육 대상자", ACCENT,
    "Le Thi Hoa가 신입사원으로 입사했습니다.\n4단계 TQC 교육 프로세스를 진행해야 합니다.",
    ["TQC 팀 배정", "4단계 교육 진행", "면담 기록", "현장 평가"],
    [NAVY, BLUE, GREEN, ACCENT])

# Scenario 4 Detail: TQC 4 stages
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "시나리오 4: TQC 4단계 교육 — 상세", "신입사원 온보딩 과정 (/new-tqc)")

stages = [
    ("Stage 1", "오리엔테이션\nOrientation", "공장 안전 교육\n사내 규정 안내\n기본 품질 개념", GREEN),
    ("Stage 2", "기초 교육\nBasic Training", "QIP 교육 수강\n시험 응시\n합격 기준 충족", BLUE),
    ("Stage 3", "라인 배치\nLine Assignment", "실제 라인에서 OJT\n선임자 지도하에\n실무 능력 습득", ORANGE),
    ("Stage 4", "현장 평가\nField Evaluation", "독립 작업 능력\n최종 합격 판정\n자격증 발급", ACCENT),
]

for i, (label, stage_title, desc, color) in enumerate(stages):
    x = Inches(0.5) + i * Inches(3.15)
    y = Inches(2.3)

    card = add_shape(slide, x, y, Inches(2.95), Inches(2.8), WHITE, radius=0.04)
    add_shape(slide, x, y, Inches(2.95), Inches(0.6), color, radius=0.04)
    add_shape(slide, x, y + Inches(0.4), Inches(2.95), Inches(0.2), color)

    add_text(slide, x + Inches(0.1), y + Inches(0.05), Inches(2.75), Inches(0.2),
             label, size=10, color=RGBColor(0xBF, 0xDB, 0xFE), bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(0.22), Inches(2.75), Inches(0.4),
             stage_title, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    txBox = add_text(slide, x + Inches(0.15), y + Inches(0.75), Inches(2.65), Inches(1.8), "", size=11, color=DARK)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for line in desc.split('\n'):
        add_para(tf, "  " + line, size=11, color=DARK, space_before=Pt(6))

    if i < 3:
        add_text(slide, x + Inches(2.95), y + Inches(1.2), Inches(0.2), Inches(0.3),
                 "->", size=18, color=GRAY, bold=True)

# Additional TQC features
add_shape(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8), LIGHT_BLUE, radius=0.04)
add_text(slide, Inches(0.7), Inches(5.4), Inches(11.5), Inches(0.3),
         "TQC 추가 기능", size=14, color=NAVY, bold=True)

tqc_features = [
    "정기 면담: 1주 / 1개월 / 3개월 면담 자동 스케줄링 (/new-tqc/meetings)",
    "색맹 검사: 입사 시 PASS/FAIL 기록 (품질 검사 직무 배정 기준)",
    "퇴사 추적: 퇴사 사유 8가지 카테고리 분류, 교육 기간 중 이직률 분석",
    "TQC 대시보드: 신입 인원, 완료율, 이직률, 평균 온보딩 기간 종합 모니터링",
]

txBox = add_text(slide, Inches(0.7), Inches(5.75), Inches(11.5), Inches(1.1), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for f in tqc_features:
    add_para(tf, "  " + f, size=10, color=DARK, space_before=Pt(4))


# ════════════════════════════════════════════════════════════
# SCENARIO 5: 5PRS 교육 추천 시스템 활용
# ════════════════════════════════════════════════════════════
scenario_overview(prs, "5", "5PRS 교육 추천 시스템 활용",
    "문경서 차장", "QIP 관리자 (ADMIN)", "ADMIN", BLUE,
    "5PRS 불량 분석 결과에 따라 교육이 필요한\n직원과 프로그램을 자동으로 추천받습니다.",
    ["분석 실행", "추천 목록 확인", "교육 등록", "결과 추적"],
    [NAVY, ACCENT, GREEN, BLUE])

# Scenario 5 Detail: 5PRS Data Example
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "시나리오 5: 5PRS 불량률 데이터 — 상세", "불량 분석 기반 교육 추천 예시")

prs_data = [
    ["동/라인", "불량 유형", "불량률", "기준", "상태", "추천 교육"],
    ["A동 L3", "외관 스크래치", "3.2%", "1.5%", "초과", "QIP-001 외관검사 기초"],
    ["A동 L3", "봉제 불량", "1.8%", "2.0%", "정상", "-"],
    ["B동 L1", "접착 불량", "2.5%", "1.0%", "초과", "QIP-003 접착 검사"],
    ["B동 L2", "색상 편차", "0.8%", "1.5%", "정상", "-"],
    ["C동 L1", "치수 불량", "4.1%", "2.0%", "초과", "QIP-005 치수 검사"],
]

add_table_slide(slide, Inches(0.5), Inches(2.3), Inches(12.3), prs_data,
                [Inches(1.5), Inches(2), Inches(1.2), Inches(1.2), Inches(1.2), Inches(5.2)])

# Recommendation flow
add_shape(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(2.0), GREEN_BG, radius=0.04)
add_text(slide, Inches(0.7), Inches(5.3), Inches(11.5), Inches(0.3),
         "5PRS 추천 시스템 작동 원리", size=13, color=GREEN, bold=True)

txBox = add_text(slide, Inches(0.7), Inches(5.7), Inches(11.5), Inches(1.3), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
add_para(tf, "  1. 5PRS 불량 데이터 수집 (동/라인/불량유형별)", size=10, color=DARK, space_before=Pt(4))
add_para(tf, "  2. 기준 대비 초과 항목 자동 식별", size=10, color=DARK, space_before=Pt(4))
add_para(tf, "  3. 해당 불량 유형에 매핑된 교육 프로그램 추천", size=10, color=DARK, space_before=Pt(4))
add_para(tf, "  4. 해당 라인 작업자 자동 선별 -> 교육 세션 생성 -> 결과 추적까지 연동", size=10, color=DARK, space_before=Pt(4))


# ════════════════════════════════════════════════════════════
# SCENARIO 6: CAPA 프로세스 연동
# ════════════════════════════════════════════════════════════
scenario_overview(prs, "6", "CAPA 프로세스 연동",
    "문경서 차장", "QIP 관리자 (ADMIN)", "ADMIN", BLUE,
    "A동 L3에서 외관검사 불량률이 기준을 초과했습니다.\nCAPA 프로세스를 통해 시정/예방 조치를 진행합니다.",
    ["CAPA 등록", "5-Why 분석", "시정조치 실행", "교육 연계", "검증 완료"],
    [RED, ORANGE, BLUE, GREEN, ACCENT])

# Scenario 6 Detail: CAPA 5 stages
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "시나리오 6: CAPA 5단계 프로세스 — 상세", "문제 발생 시 체계적 시정/예방 조치 (/capa)")

capa_stages = [
    ("1", "발견\nDiscovery", "문제 기술\n출처 분류\n즉시 조치", RED),
    ("2", "조사\nInvestigation", "근본 원인 분석\n5-Why 분석\n영향 평가", ORANGE),
    ("3", "조치\nAction", "시정 조치 계획\n예방 조치 계획\n담당자/기한", BLUE),
    ("4", "검증\nVerification", "효과성 검증\n모니터링\n재발 확인", ACCENT),
    ("5", "종료\nClosure", "교훈 기록\n지식 공유\n문서화 완료", GREEN),
]

for i, (num, stage_title, desc, color) in enumerate(capa_stages):
    x = Inches(0.3) + i * Inches(2.55)
    y = Inches(2.3)

    add_shape(slide, x, y, Inches(2.35), Inches(2.5), WHITE, radius=0.04)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85), y + Inches(0.12), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.paragraphs[0].text = num
    ctf.paragraphs[0].font.size = Pt(16)
    ctf.paragraphs[0].font.color.rgb = WHITE
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, x + Inches(0.1), y + Inches(0.75), Inches(2.15), Inches(0.45),
             stage_title, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)

    txBox = add_text(slide, x + Inches(0.1), y + Inches(1.25), Inches(2.15), Inches(1.0), "", size=10, color=DARK)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for line in desc.split('\n'):
        add_para(tf, line, size=10, color=DARK, space_before=Pt(4), align=PP_ALIGN.CENTER)

    if i < 4:
        add_text(slide, x + Inches(2.35), y + Inches(1.0), Inches(0.2), Inches(0.3),
                 "->", size=16, color=GRAY, bold=True)

# CAPA Example
add_shape(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.2), RED_BG, radius=0.04)
add_text(slide, Inches(0.7), Inches(5.1), Inches(11.5), Inches(0.3),
         "시뮬레이션 예시: CAPA-2026-003", size=14, color=RED, bold=True)

txBox = add_text(slide, Inches(0.7), Inches(5.45), Inches(11.5), Inches(1.6), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
add_para(tf, "  문제: Building A L3 외관검사 불량률 3.2% (기준 1.5% 초과)", size=10, color=DARK, space_before=Pt(3))
add_para(tf, "  원인: QIP-001 교육이 이론 위주 -> 실무 적용 능력 부족 (5-Why 분석 결과)", size=10, color=DARK, space_before=Pt(3))
add_para(tf, "  시정조치: 보충교육 + 미세 스크래치 판별 가이드 제작 + 커리큘럼에 실습 2시간 추가", size=10, color=DARK, space_before=Pt(3))
add_para(tf, "  교육 연계: QIP-001 재교육 세션 생성 -> 해당 라인 작업자 30명 등록", size=10, color=BLUE, bold=True, space_before=Pt(3))
add_para(tf, "  결과: 불량률 3.2% -> 1.1% (기준 이내) | 효과성 85/100 | CAPA -> 프로그램 개선 완료", size=10, color=GREEN, bold=True, space_before=Pt(3))


# ════════════════════════════════════════════════════════════
# SCENARIO 7: 역량 매트릭스 관리
# ════════════════════════════════════════════════════════════
scenario_overview(prs, "7", "역량 매트릭스 관리",
    "문경서 차장", "QIP 관리자 (ADMIN)", "ADMIN", BLUE,
    "부서별 역량 수준을 평가하고 갭 분석을 통해\n학습 경로를 생성합니다.",
    ["역량 정의", "직원 할당", "갭 분석", "학습 경로 생성"],
    [NAVY, BLUE, ORANGE, GREEN])

# Scenario 7 Detail: Competency Tables
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "시나리오 7: 역량 & 스킬 갭 분석 — 상세", "역량 평가 -> 갭 분석 -> 학습 경로 추천")

add_text(slide, Inches(0.5), Inches(2.2), Inches(6), Inches(0.3),
         "직원별 역량 평가 (/competency)", size=14, color=NAVY, bold=True)

comp_data = [
    ["역량", "현재 수준", "목표 수준", "갭"],
    ["외관 검사", "COMPETENT (3)", "PROFICIENT (4)", "-1"],
    ["봉제 품질", "BEGINNER (2)", "COMPETENT (3)", "-1"],
    ["안전 절차", "COMPETENT (3)", "COMPETENT (3)", "0"],
    ["팀 리더십", "NOVICE (1)", "NOVICE (1)", "0"],
]
add_table_slide(slide, Inches(0.5), Inches(2.55), Inches(6), comp_data,
                [Inches(1.5), Inches(1.5), Inches(1.5), Inches(0.8)])

add_text(slide, Inches(6.8), Inches(2.2), Inches(6), Inches(0.3),
         "팀/부서 스킬 갭 분석 (/skill-gap)", size=14, color=NAVY, bold=True)

gap_data = [
    ["역량", "팀 평균", "목표", "갭", "미달 인원"],
    ["외관 검사", "2.8", "4.0", "-1.2", "15명"],
    ["봉제 품질", "2.5", "3.0", "-0.5", "8명"],
    ["안전 절차", "3.2", "3.0", "+0.2", "0명"],
]
add_table_slide(slide, Inches(6.8), Inches(2.55), Inches(6), gap_data,
                [Inches(1.3), Inches(1), Inches(1), Inches(0.8), Inches(1)])

# Learning path
add_shape(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.2), LIGHT_BLUE, radius=0.04)
add_text(slide, Inches(0.7), Inches(5.1), Inches(11.5), Inches(0.3),
         "학습 경로 & 개발 계획", size=14, color=NAVY, bold=True)

lp_items = [
    "5가지 역량 수준: NOVICE (1) -> BEGINNER (2) -> COMPETENT (3) -> PROFICIENT (4) -> EXPERT (5)",
    "6가지 역량 카테고리: TECHNICAL / QUALITY / SAFETY / LEADERSHIP / COMMUNICATION / PROCESS",
    "학습 경로 유형: ONBOARDING (입사) / POSITION (직무) / PROMOTION (승진) / SPECIALIZATION (전문) / REMEDIAL (보충)",
    "개인별 개발 계획(IDP): 목표 역량, 할당된 학습 경로, 관리자/본인 코멘트, 진행률 추적",
    "역량-프로그램 매핑: 각 역량에 필요한 교육 프로그램 자동 추천",
]

txBox = add_text(slide, Inches(0.7), Inches(5.5), Inches(11.5), Inches(1.5), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in lp_items:
    add_para(tf, "  " + item, size=10, color=DARK, space_before=Pt(5))


# ════════════════════════════════════════════════════════════
# SCENARIO 8: 대시보드 & 리포트 활용
# ════════════════════════════════════════════════════════════
scenario_overview(prs, "8", "대시보드 & 리포트 활용",
    "Tran Van Duc", "생산 1동 라인장 (VIEWER)", "VIEWER", ORANGE,
    "Tran Van Duc 라인장이 자신의 팀 교육 현황을\n대시보드에서 확인하고 리포트를 생성합니다.",
    ["대시보드 확인", "교육 이수율/합격률", "부서별 분석", "리포트 생성"],
    [NAVY, BLUE, GREEN, ACCENT])

# Scenario 8 Detail: Dashboard & KPIs
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "시나리오 8: 대시보드 & 리포트 — 상세", "실시간 KPI 모니터링과 리포트 활용")

# KPI Cards
add_card(slide, Inches(0.5), Inches(2.3), Inches(2.9), Inches(1.3),
         "전체 직원", "1,200명", "Active employees", LIGHT_BLUE, BLUE, NAVY)
add_card(slide, Inches(3.6), Inches(2.3), Inches(2.9), Inches(1.3),
         "이번달 완료", "3건", "Monthly sessions", LIGHT_BLUE, BLUE, NAVY)
add_card(slide, Inches(6.7), Inches(2.3), Inches(2.9), Inches(1.3),
         "전체 이수율", "67.2%", "Target: 90%", LIGHT_BLUE, BLUE, NAVY)
add_card(slide, Inches(9.8), Inches(2.3), Inches(2.9), Inches(1.3),
         "재교육 대상", "45명", "Needs retraining", RED_BG, RED, RED)

# Dashboard + Executive
add_shape(slide, Inches(0.5), Inches(3.9), Inches(6), Inches(3.3), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(4.05), Inches(5.5), Inches(0.3),
         "메인 대시보드 (/dashboard)", size=14, color=NAVY, bold=True)

dash_items = [
    "등급 분포 파이차트 (AA/A/B/C)",
    "월별 교육 추세 (계획 vs 완료) 막대 그래프",
    "동별 재교육/만료 현황 차트",
    "30일 내 만료 예정 자격 목록",
    "다가오는 교육 세션 일정",
    "최근 재교육 대상자 목록",
]
txBox = add_text(slide, Inches(0.7), Inches(4.45), Inches(5.5), Inches(2.5), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in dash_items:
    add_para(tf, "  " + item, size=11, color=DARK, space_before=Pt(5))

add_shape(slide, Inches(6.8), Inches(3.9), Inches(6), Inches(3.3), WHITE, radius=0.04)
add_text(slide, Inches(7), Inches(4.05), Inches(5.5), Inches(0.3),
         "경영진 대시보드 & 리포트", size=14, color=NAVY, bold=True)

exec_data = [
    ["KPI", "현재", "목표", "달성률"],
    ["교육 이수율", "67.2%", "90%", "74.7%"],
    ["자격 보유율", "58.3%", "85%", "68.6%"],
    ["신입 이직률", "8.2%", "5%", "-"],
    ["교육 ROI", "187%", "200%", "93.5%"],
]
add_table_slide(slide, Inches(7), Inches(4.5), Inches(5.5), exec_data,
                [Inches(1.5), Inches(1), Inches(1), Inches(1)])

add_text(slide, Inches(7), Inches(6.7), Inches(5.5), Inches(0.4),
         "리포트: 부서별, 프로그램별, 직원별\n기간 필터 (1M/3M/6M/1Y) | CSV/PDF 내보내기",
         size=10, color=GRAY)


# ════════════════════════════════════════════════════════════
# SLIDE: 데이터 업데이트 체크리스트
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "데이터 업데이트 체크리스트", "관리자별 일/주/월 업무 체크리스트")

checklist_data = [
    ["주기", "담당", "작업", "시스템 위치", "비고"],
    ["매일", "TRAINER", "출석 체크 & 결과 입력", "/attendance, /results", "교육 당일 즉시"],
    ["매일", "ADMIN", "대시보드 확인 (이상 징후)", "/dashboard", "KPI 모니터링"],
    ["매주", "TRAINER", "재교육 대상자 확인 & 세션 배정", "/retraining", "FAIL/ABSENT 처리"],
    ["매주", "ADMIN", "프로그램 변경 이력 리뷰", "/programs", "change_logs 확인"],
    ["매주", "TRAINER", "TQC 신입 교육 면담 실시", "/new-tqc/meetings", "1주 면담"],
    ["매월", "ADMIN", "경영진 대시보드 리뷰", "/executive", "KPI 리포트"],
    ["매월", "ADMIN", "만료 예정 자격 확인 & 재교육 계획", "/retraining", "30일 이내 만료"],
    ["매월", "ADMIN", "CAPA 현황 점검", "/capa", "진행 중 건 확인"],
    ["매월", "ADMIN", "역량 갭 분석 & 학습 경로 업데이트", "/skill-gap", "분기별 집중 분석"],
    ["매월", "ADMIN", "리포트 생성 & 공유", "/reports", "CSV/PDF 내보내기"],
]

add_table_slide(slide, Inches(0.5), Inches(2.3), Inches(12.3), checklist_data,
                [Inches(0.9), Inches(1.2), Inches(3.5), Inches(2.5), Inches(4.2)])


# ════════════════════════════════════════════════════════════
# SLIDE: FAQ / 문의처
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "FAQ / 문의처", "자주 묻는 질문과 지원 연락처")

# FAQ section
add_shape(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.0), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(2.4), Inches(11.5), Inches(0.3),
         "자주 묻는 질문 (FAQ)", size=16, color=NAVY, bold=True)

faqs = [
    ("Q: 교육 결과를 잘못 입력했으면 어떻게 하나요?",
     "A: 결과 수정은 가능하지만 반드시 수정 사유를 입력해야 합니다. 모든 수정 이력은 result_edit_logs에 영구 보존됩니다."),
    ("Q: 결과를 삭제할 수 있나요?",
     "A: 아니요. No Delete Policy에 따라 교육 결과는 절대 삭제할 수 없습니다. 감사 추적을 위한 핵심 원칙입니다."),
    ("Q: 자격증 유효기간이 만료되면 어떻게 되나요?",
     "A: 시스템이 자동으로 EXPIRING(30일 전)/EXPIRED 상태로 변경하고, 해당 직원과 관리자에게 알림을 발송합니다."),
    ("Q: 5PRS 추천 데이터는 어디서 가져오나요?",
     "A: 현재는 관리자가 불량 데이터를 입력합니다. 향후 MES 시스템과 자동 연동을 계획하고 있습니다."),
    ("Q: 여러 명의 결과를 한번에 입력할 수 있나요?",
     "A: 네. 세션별 일괄 입력이 가능합니다. 또한 엑셀 템플릿을 통한 대량 업로드도 지원합니다."),
]

y_pos = Inches(2.8)
for q, a in faqs:
    add_text(slide, Inches(0.7), y_pos, Inches(11.5), Inches(0.25),
             q, size=11, color=NAVY, bold=True)
    add_text(slide, Inches(0.7), y_pos + Inches(0.25), Inches(11.5), Inches(0.25),
             a, size=10, color=DARK)
    y_pos += Inches(0.6)

# Contact section
add_shape(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.8), LIGHT_BLUE, radius=0.04)
add_text(slide, Inches(0.7), Inches(6.55), Inches(11.5), Inches(0.3),
         "문의처", size=14, color=NAVY, bold=True)
add_text(slide, Inches(0.7), Inches(6.85), Inches(11.5), Inches(0.3),
         "시스템 관련 문의: ksmoon@hsvina.com  |  Q-TRAIN 시스템: https://q-train-web.web.app",
         size=11, color=DARK)


# ════════════════════════════════════════════════════════════
# SLIDE: Thank You
# ════════════════════════════════════════════════════════════
thank_you_slide(prs, "https://q-train-web.web.app", "ksmoon@hsvina.com")


# ════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "Q-TRAIN_Training_Guide_KO.pptx")
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
