#!/usr/bin/env python3
"""
Q-TRAIN Training Guide — English (Admin Module Update Guide)
8 Scenarios with detailed steps and data examples
"""

import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from pptx_helpers import (
    NAVY, BLUE, LIGHT_BLUE, WHITE, DARK, GRAY, LIGHT_GRAY, GREEN, ORANGE, RED, ACCENT,
    add_bg, add_shape, add_text, add_para, add_card, slide_header, add_table_slide,
    create_presentation, title_slide, thank_you_slide,
)

# ── Additional Colors ──
AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)
RED_BG = RGBColor(0xFE, 0xF2, 0xF2)
GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)
TEAL = RGBColor(0x06, 0x69, 0xB4)
PINK = RGBColor(0xDB, 0x27, 0x77)

prs, W, H = create_presentation()


# ════════════════════════════════════════════════════════════
# Helper: Scenario overview slide
# ════════════════════════════════════════════════════════════
def scenario_overview(prs, num, title, character_name, character_role, character_badge, badge_color,
                      situation, steps, step_colors=None):
    """Create a scenario overview slide with character avatar & numbered step flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, f"Scenario {num}: {title}")

    # Character card (left)
    add_shape(slide, Inches(0.5), Inches(2.3), Inches(3.2), Inches(2.6), WHITE, radius=0.04)
    add_shape(slide, Inches(0.5), Inches(2.3), Inches(3.2), Inches(0.45), badge_color, radius=0.04)
    add_shape(slide, Inches(0.5), Inches(2.55), Inches(3.2), Inches(0.2), badge_color)

    # Badge
    b = add_shape(slide, Inches(0.65), Inches(2.38), Inches(1.1), Inches(0.28), WHITE, radius=0.1)
    tf = b.text_frame
    tf.paragraphs[0].text = character_badge
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = badge_color
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, Inches(0.65), Inches(2.9), Inches(2.9), Inches(0.3),
             character_name, size=15, color=NAVY, bold=True)
    add_text(slide, Inches(0.65), Inches(3.2), Inches(2.9), Inches(0.25),
             character_role, size=10, color=GRAY)

    add_shape(slide, Inches(0.65), Inches(3.55), Inches(2.9), Inches(0.02), LIGHT_BLUE)

    # Situation box
    txBox = add_text(slide, Inches(0.65), Inches(3.7), Inches(2.9), Inches(1.0),
                     "", size=10, color=DARK)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    add_para(tf, "Situation:", size=10, color=badge_color, bold=True, space_before=Pt(2))
    add_para(tf, situation, size=10, color=DARK, space_before=Pt(4))

    # Step flow (right area)
    add_shape(slide, Inches(4.0), Inches(2.3), Inches(8.8), Inches(2.6), LIGHT_BLUE, radius=0.04)
    add_text(slide, Inches(4.2), Inches(2.4), Inches(8.4), Inches(0.3),
             "Steps", size=14, color=NAVY, bold=True)

    cols = min(len(steps), 4)
    rows_needed = (len(steps) + cols - 1) // cols

    for i, step_text in enumerate(steps):
        col = i % cols
        row = i // cols
        sx = Inches(4.3) + col * Inches(2.15)
        sy = Inches(2.85) + row * Inches(0.9)

        sc = step_colors[i] if step_colors and i < len(step_colors) else BLUE

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, sx, sy, Inches(0.3), Inches(0.3))
        circle.fill.solid()
        circle.fill.fore_color.rgb = sc
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.paragraphs[0].text = str(i + 1)
        ctf.paragraphs[0].font.size = Pt(10)
        ctf.paragraphs[0].font.color.rgb = WHITE
        ctf.paragraphs[0].font.bold = True
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_text(slide, sx + Inches(0.38), sy, Inches(1.7), Inches(0.35),
                 step_text, size=10, color=DARK)

        if col < cols - 1 and i < len(steps) - 1:
            add_text(slide, sx + Inches(1.95), sy, Inches(0.2), Inches(0.3),
                     "->", size=14, color=GRAY, bold=True)

    return slide


def scenario_detail_table(prs, num, title, subtitle, rows_data, col_widths, note=""):
    """Create a scenario detail slide with a data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, f"Scenario {num}: {title} — Details", subtitle)

    add_table_slide(slide, Inches(0.5), Inches(2.3), Inches(12.3), rows_data, col_widths)

    if note:
        y_note = Inches(2.3) + Inches(0.4 * len(rows_data)) + Inches(0.2)
        add_text(slide, Inches(0.5), y_note, Inches(12), Inches(0.4),
                 note, size=10, color=GRAY)

    return slide


# ════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════
title_slide(prs,
            "Admin Module Update Guide",
            "Learn Q-TRAIN management features through 8 scenarios",
            "HWK Vietnam | QIP Training Management",
            "https://q-train-web.web.app",
            "February 2026 | Version 1.0")


# ════════════════════════════════════════════════════════════
# SLIDE 2: Table of Contents
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Table of Contents", "8 Scenario-Based Guide Chapters")

scenarios = [
    ("01", "New Training Program Registration", "Create a new quality inspection training program"),
    ("02", "Session Scheduling & Result Entry", "Create sessions, register attendees, record results"),
    ("03", "Failed Trainee Retraining Management", "Auto-detect -> Retrain -> Re-evaluate -> Pass"),
    ("04", "TQC New Employee Onboarding", "4-stage training -> Interviews -> Field evaluation"),
    ("05", "5PRS Training Recommendation System", "Defect analysis -> Recommend -> Train -> Track results"),
    ("06", "CAPA Process Integration", "Problem -> 5-Why -> Corrective action -> Training link"),
    ("07", "Competency Matrix Management", "Define competencies -> Gap analysis -> Learning paths"),
    ("08", "Dashboard & Report Utilization", "KPI monitoring -> Analysis -> Report generation"),
]

y = Inches(2.3)
for i, (num, title, desc) in enumerate(scenarios):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + col * Inches(6.2)
    cy = y + row * Inches(0.85)

    colors = [BLUE, GREEN, RED, ORANGE, ACCENT, RGBColor(0xEF, 0x44, 0x44), TEAL, NAVY]
    sc = colors[i]

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, cy, Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = sc
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, x + Inches(0.55), cy - Inches(0.02), Inches(5), Inches(0.25),
             title, size=14, color=NAVY, bold=True)
    add_text(slide, x + Inches(0.55), cy + Inches(0.25), Inches(5), Inches(0.2),
             desc, size=10, color=GRAY)


# ════════════════════════════════════════════════════════════
# SCENARIO 1: New Training Program Registration
# ════════════════════════════════════════════════════════════
scenario_overview(prs, "1", "New Training Program Registration",
    "Mr. Moon", "QIP Manager (ADMIN)", "ADMIN", BLUE,
    "A new quality inspection training is needed.\nThe Visual Inspection Basic program must be registered.",
    ["Create program code", "Set evaluation criteria", "Set validity period", "Assign trainer"],
    [NAVY, BLUE, GREEN, ORANGE])

scenario_detail_table(prs, "1", "New Training Program Registration", "Program registration example data",
    [
        ["Field", "Value", "Description"],
        ["Program Code", "QIP-001", "Unique identifier (auto-generated)"],
        ["Name (EN)", "Visual Inspection Basic", "English program name"],
        ["Name (VI)", "Kiem tra truc quan co ban", "Vietnamese program name"],
        ["Name (KR)", "Visual Inspection Basic (KR)", "Korean program name"],
        ["Category", "QIP", "QIP / Production / Retraining / Newcomer"],
        ["Target Positions", "Worker, Line Leader", "Eligible job positions"],
        ["Evaluation Type", "SCORE (Score-based)", "Score-based or Pass/Fail"],
        ["Passing Score", "70 points", "Minimum score for PASS"],
        ["Grade Thresholds", "AA>=95 / A>=85 / B>=70", "Auto-calculated grade bands"],
        ["Duration", "4 hours", "Training session length"],
        ["Validity", "12 months", "Certificate expiration period"],
    ],
    [Inches(2.2), Inches(3.5), Inches(6.6)],
    "* All program changes are automatically logged in program_change_logs with before/after data.")


# ════════════════════════════════════════════════════════════
# SCENARIO 2: Session Scheduling & Result Entry
# ════════════════════════════════════════════════════════════
scenario_overview(prs, "2", "Session Scheduling & Result Entry",
    "Ms. Mai", "QA Team Lead / Trainer (TRAINER)", "TRAINER", GREEN,
    "QIP-001 Visual Inspection Basic training needs to be\nconducted for Building A workers.",
    ["Create session", "Register attendees", "Record attendance", "Enter scores/grades"],
    [BLUE, GREEN, ORANGE, ACCENT])

# Scenario 2 Detail: Attendance & Results
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Scenario 2: Attendance & Results — Details", "Training day data recording examples")

add_text(slide, Inches(0.5), Inches(2.2), Inches(5.5), Inches(0.3),
         "Attendance Records", size=14, color=NAVY, bold=True)

att_data = [
    ["Employee ID", "Name", "Status", "Check-In", "Check-Out"],
    ["HV-2024-0847", "Le Thi Hoa", "PRESENT", "08:50", "13:05"],
    ["HV-2023-0312", "Pham Van Thanh", "PRESENT", "08:55", "13:00"],
    ["HV-2022-0156", "Vo Thi Lan", "LATE", "09:20", "13:00"],
    ["HV-2024-0901", "Bui Quang Minh", "ABSENT", "-", "-"],
]
add_table_slide(slide, Inches(0.5), Inches(2.55), Inches(6), att_data,
                [Inches(1.4), Inches(1.5), Inches(1.1), Inches(1), Inches(1)])

add_text(slide, Inches(6.8), Inches(2.2), Inches(6), Inches(0.3),
         "Result Entry", size=14, color=NAVY, bold=True)

res_data = [
    ["Name", "Score", "Grade", "Result", "Retrain"],
    ["Le Thi Hoa", "92", "A", "PASS", "NO"],
    ["Pham Van Thanh", "78", "B", "PASS", "NO"],
    ["Vo Thi Lan", "55", "C", "FAIL", "YES"],
    ["Bui Quang Minh", "-", "-", "ABSENT", "YES"],
]
add_table_slide(slide, Inches(6.8), Inches(2.55), Inches(6), res_data,
                [Inches(1.5), Inches(0.8), Inches(0.8), Inches(1.1), Inches(0.8)])

# Key rules
add_shape(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.2), AMBER_BG, radius=0.04)
add_text(slide, Inches(0.7), Inches(5.1), Inches(12), Inches(0.3),
         "Key Business Rules", size=14, color=RGBColor(0x92, 0x40, 0x0E), bold=True)

rules = [
    "Auto-grade calculation: AA (>=95) -> A (>=85) -> B (>=70) -> C (<70)",
    "Score below passing (70) = FAIL -> needs_retraining auto-activated",
    "ABSENT also automatically flagged for retraining",
    "Result edits require mandatory reason -> permanently stored in result_edit_logs",
    "Results can NEVER be deleted (No Delete Policy)",
]

txBox = add_text(slide, Inches(0.7), Inches(5.45), Inches(11.5), Inches(1.6), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for r in rules:
    add_para(tf, "  " + r, size=11, color=RGBColor(0x78, 0x35, 0x0F), space_before=Pt(5))


# ════════════════════════════════════════════════════════════
# SCENARIO 3: Failed Trainee Retraining Management
# ════════════════════════════════════════════════════════════
scenario_overview(prs, "3", "Failed Trainee Retraining Management",
    "Ms. Mai", "QA Team Lead / Trainer (TRAINER)", "TRAINER", GREEN,
    "From Scenario 2, Vo Thi Lan failed (55 pts) and\nBui Quang Minh was absent.\nBoth need retraining.",
    ["Auto-detect failures", "Create retraining session", "Re-evaluate", "Pass processing"],
    [RED, ORANGE, BLUE, GREEN])

# Scenario 3 Detail
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Scenario 3: Retraining Management — Details", "Automatic detection and tracking flow")

# Left: Retraining targets
add_shape(slide, Inches(0.5), Inches(2.3), Inches(5.8), Inches(4.5), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(2.45), Inches(5), Inches(0.3),
         "Auto-Identified Retraining Targets (/retraining)", size=14, color=NAVY, bold=True)

retrain_items = [
    ("FAILED — Exam failures", RED, True),
    ("  Vo Thi Lan: 55 pts (passing: 70 pts)", DARK, False),
    ("", DARK, False),
    ("ABSENT — Training absences", RED, True),
    ("  Bui Quang Minh: Did not attend", DARK, False),
    ("", DARK, False),
    ("EXPIRED — Validity period expired", ORANGE, True),
    ("  (N/A — First training)", GRAY, False),
    ("", DARK, False),
    ("EXPIRING — Expiring within 30 days", ORANGE, True),
    ("  (N/A)", GRAY, False),
]

txBox = add_text(slide, Inches(0.7), Inches(2.9), Inches(5.2), Inches(3.5), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item, color, bold in retrain_items:
    add_para(tf, "  " + item, size=11, color=color, bold=bold, space_before=Pt(3))

# Right: Retraining flow
add_shape(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(4.5), LIGHT_BLUE, radius=0.04)
add_text(slide, Inches(7), Inches(2.45), Inches(5.5), Inches(0.3),
         "Retraining Process Flow", size=14, color=NAVY, bold=True)

flow_items = [
    "1. System auto-identifies FAIL/ABSENT",
    "2. Admin creates retraining session (same program)",
    "3. Notification sent (RETRAINING_REQUIRED)",
    "4. Conduct retraining + attendance tracking",
    "5. Re-exam + result entry",
    "6. On pass: new certificate issued",
    "7. Full retraining history tracked",
    "",
    "Retraining Results Example:",
    "  Vo Thi Lan: 82 pts (Grade B) -> PASS",
    "  Bui Quang Minh: 75 pts (Grade B) -> PASS",
]

txBox = add_text(slide, Inches(7), Inches(2.9), Inches(5.5), Inches(3.5), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in flow_items:
    c = GREEN if "PASS" in item else DARK
    add_para(tf, "  " + item, size=11, color=c, space_before=Pt(4),
             bold=("PASS" in item))


# ════════════════════════════════════════════════════════════
# SCENARIO 4: TQC New Employee Onboarding
# ════════════════════════════════════════════════════════════
scenario_overview(prs, "4", "TQC New Employee Onboarding",
    "Ms. Hoa", "New Worker (Day 3)", "TRAINEE", ACCENT,
    "Ms. Hoa has joined as a new employee.\nShe must complete the 4-stage TQC onboarding process.",
    ["Assign TQC team", "4-stage training", "Interview records", "Field evaluation"],
    [NAVY, BLUE, GREEN, ACCENT])

# Scenario 4 Detail: TQC 4 stages
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Scenario 4: TQC 4-Stage Training — Details", "New employee onboarding process (/new-tqc)")

stages = [
    ("Stage 1", "Orientation", "Factory safety rules\nCompany policies\nBasic quality concepts", GREEN),
    ("Stage 2", "Basic Training", "QIP program courses\nExam & scoring\nMust pass all required", BLUE),
    ("Stage 3", "Line Assignment", "On-the-job training\nSenior worker mentoring\nHands-on skill building", ORANGE),
    ("Stage 4", "Field Evaluation", "Independent work test\nFinal qualification\nCertificate issuance", ACCENT),
]

for i, (label, stage_title, desc, color) in enumerate(stages):
    x = Inches(0.5) + i * Inches(3.15)
    y = Inches(2.3)

    card = add_shape(slide, x, y, Inches(2.95), Inches(2.8), WHITE, radius=0.04)
    add_shape(slide, x, y, Inches(2.95), Inches(0.6), color, radius=0.04)
    add_shape(slide, x, y + Inches(0.4), Inches(2.95), Inches(0.2), color)

    add_text(slide, x + Inches(0.1), y + Inches(0.05), Inches(2.75), Inches(0.2),
             label, size=10, color=RGBColor(0xBF, 0xDB, 0xFE), bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(0.22), Inches(2.75), Inches(0.4),
             stage_title, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    txBox = add_text(slide, x + Inches(0.15), y + Inches(0.75), Inches(2.65), Inches(1.8), "", size=11, color=DARK)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for line in desc.split('\n'):
        add_para(tf, "  " + line, size=11, color=DARK, space_before=Pt(6))

    if i < 3:
        add_text(slide, x + Inches(2.95), y + Inches(1.2), Inches(0.2), Inches(0.3),
                 "->", size=18, color=GRAY, bold=True)

# Additional TQC features
add_shape(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8), LIGHT_BLUE, radius=0.04)
add_text(slide, Inches(0.7), Inches(5.4), Inches(11.5), Inches(0.3),
         "Additional TQC Features", size=14, color=NAVY, bold=True)

tqc_features = [
    "Periodic Interviews: Auto-scheduled at 1 week / 1 month / 3 months (/new-tqc/meetings)",
    "Color-blind Test: PASS/FAIL recorded at hire (quality inspection job assignment criteria)",
    "Resignation Tracking: 8 reason categories, turnover rate analysis during training period",
    "TQC Dashboard: New hire count, completion rate, turnover rate, avg onboarding duration",
]

txBox = add_text(slide, Inches(0.7), Inches(5.75), Inches(11.5), Inches(1.1), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for f in tqc_features:
    add_para(tf, "  " + f, size=10, color=DARK, space_before=Pt(4))


# ════════════════════════════════════════════════════════════
# SCENARIO 5: 5PRS Training Recommendation System
# ════════════════════════════════════════════════════════════
scenario_overview(prs, "5", "5PRS Training Recommendation System",
    "Mr. Moon", "QIP Manager (ADMIN)", "ADMIN", BLUE,
    "Based on 5PRS defect analysis results, the system\nautomatically recommends training programs and employees.",
    ["Run analysis", "Review recommendations", "Register training", "Track results"],
    [NAVY, ACCENT, GREEN, BLUE])

# Scenario 5 Detail: 5PRS Data Example
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Scenario 5: 5PRS Defect Rate Data — Details", "Defect analysis-based training recommendations")

prs_data = [
    ["Bldg/Line", "Defect Type", "Rate", "Threshold", "Status", "Recommended Training"],
    ["Bldg A L3", "Surface Scratch", "3.2%", "1.5%", "EXCEEDED", "QIP-001 Visual Inspection Basic"],
    ["Bldg A L3", "Stitching Defect", "1.8%", "2.0%", "NORMAL", "-"],
    ["Bldg B L1", "Adhesion Defect", "2.5%", "1.0%", "EXCEEDED", "QIP-003 Adhesion Inspection"],
    ["Bldg B L2", "Color Variance", "0.8%", "1.5%", "NORMAL", "-"],
    ["Bldg C L1", "Dimension Defect", "4.1%", "2.0%", "EXCEEDED", "QIP-005 Dimension Inspection"],
]

add_table_slide(slide, Inches(0.5), Inches(2.3), Inches(12.3), prs_data,
                [Inches(1.5), Inches(2), Inches(1.0), Inches(1.2), Inches(1.4), Inches(5.2)])

# Recommendation flow
add_shape(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(2.0), GREEN_BG, radius=0.04)
add_text(slide, Inches(0.7), Inches(5.3), Inches(11.5), Inches(0.3),
         "5PRS Recommendation System Workflow", size=13, color=GREEN, bold=True)

txBox = add_text(slide, Inches(0.7), Inches(5.7), Inches(11.5), Inches(1.3), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
add_para(tf, "  1. Collect 5PRS defect data (by building/line/defect type)", size=10, color=DARK, space_before=Pt(4))
add_para(tf, "  2. Auto-identify items exceeding threshold", size=10, color=DARK, space_before=Pt(4))
add_para(tf, "  3. Recommend mapped training programs for each defect type", size=10, color=DARK, space_before=Pt(4))
add_para(tf, "  4. Auto-select affected line workers -> Create training session -> Track results end-to-end", size=10, color=DARK, space_before=Pt(4))


# ════════════════════════════════════════════════════════════
# SCENARIO 6: CAPA Process Integration
# ════════════════════════════════════════════════════════════
scenario_overview(prs, "6", "CAPA Process Integration",
    "Mr. Moon", "QIP Manager (ADMIN)", "ADMIN", BLUE,
    "Defect rate in Building A L3 exceeds threshold.\nCAPA process initiated for corrective/preventive action.",
    ["Register CAPA", "5-Why Analysis", "Execute corrective action", "Link to training", "Verify results"],
    [RED, ORANGE, BLUE, GREEN, ACCENT])

# Scenario 6 Detail: CAPA 5 stages
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Scenario 6: CAPA 5-Stage Process — Details", "Systematic corrective/preventive action (/capa)")

capa_stages = [
    ("1", "Discovery", "Problem description\nSource classification\nImmediate action", RED),
    ("2", "Investigation", "Root cause analysis\n5-Why Analysis\nImpact assessment", ORANGE),
    ("3", "Action", "Corrective plan\nPreventive plan\nAssignees/deadlines", BLUE),
    ("4", "Verification", "Effectiveness test\nMonitoring period\nRecurrence check", ACCENT),
    ("5", "Closure", "Lessons learned\nKnowledge sharing\nDocumentation", GREEN),
]

for i, (num, stage_title, desc, color) in enumerate(capa_stages):
    x = Inches(0.3) + i * Inches(2.55)
    y = Inches(2.3)

    add_shape(slide, x, y, Inches(2.35), Inches(2.5), WHITE, radius=0.04)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85), y + Inches(0.12), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.paragraphs[0].text = num
    ctf.paragraphs[0].font.size = Pt(16)
    ctf.paragraphs[0].font.color.rgb = WHITE
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, x + Inches(0.1), y + Inches(0.75), Inches(2.15), Inches(0.45),
             stage_title, size=12, color=color, bold=True, align=PP_ALIGN.CENTER)

    txBox = add_text(slide, x + Inches(0.1), y + Inches(1.25), Inches(2.15), Inches(1.0), "", size=10, color=DARK)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for line in desc.split('\n'):
        add_para(tf, line, size=10, color=DARK, space_before=Pt(4), align=PP_ALIGN.CENTER)

    if i < 4:
        add_text(slide, x + Inches(2.35), y + Inches(1.0), Inches(0.2), Inches(0.3),
                 "->", size=16, color=GRAY, bold=True)

# CAPA Example
add_shape(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.2), RED_BG, radius=0.04)
add_text(slide, Inches(0.7), Inches(5.1), Inches(11.5), Inches(0.3),
         "Scenario Example: CAPA-2026-003", size=14, color=RED, bold=True)

txBox = add_text(slide, Inches(0.7), Inches(5.45), Inches(11.5), Inches(1.6), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
add_para(tf, "  Problem: Building A L3 visual inspection defect rate 3.2% (threshold: 1.5%)", size=10, color=DARK, space_before=Pt(3))
add_para(tf, "  Root Cause: QIP-001 training was theory-only -> lacking hands-on application (5-Why result)", size=10, color=DARK, space_before=Pt(3))
add_para(tf, "  Corrective Action: Supplementary training + Visual guide created + 2hrs practical added to curriculum", size=10, color=DARK, space_before=Pt(3))
add_para(tf, "  Training Link: QIP-001 retraining session created -> 30 line workers enrolled", size=10, color=BLUE, bold=True, space_before=Pt(3))
add_para(tf, "  Result: Defect rate 3.2% -> 1.1% (within threshold) | Effectiveness: 85/100 | CAPA -> Program improvement complete", size=10, color=GREEN, bold=True, space_before=Pt(3))


# ════════════════════════════════════════════════════════════
# SCENARIO 7: Competency Matrix Management
# ════════════════════════════════════════════════════════════
scenario_overview(prs, "7", "Competency Matrix Management",
    "Mr. Moon", "QIP Manager (ADMIN)", "ADMIN", BLUE,
    "Assess department-level competency and create\nlearning paths based on gap analysis.",
    ["Define competencies", "Assign employees", "Gap analysis", "Create learning paths"],
    [NAVY, BLUE, ORANGE, GREEN])

# Scenario 7 Detail: Competency Tables
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Scenario 7: Competency & Skill Gap — Details", "Competency assessment -> Gap analysis -> Learning path recommendations")

add_text(slide, Inches(0.5), Inches(2.2), Inches(6), Inches(0.3),
         "Individual Competency Assessment (/competency)", size=14, color=NAVY, bold=True)

comp_data = [
    ["Competency", "Current Level", "Target Level", "Gap"],
    ["Visual Inspection", "COMPETENT (3)", "PROFICIENT (4)", "-1"],
    ["Stitching QC", "BEGINNER (2)", "COMPETENT (3)", "-1"],
    ["Safety Procedures", "COMPETENT (3)", "COMPETENT (3)", "0"],
    ["Team Leadership", "NOVICE (1)", "NOVICE (1)", "0"],
]
add_table_slide(slide, Inches(0.5), Inches(2.55), Inches(6), comp_data,
                [Inches(1.5), Inches(1.5), Inches(1.5), Inches(0.8)])

add_text(slide, Inches(6.8), Inches(2.2), Inches(6), Inches(0.3),
         "Team/Dept Skill Gap Analysis (/skill-gap)", size=14, color=NAVY, bold=True)

gap_data = [
    ["Competency", "Team Avg", "Target", "Gap", "Below Target"],
    ["Visual Inspection", "2.8", "4.0", "-1.2", "15 employees"],
    ["Stitching QC", "2.5", "3.0", "-0.5", "8 employees"],
    ["Safety Procedures", "3.2", "3.0", "+0.2", "0 employees"],
]
add_table_slide(slide, Inches(6.8), Inches(2.55), Inches(6), gap_data,
                [Inches(1.3), Inches(1), Inches(1), Inches(0.8), Inches(1)])

# Learning path
add_shape(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.2), LIGHT_BLUE, radius=0.04)
add_text(slide, Inches(0.7), Inches(5.1), Inches(11.5), Inches(0.3),
         "Learning Paths & Development Plans", size=14, color=NAVY, bold=True)

lp_items = [
    "5 Competency Levels: NOVICE (1) -> BEGINNER (2) -> COMPETENT (3) -> PROFICIENT (4) -> EXPERT (5)",
    "6 Competency Categories: TECHNICAL / QUALITY / SAFETY / LEADERSHIP / COMMUNICATION / PROCESS",
    "Learning Path Types: ONBOARDING / POSITION / PROMOTION / SPECIALIZATION / REMEDIAL",
    "Individual Development Plans (IDP): Target competencies, assigned paths, manager comments, progress tracking",
    "Competency-Program Mapping: Auto-recommend training programs for each competency gap",
]

txBox = add_text(slide, Inches(0.7), Inches(5.5), Inches(11.5), Inches(1.5), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in lp_items:
    add_para(tf, "  " + item, size=10, color=DARK, space_before=Pt(5))


# ════════════════════════════════════════════════════════════
# SCENARIO 8: Dashboard & Report Utilization
# ════════════════════════════════════════════════════════════
scenario_overview(prs, "8", "Dashboard & Report Utilization",
    "Mr. Duc", "Line Manager, Building A (VIEWER)", "VIEWER", ORANGE,
    "Mr. Duc checks his team's training status\non the dashboard and generates reports.",
    ["Check dashboard", "Completion/pass rates", "Department analysis", "Generate reports"],
    [NAVY, BLUE, GREEN, ACCENT])

# Scenario 8 Detail: Dashboard & KPIs
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Scenario 8: Dashboard & Reports — Details", "Real-time KPI monitoring and reporting")

# KPI Cards
add_card(slide, Inches(0.5), Inches(2.3), Inches(2.9), Inches(1.3),
         "Total Employees", "1,200", "Active employees", LIGHT_BLUE, BLUE, NAVY)
add_card(slide, Inches(3.6), Inches(2.3), Inches(2.9), Inches(1.3),
         "Monthly Completed", "3", "Monthly sessions", LIGHT_BLUE, BLUE, NAVY)
add_card(slide, Inches(6.7), Inches(2.3), Inches(2.9), Inches(1.3),
         "Completion Rate", "67.2%", "Target: 90%", LIGHT_BLUE, BLUE, NAVY)
add_card(slide, Inches(9.8), Inches(2.3), Inches(2.9), Inches(1.3),
         "Needs Retraining", "45", "Retraining required", RED_BG, RED, RED)

# Dashboard + Executive
add_shape(slide, Inches(0.5), Inches(3.9), Inches(6), Inches(3.3), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(4.05), Inches(5.5), Inches(0.3),
         "Main Dashboard (/dashboard)", size=14, color=NAVY, bold=True)

dash_items = [
    "Grade distribution pie chart (AA/A/B/C)",
    "Monthly training trend (planned vs completed)",
    "Building-level retraining/expiry chart",
    "Certifications expiring within 30 days",
    "Upcoming training sessions schedule",
    "Recent retraining candidates list",
]
txBox = add_text(slide, Inches(0.7), Inches(4.45), Inches(5.5), Inches(2.5), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in dash_items:
    add_para(tf, "  " + item, size=11, color=DARK, space_before=Pt(5))

add_shape(slide, Inches(6.8), Inches(3.9), Inches(6), Inches(3.3), WHITE, radius=0.04)
add_text(slide, Inches(7), Inches(4.05), Inches(5.5), Inches(0.3),
         "Executive Dashboard & Reports", size=14, color=NAVY, bold=True)

exec_data = [
    ["KPI", "Current", "Target", "Achievement"],
    ["Completion Rate", "67.2%", "90%", "74.7%"],
    ["Qualification Rate", "58.3%", "85%", "68.6%"],
    ["New Hire Turnover", "8.2%", "5%", "-"],
    ["Training ROI", "187%", "200%", "93.5%"],
]
add_table_slide(slide, Inches(7), Inches(4.5), Inches(5.5), exec_data,
                [Inches(1.5), Inches(1), Inches(1), Inches(1)])

add_text(slide, Inches(7), Inches(6.7), Inches(5.5), Inches(0.4),
         "Reports: By department, program, employee\nPeriod filters (1M/3M/6M/1Y) | CSV/PDF export",
         size=10, color=GRAY)


# ════════════════════════════════════════════════════════════
# SLIDE: Data Update Checklist
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Data Update Checklist", "Daily / Weekly / Monthly tasks by role")

checklist_data = [
    ["Frequency", "Role", "Task", "System Location", "Notes"],
    ["Daily", "TRAINER", "Record attendance & enter results", "/attendance, /results", "Same day as training"],
    ["Daily", "ADMIN", "Check dashboard for anomalies", "/dashboard", "KPI monitoring"],
    ["Weekly", "TRAINER", "Review retraining list & assign sessions", "/retraining", "Process FAIL/ABSENT"],
    ["Weekly", "ADMIN", "Review program change history", "/programs", "Check change_logs"],
    ["Weekly", "TRAINER", "Conduct TQC new hire interviews", "/new-tqc/meetings", "Week 1 interview"],
    ["Monthly", "ADMIN", "Executive dashboard review", "/executive", "KPI report"],
    ["Monthly", "ADMIN", "Check expiring certifications & plan retraining", "/retraining", "Within 30 days"],
    ["Monthly", "ADMIN", "Review CAPA status", "/capa", "Check in-progress items"],
    ["Monthly", "ADMIN", "Skill gap analysis & learning path updates", "/skill-gap", "Quarterly deep dive"],
    ["Monthly", "ADMIN", "Generate & share reports", "/reports", "CSV/PDF export"],
]

add_table_slide(slide, Inches(0.5), Inches(2.3), Inches(12.3), checklist_data,
                [Inches(0.9), Inches(1.2), Inches(3.5), Inches(2.5), Inches(4.2)])


# ════════════════════════════════════════════════════════════
# SLIDE: FAQ / Contact
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "FAQ / Contact", "Frequently asked questions and support")

# FAQ section
add_shape(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.0), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(2.4), Inches(11.5), Inches(0.3),
         "Frequently Asked Questions", size=16, color=NAVY, bold=True)

faqs = [
    ("Q: What if I entered a wrong result?",
     "A: Results can be edited, but a mandatory reason must be provided. All edit history is permanently stored in result_edit_logs."),
    ("Q: Can results be deleted?",
     "A: No. Per the No Delete Policy, training results can never be deleted. This is a core principle for audit traceability."),
    ("Q: What happens when a certificate expires?",
     "A: The system auto-changes status to EXPIRING (30 days before) / EXPIRED and sends notifications to the employee and manager."),
    ("Q: Where does 5PRS recommendation data come from?",
     "A: Currently, admins input defect data manually. Future plans include automatic MES system integration."),
    ("Q: Can I enter results for multiple trainees at once?",
     "A: Yes. Batch entry per session is available. Excel template upload for bulk data entry is also supported."),
]

y_pos = Inches(2.8)
for q, a in faqs:
    add_text(slide, Inches(0.7), y_pos, Inches(11.5), Inches(0.25),
             q, size=11, color=NAVY, bold=True)
    add_text(slide, Inches(0.7), y_pos + Inches(0.25), Inches(11.5), Inches(0.25),
             a, size=10, color=DARK)
    y_pos += Inches(0.6)

# Contact section
add_shape(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.8), LIGHT_BLUE, radius=0.04)
add_text(slide, Inches(0.7), Inches(6.55), Inches(11.5), Inches(0.3),
         "Contact", size=14, color=NAVY, bold=True)
add_text(slide, Inches(0.7), Inches(6.85), Inches(11.5), Inches(0.3),
         "System inquiries: ksmoon@hsvina.com  |  Q-TRAIN System: https://q-train-web.web.app",
         size=11, color=DARK)


# ════════════════════════════════════════════════════════════
# SLIDE: Thank You
# ════════════════════════════════════════════════════════════
thank_you_slide(prs, "https://q-train-web.web.app", "ksmoon@hsvina.com")


# ════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "Q-TRAIN_Training_Guide_EN.pptx")
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
