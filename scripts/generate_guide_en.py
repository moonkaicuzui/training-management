#!/usr/bin/env python3
"""
Q-TRAIN System Complete Guide — English PowerPoint Presentation
8 Feature Categories with scenario-based storytelling
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors ──
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
BLUE = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
GREEN = RGBColor(0x05, 0x96, 0x69)
LIGHT_GREEN = RGBColor(0xEC, 0xFD, 0xF5)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PINK = RGBColor(0xDB, 0x27, 0x77)
TEAL = RGBColor(0x06, 0x69, 0xB4)
AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)
RED_BG = RGBColor(0xFE, 0xF2, 0xF2)
BLUE_BG = RGBColor(0xEF, 0xF6, 0xFF)
PURPLE_BG = RGBColor(0xF5, 0xF3, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ════════════════════════════════════════
# Helpers
# ════════════════════════════════════════

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, l, t, w, h, color, radius=None):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(st, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    if radius:
        try: s.adjustments[0] = radius
        except: pass
    return s

def add_text(slide, l, t, w, h, text, sz=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Arial'
    p.alignment = align
    return tb

def add_para(tf, text, sz=14, color=DARK, bold=False, sp_before=Pt(4), sp_after=Pt(2), align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Arial'
    p.alignment = align
    p.space_before = sp_before
    p.space_after = sp_after
    return p

def header(slide, title, subtitle="", bg=WHITE):
    add_bg(slide, bg)
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), BLUE)
    add_text(slide, Inches(0.5), Inches(0.22), Inches(3), Inches(0.35), "Q-TRAIN", sz=13, color=BLUE, bold=True)
    add_text(slide, Inches(0.5), Inches(0.5), Inches(6), Inches(0.25), "HWK Vietnam QIP Training Management System", sz=8, color=GRAY)
    add_text(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.5), title, sz=26, color=NAVY, bold=True)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(0.35), subtitle, sz=13, color=GRAY)
    add_shape(slide, Inches(0.5), Inches(1.9), Inches(1.5), Inches(0.04), BLUE)

def add_table(slide, l, t, w, rows, col_w, hdr_color=NAVY):
    nrows = len(rows)
    ncols = len(rows[0])
    ts = slide.shapes.add_table(nrows, ncols, l, t, w, Inches(0.38 * nrows))
    tbl = ts.table
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = cw
    for r, rd in enumerate(rows):
        for c, ct in enumerate(rd):
            cell = tbl.cell(r, c)
            cell.text = str(ct)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = 'Arial'
                if r == 0:
                    p.font.color.rgb = WHITE
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
                else:
                    p.font.color.rgb = DARK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hdr_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    return ts

def card(slide, l, t, w, h, title, lines, bar_color=BLUE, bg_color=WHITE):
    add_shape(slide, l, t, w, h, bg_color, radius=0.04)
    add_shape(slide, l, t, Inches(0.06), h, bar_color)
    add_text(slide, l + Inches(0.2), t + Inches(0.1), w - Inches(0.3), Inches(0.3), title, sz=13, color=bar_color, bold=True)
    tb = add_text(slide, l + Inches(0.2), t + Inches(0.45), w - Inches(0.3), h - Inches(0.5), "", sz=10, color=DARK)
    tf = tb.text_frame
    tf.paragraphs[0].text = ""
    for line in lines:
        add_para(tf, line, sz=10, color=DARK, sp_before=Pt(3))
    return tb

def category_header(slide, num, title, subtitle, color):
    """Category intro slide header"""
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), W, Inches(3.2), color)
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), WHITE)
    # Category number
    add_text(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4), f"CATEGORY {num}", sz=12, color=RGBColor(0xBF, 0xDB, 0xFE), bold=True)
    add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7), title, sz=36, color=WHITE, bold=True)
    add_text(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5), subtitle, sz=16, color=RGBColor(0xBF, 0xDB, 0xFE))
    # Accent line
    add_shape(slide, Inches(0.8), Inches(2.5), Inches(2), Inches(0.04), WHITE)


def scenario_box(slide, l, t, w, h, scene_title, who, what, result, color=BLUE):
    """A scenario storytelling box"""
    add_shape(slide, l, t, w, h, WHITE, radius=0.04)
    # Scene icon bar
    add_shape(slide, l, t, w, Inches(0.4), color, radius=0.04)
    add_shape(slide, l, t + Inches(0.25), w, Inches(0.15), color)
    add_text(slide, l + Inches(0.15), t + Inches(0.05), w - Inches(0.3), Inches(0.3), scene_title, sz=11, color=WHITE, bold=True)
    # Who
    add_text(slide, l + Inches(0.15), t + Inches(0.5), w - Inches(0.3), Inches(0.2), who, sz=9, color=color, bold=True)
    # What
    tb = add_text(slide, l + Inches(0.15), t + Inches(0.75), w - Inches(0.3), h - Inches(1.2), "", sz=10, color=DARK)
    tf = tb.text_frame
    tf.paragraphs[0].text = ""
    for line in what:
        add_para(tf, "  " + line, sz=10, color=DARK, sp_before=Pt(3))
    # Result
    add_text(slide, l + Inches(0.15), t + h - Inches(0.35), w - Inches(0.3), Inches(0.3), result, sz=9, color=GREEN, bold=True)


# ════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, NAVY)
add_shape(s, Inches(0), Inches(0), W, Inches(0.08), BLUE)
add_shape(s, Inches(0), Inches(7.0), W, Inches(0.5), RGBColor(0x15, 0x2A, 0x4A))
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(1), "Q-TRAIN", sz=60, color=WHITE, bold=True)
add_text(s, Inches(1), Inches(2.2), Inches(11), Inches(0.6), "Complete System Guide", sz=30, color=RGBColor(0x93, 0xC5, 0xFD))
add_shape(s, Inches(1), Inches(3.1), Inches(2), Inches(0.05), BLUE)
add_text(s, Inches(1), Inches(3.5), Inches(10), Inches(0.5), "8 Feature Categories — Scenario-Based Learning", sz=20, color=RGBColor(0xBF, 0xDB, 0xFE))
add_text(s, Inches(1), Inches(4.3), Inches(5), Inches(0.3), "HWK Vietnam | QIP Training Management", sz=14, color=RGBColor(0x93, 0xC5, 0xFD))
add_text(s, Inches(1), Inches(4.7), Inches(5), Inches(0.3), "https://q-train-web.web.app", sz=13, color=RGBColor(0x60, 0xA5, 0xFA))
add_text(s, Inches(1), Inches(6.5), Inches(5), Inches(0.3), "February 2026 | Version 1.0", sz=11, color=RGBColor(0x64, 0x78, 0x96))


# ════════════════════════════════════════════════════════════
# SLIDE 2: SYSTEM OVERVIEW — 8 CATEGORIES
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "System Overview", "Q-TRAIN covers 8 functional categories across 30+ pages")

cats = [
    ("1", "Training\nManagement", "Programs, Sessions\nAttendance, Results", BLUE),
    ("2", "Certificates &\nCompliance", "Issuance, Audit Trail\nAdidas Compliance", GREEN),
    ("3", "New Employee\nOnboarding", "TQC 4-Stage Process\nMeetings, Tracking", ORANGE),
    ("4", "Quality\nManagement", "CAPA 5-Stage\nCorrective & Preventive", RED),
    ("5", "Project\nManagement", "Tasks, Kanban, Gantt\nCalendar, Automation", PURPLE),
    ("6", "Employee\nDevelopment", "Competency, Skill Gap\nEvaluation, IDP", TEAL),
    ("7", "Knowledge &\nPlanning", "Materials Library\nAnnual Training Plans", PINK),
    ("8", "Analytics &\nDashboards", "KPIs, ROI, Reports\nExecutive Dashboard", NAVY),
]

for i, (num, title, desc, color) in enumerate(cats):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + col * Inches(3.1)
    y = Inches(2.2) + row * Inches(2.5)
    add_shape(s, x, y, Inches(2.9), Inches(2.2), WHITE, radius=0.04)
    # Color top bar
    add_shape(s, x, y, Inches(2.9), Inches(0.5), color, radius=0.04)
    add_shape(s, x, y + Inches(0.3), Inches(2.9), Inches(0.2), color)
    # Number
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.1), y + Inches(0.08), Inches(0.35), Inches(0.35))
    c.fill.solid(); c.fill.fore_color.rgb = WHITE; c.line.fill.background()
    tf = c.text_frame; tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(12); tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(s, x + Inches(0.55), y + Inches(0.08), Inches(2.2), Inches(0.4), title, sz=12, color=WHITE, bold=True)
    add_text(s, x + Inches(0.15), y + Inches(0.6), Inches(2.6), Inches(1.4), desc, sz=11, color=DARK)


# ════════════════════════════════════════════════════════════
# SLIDE 3: CHARACTERS
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "Meet the Team", "Four characters guide you through every scenario")

chars = [
    ("Mr. Moon", "QIP Manager", "ADMIN", BLUE,
     ["Full system control", "Create programs & settings", "Access all reports & audit logs", "Manage users & permissions"]),
    ("Ms. Mai", "QA Team Lead / Trainer", "TRAINER", GREEN,
     ["Create & conduct sessions", "Record attendance & results", "Request certificates", "Evaluate training effectiveness"]),
    ("Mr. Duc", "Line Manager, Bldg A", "VIEWER", ORANGE,
     ["View dashboards & reports", "Monitor team progress", "Receive notifications", "Check training schedules"]),
    ("Ms. Hoa", "New Worker (Day 3)", "TRAINEE", PURPLE,
     ["Attend training sessions", "Take exams", "Complete TQC onboarding", "Receive certificates"]),
]

for i, (name, role, badge, color, perms) in enumerate(chars):
    x = Inches(0.5) + i * Inches(3.15)
    y = Inches(2.2)
    add_shape(s, x, y, Inches(2.95), Inches(4.8), WHITE, radius=0.04)
    add_shape(s, x, y, Inches(2.95), Inches(0.5), color, radius=0.04)
    add_shape(s, x, y + Inches(0.3), Inches(2.95), Inches(0.2), color)
    # Badge
    b = add_shape(s, x + Inches(0.12), y + Inches(0.1), Inches(1), Inches(0.28), WHITE, radius=0.1)
    tf = b.text_frame; tf.paragraphs[0].text = badge
    tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(s, x + Inches(0.15), y + Inches(0.65), Inches(2.65), Inches(0.3), name, sz=16, color=NAVY, bold=True)
    add_text(s, x + Inches(0.15), y + Inches(0.95), Inches(2.65), Inches(0.25), role, sz=11, color=GRAY)
    add_shape(s, x + Inches(0.15), y + Inches(1.3), Inches(2.65), Inches(0.02), LIGHT_BLUE)
    tb = add_text(s, x + Inches(0.15), y + Inches(1.45), Inches(2.65), Inches(3), "", sz=10, color=DARK)
    tf = tb.text_frame; tf.paragraphs[0].text = ""
    for p in perms:
        add_para(tf, "  " + p, sz=11, color=DARK, sp_before=Pt(8))


# ════════════════════════════════════════════════════════════
# CATEGORY 1: TRAINING MANAGEMENT (2 slides)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
category_header(s, "01", "Training Management", "The core workflow: Programs → Sessions → Attendance → Results", BLUE)

# Scenario cards below the header
scenario_box(s, Inches(0.5), Inches(3.5), Inches(6.1), Inches(3.7),
    "Scene: Mr. Moon creates a training program",
    "WHO: Mr. Moon (ADMIN) | WHERE: /programs",
    [
        "Creates QIP-001 'Visual Inspection Basic'",
        "Sets trilingual names (EN / VI / KR)",
        "Category: QIP | Target: Workers & Line Leaders",
        "Evaluation: Score-based | Pass: 70 pts",
        "Grade thresholds: AA≥95, A≥85, B≥70",
        "Validity: 12 months (auto-triggers retraining)",
        "All changes logged in program_change_logs",
    ],
    "Result: Program ready for scheduling", BLUE)

scenario_box(s, Inches(6.8), Inches(3.5), Inches(6.1), Inches(3.7),
    "Scene: Ms. Mai schedules & conducts training",
    "WHO: Ms. Mai (TRAINER) | WHERE: /schedule → /attendance → /results",
    [
        "Creates session: Feb 20, 09:00-13:00, Room A",
        "Assigns 28 workers from Building A",
        "Participants receive TRAINING_REMINDER alert",
        "On training day: records attendance (Present/Late/Absent)",
        "After exam: enters scores — grades auto-calculated",
        "FAIL/ABSENT → needs_retraining auto-flagged",
        "Result edits require mandatory reason (audit trail)",
    ],
    "Result: 26 passed, 1 failed, 1 absent → retraining list updated", GREEN)

# Slide: Training Rules & Flow
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "Training Management — Key Rules", "Critical business rules enforced by the system")

# Flow steps
steps = [
    ("Create\nProgram", NAVY), ("Schedule\nSession", BLUE), ("Record\nAttendance", TEAL),
    ("Enter\nResults", GREEN), ("Auto-Grade\nCalculation", PURPLE), ("Retraining\nFlag", ORANGE),
]
for i, (label, color) in enumerate(steps):
    x = Inches(0.4) + i * Inches(2.1)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.35), Inches(2.2), Inches(0.9), Inches(0.9))
    c.fill.solid(); c.fill.fore_color.rgb = color; c.line.fill.background()
    tf = c.text_frame; tf.paragraphs[0].text = str(i+1)
    tf.paragraphs[0].font.size = Pt(18); tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(s, x, Inches(3.2), Inches(1.6), Inches(0.5), label, sz=10, color=color, bold=True, align=PP_ALIGN.CENTER)
    if i < 5:
        add_text(s, x + Inches(1.5), Inches(2.4), Inches(0.6), Inches(0.4), "→", sz=20, color=LIGHT_GRAY, bold=True, align=PP_ALIGN.CENTER)

# Rules box
add_shape(s, Inches(0.5), Inches(3.9), Inches(12.3), Inches(3.3), AMBER_BG, radius=0.04)
add_text(s, Inches(0.7), Inches(4.0), Inches(11.5), Inches(0.3), "Business Rules", sz=15, color=RGBColor(0x92, 0x40, 0x0E), bold=True)
rules = [
    ("No Delete Policy", "Training results can NEVER be deleted — only edited with mandatory reason logging"),
    ("Auto-Grade", "Score → Grade auto-calculated: AA (≥95) | A (≥85) | B (≥70) | C (<70)"),
    ("Auto-Retraining", "FAIL or ABSENT → automatically flagged for retraining"),
    ("Audit Trail", "Every program change and result edit is permanently logged with before/after data"),
    ("Validity Tracking", "Programs with validity period trigger EXPIRING (30 days) and EXPIRED alerts"),
    ("Edit Logging", "Result updates use atomic transactions — result + edit log written together"),
]
tb = add_text(s, Inches(0.7), Inches(4.35), Inches(11.5), Inches(2.7), "", sz=10, color=DARK)
tf = tb.text_frame; tf.paragraphs[0].text = ""
for title, desc in rules:
    add_para(tf, f"  {title}:  {desc}", sz=10, color=RGBColor(0x78, 0x35, 0x0F), sp_before=Pt(6))


# ════════════════════════════════════════════════════════════
# CATEGORY 2: CERTIFICATES & COMPLIANCE
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
category_header(s, "02", "Certificates & Compliance", "Automated certification, audit trails, and Adidas compliance tracking", GREEN)

scenario_box(s, Inches(0.5), Inches(3.5), Inches(4), Inches(3.7),
    "Scene: Batch certificate issuance",
    "WHO: Mr. Moon (ADMIN) | WHERE: /certificates",
    [
        "Selects session: QIP-001, Feb 20",
        "System auto-filters: 26 PASS results",
        "Clicks 'Batch Issue' → 26 certs generated",
        "Each cert: CERT-2026-XXXXXX (auto-number)",
        "Includes: score, grade, trainer, expiry date",
        "Custom template: company logo & signature",
        "Recipients get CERTIFICATE_READY alert",
    ],
    "Result: 26 PDF certificates downloadable", GREEN)

scenario_box(s, Inches(4.7), Inches(3.5), Inches(4), Inches(3.7),
    "Scene: Adidas audit preparation",
    "WHO: Mr. Moon (ADMIN) | WHERE: /audit",
    [
        "6 audit categories tracked:",
        "  Training Coverage — 95%",
        "  Certification — 92%",
        "  Documentation — 88%",
        "  Competency Evaluation — 72%",
        "  Retraining Management — 90%",
        "  Records Retention — 100%",
    ],
    "Result: Compliance report exported for auditor", NAVY)

scenario_box(s, Inches(8.9), Inches(3.5), Inches(4), Inches(3.7),
    "Scene: Investigating a past edit",
    "WHO: Mr. Moon (ADMIN) | WHERE: /audit-log",
    [
        "Full audit trail for every system action",
        "Entities: Program, Result, Session, Employee",
        "Actions: CREATE, UPDATE, DELETE, LOGIN",
        "Filter by entity, action, user, date range",
        "Before/after data diff view",
        "IP address tracking",
        "Append-only — logs can never be deleted",
    ],
    "Result: Complete traceability for any audit", PURPLE)


# ════════════════════════════════════════════════════════════
# CATEGORY 3: NEW EMPLOYEE ONBOARDING (TQC)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
category_header(s, "03", "New Employee Onboarding", "TQC 4-stage process: From Day 1 to fully qualified worker", ORANGE)

# 4 stages
stages_data = [
    ("Stage 1", "Orientation", "Factory safety rules\nCompany policies\nBasic quality concepts\nWorkplace tour", GREEN),
    ("Stage 2", "Basic Training", "QIP program courses\nExam & scoring\nMust pass all required\nprograms to advance", BLUE),
    ("Stage 3", "Line Assignment", "On-the-job training\nSenior worker mentoring\nReal production line\nHands-on skill building", ORANGE),
    ("Stage 4", "Field Evaluation", "Independent work test\nFinal qualification exam\nCertificate issuance\nFull worker status", PURPLE),
]

for i, (label, title, desc, color) in enumerate(stages_data):
    x = Inches(0.5) + i * Inches(3.15)
    y = Inches(3.5)
    add_shape(s, x, y, Inches(2.9), Inches(2.5), WHITE, radius=0.04)
    add_shape(s, x, y, Inches(2.9), Inches(0.6), color, radius=0.04)
    add_shape(s, x, y + Inches(0.4), Inches(2.9), Inches(0.2), color)
    add_text(s, x + Inches(0.1), y + Inches(0.05), Inches(2.7), Inches(0.2), label, sz=9, color=RGBColor(0xBF,0xDB,0xFE), bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(0.22), Inches(2.7), Inches(0.3), title, sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb = add_text(s, x + Inches(0.15), y + Inches(0.7), Inches(2.6), Inches(1.6), "", sz=10, color=DARK)
    tf = tb.text_frame; tf.paragraphs[0].text = ""
    for line in desc.split('\n'):
        add_para(tf, "  " + line, sz=10, color=DARK, sp_before=Pt(5))
    if i < 3:
        add_text(s, x + Inches(2.9), Inches(4.5), Inches(0.25), Inches(0.3), "→", sz=16, color=GRAY, bold=True)

# Additional features
add_shape(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(1.0), LIGHT_BLUE, radius=0.04)
add_text(s, Inches(0.7), Inches(6.25), Inches(11.5), Inches(0.8),
    "Additional: Periodic interviews (1 week / 1 month / 3 months)  |  Color-blind test tracking  |  "
    "Resignation tracking (8 reason categories)  |  TQC Dashboard with onboarding KPIs",
    sz=10, color=NAVY)


# ════════════════════════════════════════════════════════════
# CATEGORY 4: QUALITY MANAGEMENT (CAPA)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
category_header(s, "04", "Quality Management (CAPA)", "5-stage Corrective & Preventive Action process", RED)

capa_stages = [
    ("1", "Discovery", "Problem description\nSource classification\nImmediate action", RED),
    ("2", "Investigation", "Root cause analysis\n5-Why / Fishbone\nImpact assessment", ORANGE),
    ("3", "Action", "Corrective actions\nPreventive actions\nAssignees & deadlines", BLUE),
    ("4", "Verification", "Effectiveness test\nMonitoring period\nRecurrence check", PURPLE),
    ("5", "Closure", "Lessons learned\nKnowledge sharing\nDocumentation", GREEN),
]

for i, (num, title, desc, color) in enumerate(capa_stages):
    x = Inches(0.3) + i * Inches(2.55)
    y = Inches(3.4)
    add_shape(s, x, y, Inches(2.35), Inches(2.4), WHITE, radius=0.04)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85), y + Inches(0.12), Inches(0.55), Inches(0.55))
    c.fill.solid(); c.fill.fore_color.rgb = color; c.line.fill.background()
    tf = c.text_frame; tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16); tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(s, x + Inches(0.1), y + Inches(0.75), Inches(2.15), Inches(0.3), title, sz=12, color=color, bold=True, align=PP_ALIGN.CENTER)
    tb = add_text(s, x + Inches(0.1), y + Inches(1.1), Inches(2.15), Inches(1.1), "", sz=10, color=DARK)
    tf = tb.text_frame; tf.paragraphs[0].text = ""
    for line in desc.split('\n'):
        add_para(tf, line, sz=10, color=DARK, sp_before=Pt(4), align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(s, x + Inches(2.35), y + Inches(0.85), Inches(0.2), Inches(0.3), "→", sz=16, color=GRAY, bold=True)

# Example scenario
add_shape(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.2), RED_BG, radius=0.04)
add_text(s, Inches(0.7), Inches(6.1), Inches(11.5), Inches(0.25), "Scenario: CAPA-2026-003 — Defect rate exceeded threshold on Line L3", sz=12, color=RED, bold=True)
add_text(s, Inches(0.7), Inches(6.35), Inches(11.5), Inches(0.7),
    "Problem: Visual inspection defect rate 3.2% (threshold: 1.5%) → Root cause: QIP-001 training was theory-only, lacking hands-on practice\n"
    "Action: Supplementary training + Visual guide created + 2hrs practical added to curriculum → Result: Defect rate dropped to 1.1% | Effectiveness: 85/100",
    sz=10, color=DARK)


# ════════════════════════════════════════════════════════════
# CATEGORY 5: PROJECT MANAGEMENT (2 slides)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
category_header(s, "05", "Project Management", "Task tracking, Kanban boards, Gantt timelines, and team collaboration", PURPLE)

scenario_box(s, Inches(0.5), Inches(3.5), Inches(6.1), Inches(3.7),
    "Scene: Mr. Moon manages a QIP improvement project",
    "WHO: Mr. Moon (ADMIN) | WHERE: /projects/tasks",
    [
        "Creates project: 'Q1 Quality Improvement'",
        "Adds team members with roles (Admin/Member/Guest)",
        "Creates tasks: 'Update QIP-001 curriculum'",
        "  → assigns to Ms. Mai, due date: Mar 20",
        "  → priority: HIGH, status: TODO",
        "4 view modes: List / Kanban / Gantt / Calendar",
        "Drag-and-drop on Kanban board to change status",
    ],
    "Result: All tasks tracked with assignees and deadlines", PURPLE)

scenario_box(s, Inches(6.8), Inches(3.5), Inches(6.1), Inches(3.7),
    "Scene: Ms. Mai collaborates on a task",
    "WHO: Ms. Mai (MEMBER) | WHERE: Task detail panel",
    [
        "Opens task: 'Update QIP-001 curriculum'",
        "Posts a Question message (amber-coded):",
        "  'Should we add a practical exam section?'",
        "Uses @mention to tag Mr. Moon for input",
        "Mr. Moon replies with Approval (green-coded):",
        "  'Yes, add 2 hours of hands-on practice'",
        "Ms. Mai marks the question as Resolved",
    ],
    "Result: Decision documented within the task thread", TEAL)

# Slide: Project Features detail
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "Project Management — Features in Detail", "4 task views, team messaging, calendar events, and workflow automation")

card(s, Inches(0.5), Inches(2.2), Inches(3.8), Inches(2.3), "4 Task View Modes", [
    "List View — tabular with status & priority badges",
    "Board View — Kanban with drag-and-drop",
    "  Columns: Todo → In Progress → Review → Done",
    "Timeline View — Gantt chart with date bars",
    "  Today marker, weekend shading, month nav",
    "Calendar View — full monthly/weekly/daily view",
], PURPLE)

card(s, Inches(4.5), Inches(2.2), Inches(4.2), Inches(2.3), "Task Communication (Slack-style)", [
    "6 message types: Question / Request / Approval /",
    "  Rejection / Info / Comment (color-coded)",
    "@Mentions with autocomplete popup",
    "Read receipts (who has seen each message)",
    "Resolve button for Questions & Requests",
    "Full conversation history per task",
], TEAL)

card(s, Inches(8.9), Inches(2.2), Inches(4), Inches(2.3), "Team & Access Control", [
    "3 roles: Admin / Member / Guest",
    "Admin: full CRUD on tasks, members, settings",
    "Member: create & edit tasks, view members",
    "Guest: read-only access",
    "Add members with name, email, department, role",
    "Activate / Deactivate without deleting data",
], BLUE)

card(s, Inches(0.5), Inches(4.7), Inches(6.1), Inches(2.5), "Workflow Automation (Monday.com-style)", [
    "IF-THEN rules that execute automatically",
    "",
    "7 Trigger types:",
    "  Task Status Changed, Assignee Changed, Due Date Approaching,",
    "  Task Created, Progress Changed, Scheduled Time, Custom Field",
    "",
    "8 Action types:",
    "  Change Status, Change Assignee, Send Notification, Send Email,",
    "  Create Task, Update Field, Add Comment, Extend Due Date",
    "",
    "Example: When task status → 'Review', notify Mr. Moon via email",
], ORANGE)

card(s, Inches(6.8), Inches(4.7), Inches(6.1), Inches(2.5), "Calendar & Events", [
    "Full calendar: Month / Week / Day / Agenda views",
    "Click any time slot to create an event",
    "",
    "Event fields: title, description, start/end datetime,",
    "  category (color-coded), location",
    "",
    "Category management in Settings tab",
    "  → Each category has a custom color",
    "  → Types: 'task' for tasks, 'event' for calendar",
    "",
    "Real-time sync — changes appear instantly for all members",
], GREEN)


# ════════════════════════════════════════════════════════════
# CATEGORY 6: EMPLOYEE DEVELOPMENT
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
category_header(s, "06", "Employee Development", "Competency assessment, skill gap analysis, training evaluation, and individual development plans", TEAL)

scenario_box(s, Inches(0.5), Inches(3.5), Inches(4), Inches(3.7),
    "Scene: Quarterly competency assessment",
    "WHO: Ms. Mai (TRAINER) | WHERE: /competency",
    [
        "Assesses Ms. Hoa's competencies:",
        "  Visual Inspection: COMPETENT (3/5)",
        "  Stitching QC: BEGINNER (2/5)",
        "  Safety: COMPETENT (3/5)",
        "5 levels: Novice → Beginner → Competent",
        "  → Proficient → Expert",
        "6 categories: Technical, Quality, Safety,",
        "  Leadership, Communication, Process",
    ],
    "Result: Skills mapped against target levels", TEAL)

scenario_box(s, Inches(4.7), Inches(3.5), Inches(4), Inches(3.7),
    "Scene: Skill gap analysis for Building A",
    "WHO: Mr. Moon (ADMIN) | WHERE: /skill-gap",
    [
        "Team skill gap report shows:",
        "  Visual Inspection: avg 2.8 / target 4.0",
        "    → Gap: -1.2 | 15 employees below target",
        "  Stitching QC: avg 2.5 / target 3.0",
        "    → Gap: -0.5 | 8 employees below target",
        "System auto-recommends learning paths:",
        "  'Visual Inspection Expert' pathway",
        "  QIP-001 → QIP-001A (Advanced) → OJT",
    ],
    "Result: Training priorities identified by data", BLUE)

scenario_box(s, Inches(8.9), Inches(3.5), Inches(4), Inches(3.7),
    "Scene: Post-training evaluation",
    "WHO: Ms. Mai (TRAINER) | WHERE: /evaluation",
    [
        "Kirkpatrick 4-Level Model:",
        "  L1 Reaction: Was it useful? (satisfaction)",
        "  L2 Learning: Did they learn? (pre/post test)",
        "  L3 Behavior: Applied on job? (follow-up)",
        "  L4 Results: Business impact? (ROI)",
        "6 weighted criteria (totaling 100%):",
        "  Content, Instructor, Materials, Environment,",
        "  Objectives, Job Applicability",
    ],
    "Result: Training effectiveness measured & ranked", PURPLE)


# ════════════════════════════════════════════════════════════
# CATEGORY 7: KNOWLEDGE & PLANNING
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
category_header(s, "07", "Knowledge Base & Annual Planning", "Training materials library, annual planning, and department status tracking", PINK)

scenario_box(s, Inches(0.5), Inches(3.5), Inches(6.1), Inches(3.7),
    "Scene: Organizing training materials",
    "WHO: Mr. Moon (ADMIN) | WHERE: /materials",
    [
        "Hierarchical folder structure with nested folders",
        "Uploads QIP-001 materials: PDF manual, video tutorial",
        "File metadata: version, tags, linked program, language",
        "File types: PDF, Video, PPT, Document, Link",
        "4 tabs: Browse / Favorites / Recent / By Program",
        "Star files for quick access",
        "Search by name or tags | Filter by file type",
        "Storage usage dashboard with capacity tracking",
    ],
    "Result: Centralized training resource library", PINK)

scenario_box(s, Inches(6.8), Inches(3.5), Inches(6.1), Inches(3.7),
    "Scene: Annual training plan creation",
    "WHO: Mr. Moon (ADMIN) | WHERE: /training-plan",
    [
        "Creates plan: '2026 Annual Training Plan'",
        "12-month calendar grid: visualize distribution",
        "Adds programs with monthly targets:",
        "  QIP-001: 12 sessions/year, 360 target participants",
        "  QIP-002: 6 sessions/year, 180 target participants",
        "Budget allocation per program",
        "Quality Analysis tab: pass rates, retraining needs",
        "Department Status: completion rate by department",
        "Auto-generated recommendations for underperforming depts",
    ],
    "Result: Data-driven annual plan with budget tracking", NAVY)


# ════════════════════════════════════════════════════════════
# CATEGORY 8: ANALYTICS & DASHBOARDS
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
category_header(s, "08", "Analytics & Dashboards", "Real-time KPIs, ROI tracking, multi-dimensional reports, and executive insights", NAVY)

scenario_box(s, Inches(0.5), Inches(3.5), Inches(4), Inches(3.7),
    "Scene: Daily monitoring",
    "WHO: Mr. Duc (VIEWER) | WHERE: /dashboard",
    [
        "4 KPI cards: Employees, Monthly Sessions,",
        "  Completion Rate, Retraining Count",
        "Grade distribution pie chart (AA/A/B/C)",
        "Monthly trend: planned vs. completed (6 months)",
        "Building-level alert chart",
        "Expiring certifications list (within 30 days)",
        "Upcoming sessions calendar",
        "Quick-action: navigate to schedules & reports",
    ],
    "Result: Instant overview of training health", BLUE)

scenario_box(s, Inches(4.7), Inches(3.5), Inches(4), Inches(3.7),
    "Scene: Executive monthly review",
    "WHO: Mr. Moon (ADMIN) | WHERE: /executive",
    [
        "4 strategic KPIs:",
        "  Completion Rate: 67.2% (target: 90%)",
        "  Qualification Rate: 58.3% (target: 85%)",
        "  New Hire Turnover: 8.2% (target: 5%)",
        "  Training ROI: 187% (target: 200%)",
        "ROI = (Productivity + Quality + Retention",
        "  savings) / Training Cost × 100",
        "Benchmarking: vs. Target / Industry / HWK Group",
    ],
    "Result: Strategic decisions backed by data", NAVY)

scenario_box(s, Inches(8.9), Inches(3.5), Inches(4), Inches(3.7),
    "Scene: Quarterly report generation",
    "WHO: Mr. Moon (ADMIN) | WHERE: /reports",
    [
        "3 report dimensions:",
        "  Department: completion %, avg score, pass rate",
        "  Program: sessions, trainees, retraining rate",
        "  Employee: individual training status",
        "Period filters: 1M / 3M / 6M / 1Y",
        "Export formats: CSV and PDF",
        "Notification system: expiry warnings,",
        "  session reminders, per-user settings",
    ],
    "Result: Multi-dimensional reports for any audience", GREEN)


# ════════════════════════════════════════════════════════════
# SLIDE: NOTIFICATION SYSTEM
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "Notification & Alert System", "Automated alerts keep everyone informed at the right time")

notif_data = [
    ["Notification Type", "Trigger", "Recipients", "Priority"],
    ["Training Reminder", "Upcoming session (configurable days)", "Assigned attendees", "MEDIUM"],
    ["Expiry Warning", "30 / 14 / 7 / 1 days before expiry", "Affected employee + manager", "HIGH"],
    ["Retraining Required", "FAIL result or expired training", "Employee + trainer", "HIGH"],
    ["Session Cancelled", "Session status → CANCELLED", "All assigned attendees", "MEDIUM"],
    ["Result Available", "Results entered for a session", "Session participants", "LOW"],
    ["Certificate Ready", "Certificate issued", "Certificate recipient", "LOW"],
    ["System Alert", "System-wide announcements", "All users / Department / Individual", "URGENT"],
]
add_table(s, Inches(0.5), Inches(2.2), Inches(12.3), notif_data,
          [Inches(2), Inches(3.5), Inches(3), Inches(1.5)])

add_shape(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.7), LIGHT_BLUE, radius=0.04)
add_text(s, Inches(0.7), Inches(5.6), Inches(11.5), Inches(0.3), "Per-User Settings (/notifications)", sz=13, color=NAVY, bold=True)
tb = add_text(s, Inches(0.7), Inches(5.95), Inches(11.5), Inches(1.1), "", sz=10, color=DARK)
tf = tb.text_frame; tf.paragraphs[0].text = ""
add_para(tf, "  Toggle email notifications ON/OFF per user", sz=10, color=DARK, sp_before=Pt(4))
add_para(tf, "  Toggle in-app notifications ON/OFF per user", sz=10, color=DARK, sp_before=Pt(4))
add_para(tf, "  Customize reminder days (e.g., [30, 14, 7, 1] days before expiry)", sz=10, color=DARK, sp_before=Pt(4))
add_para(tf, "  Select which notification types to receive (subscribe/unsubscribe per type)", sz=10, color=DARK, sp_before=Pt(4))


# ════════════════════════════════════════════════════════════
# SLIDE: DATA REGISTRATION SUMMARY
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "Getting Started — Data Registration", "What data needs to be registered before the system is fully operational")

reg = [
    ["#", "Data Category", "Status", "Method", "Required For"],
    ["1", "Employees", "DONE", "HR System Sync", "All modules — base entity for everything"],
    ["2", "Training Programs", "DONE", "Admin creates in /programs", "Sessions, Results, Certificates, Plans"],
    ["3", "Trainers", "REQUESTED", "Excel template → import", "Session scheduling, Result recording"],
    ["4", "Competencies", "REQUESTED", "Excel template → import", "Skill Gap analysis, Learning Paths, IDP"],
    ["5", "TQC Teams", "REQUESTED", "Excel template → import", "New Employee Onboarding module"],
    ["6", "Training Costs", "REQUESTED", "Excel template → import", "ROI calculation, Executive Dashboard"],
]
add_table(s, Inches(0.5), Inches(2.2), Inches(12.3), reg,
          [Inches(0.5), Inches(2), Inches(1.2), Inches(2.5), Inches(6.1)])

add_shape(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.4), LIGHT_GREEN, radius=0.04)
add_text(s, Inches(0.7), Inches(4.9), Inches(11.5), Inches(0.3), "After Registration — Full System Capabilities Unlocked:", sz=13, color=GREEN, bold=True)

tb = add_text(s, Inches(0.7), Inches(5.3), Inches(5.5), Inches(1.6), "", sz=10, color=DARK)
tf = tb.text_frame; tf.paragraphs[0].text = ""
add_para(tf, "  Training: Create sessions → Attendance → Results → Certificates", sz=10, color=DARK, sp_before=Pt(5))
add_para(tf, "  Retraining: Auto-detect FAIL/EXPIRED → Schedule remediation", sz=10, color=DARK, sp_before=Pt(5))
add_para(tf, "  Onboarding: Enroll new hires → 4-stage TQC → Qualification", sz=10, color=DARK, sp_before=Pt(5))
add_para(tf, "  Quality: CAPA process for any training-related issues", sz=10, color=DARK, sp_before=Pt(5))

tb = add_text(s, Inches(6.5), Inches(5.3), Inches(6), Inches(1.6), "", sz=10, color=DARK)
tf = tb.text_frame; tf.paragraphs[0].text = ""
add_para(tf, "  Projects: Manage improvement tasks with Kanban & Gantt", sz=10, color=DARK, sp_before=Pt(5))
add_para(tf, "  Development: Competency maps → Skill gap → Learning paths", sz=10, color=DARK, sp_before=Pt(5))
add_para(tf, "  Analytics: Real-time dashboards, ROI, compliance reports", sz=10, color=DARK, sp_before=Pt(5))
add_para(tf, "  Knowledge: Materials library + Annual training plans", sz=10, color=DARK, sp_before=Pt(5))


# ════════════════════════════════════════════════════════════
# SLIDE: USER ROLES & PERMISSIONS
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "User Roles & Access Control", "Three system roles with granular permission control")

perms = [
    ["Capability", "ADMIN", "TRAINER", "VIEWER"],
    ["View Dashboard & Progress", "YES", "YES", "YES"],
    ["View Programs & Schedule", "YES", "YES", "YES"],
    ["Create/Edit Programs", "YES", "—", "—"],
    ["Create/Edit Sessions", "YES", "YES", "—"],
    ["Record Attendance & Results", "YES", "YES", "—"],
    ["Issue Certificates", "YES", "YES", "—"],
    ["Create/Edit Employees", "YES", "—", "—"],
    ["View Reports & Analytics", "YES", "YES", "YES"],
    ["Access Executive Dashboard", "YES", "—", "—"],
    ["Access Audit Log", "YES", "—", "—"],
    ["Manage Users & Settings", "YES", "—", "—"],
    ["Manage CAPA Process", "YES", "YES", "—"],
    ["Manage Projects", "Separate role matrix (Admin/Member/Guest)", "", ""],
]
add_table(s, Inches(0.5), Inches(2.2), Inches(12.3), perms,
          [Inches(4), Inches(2), Inches(2), Inches(2)])

add_text(s, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
    "Authentication: Google OAuth via Firebase  |  Allowed domains: hwaseung.com, hwaseungvina.com, hsvina.com  |  Language: Vietnamese (default), Korean, English",
    sz=9, color=GRAY)


# ════════════════════════════════════════════════════════════
# SLIDE: COMPLETE SITE MAP
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "Complete Site Map", "30+ pages organized into 8 functional categories")

sitemap = [
    ("Training Management", BLUE, ["/dashboard", "/programs", "/schedule", "/attendance", "/results", "/progress", "/retraining"]),
    ("Certificates & Compliance", GREEN, ["/certificates", "/audit (Compliance)", "/audit-log"]),
    ("New Employee (TQC)", ORANGE, ["/new-tqc/dashboard", "/new-tqc/trainees", "/new-tqc/meetings", "/new-tqc/final-result", "/new-tqc/certificates", "/new-tqc/resignations"]),
    ("Quality (CAPA)", RED, ["/capa (Dashboard)", "/capa/new", "/capa/:id (Detail)"]),
    ("Project Management", PURPLE, ["/projects/dashboard", "/projects/tasks", "/projects/calendar", "/projects/members", "/projects/settings"]),
    ("Employee Development", TEAL, ["/competency", "/skill-gap", "/evaluation", "/employees"]),
    ("Knowledge & Planning", PINK, ["/materials", "/training-plan", "/department"]),
    ("Analytics & Admin", NAVY, ["/executive", "/reports", "/notifications", "/data-sync"]),
]

for i, (cat, color, pages) in enumerate(sitemap):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + col * Inches(3.1)
    y = Inches(2.2) + row * Inches(2.5)

    add_shape(s, x, y, Inches(2.9), Inches(2.3), WHITE, radius=0.04)
    add_shape(s, x, y, Inches(0.06), Inches(2.3), color)
    add_text(s, x + Inches(0.2), y + Inches(0.08), Inches(2.5), Inches(0.25), cat, sz=11, color=color, bold=True)
    add_shape(s, x + Inches(0.2), y + Inches(0.35), Inches(2.5), Inches(0.015), color)

    tb = add_text(s, x + Inches(0.2), y + Inches(0.45), Inches(2.5), Inches(1.7), "", sz=9, color=DARK)
    tf = tb.text_frame; tf.paragraphs[0].text = ""
    for pg in pages:
        add_para(tf, pg, sz=9, color=DARK, sp_before=Pt(2))


# ════════════════════════════════════════════════════════════
# SLIDE: THANK YOU
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, NAVY)
add_shape(s, Inches(0), Inches(0), W, Inches(0.08), BLUE)
add_text(s, Inches(1), Inches(1.8), Inches(11), Inches(0.8), "Q-TRAIN", sz=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.6), Inches(11), Inches(0.5), "Complete System Guide", sz=24, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)
add_shape(s, Inches(5.5), Inches(3.3), Inches(2.3), Inches(0.04), BLUE)
add_text(s, Inches(1), Inches(3.7), Inches(11), Inches(0.4), "Thank You", sz=28, color=RGBColor(0xBF,0xDB,0xFE), align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.5), Inches(11), Inches(0.3), "https://q-train-web.web.app", sz=14, color=RGBColor(0x60,0xA5,0xFA), align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
    "HWK Vietnam | QIP Training Management\nContact: ksmoon@hsvina.com",
    sz=12, color=RGBColor(0x64,0x78,0x96), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# Save
# ════════════════════════════════════════
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Q-TRAIN_Complete_Guide_EN.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
