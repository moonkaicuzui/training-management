#!/usr/bin/env python3
"""Generate Q-TRAIN User Roles & System Admin Guide PPTX"""

import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx_helpers import (
    NAVY, BLUE, LIGHT_BLUE, WHITE, DARK, GRAY, LIGHT_GRAY, GREEN, ORANGE, RED, ACCENT,
    add_bg, add_shape, add_text, add_para, add_card, slide_header,
    add_table_slide, create_presentation, title_slide, thank_you_slide,
)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUT_FILE = os.path.join(SCRIPT_DIR, "Q-TRAIN_User_Roles_Guide.pptx")

TEAL = RGBColor(0x06, 0x94, 0xA2)
PINK = RGBColor(0xDB, 0x27, 0x77)


def section_divider(prs, number, title, subtitle):
    """Dark divider slide between sections"""
    W = prs.slide_width
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), BLUE)

    add_text(slide, Inches(1), Inches(2.0), Inches(1.5), Inches(1.5),
             number, size=72, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(3), Inches(2.2), Inches(9), Inches(0.8),
             title, size=36, color=WHITE, bold=True)
    add_text(slide, Inches(3), Inches(3.1), Inches(9), Inches(0.5),
             subtitle, size=16, color=RGBColor(0x93, 0xC5, 0xFD))
    add_shape(slide, Inches(3), Inches(3.8), Inches(2), Inches(0.04), BLUE)
    return slide


def role_card(slide, x, y, w, h, icon, role_name, role_desc, color, examples):
    """Render a role card with icon, name, description, examples"""
    add_shape(slide, x, y, w, h, LIGHT_GRAY, radius=0.04)
    add_shape(slide, x, y, w, Inches(0.06), color)
    # Icon
    add_shape(slide, x + Inches(0.2), y + Inches(0.25), Inches(0.7), Inches(0.7), color, radius=0.15)
    add_text(slide, x + Inches(0.2), y + Inches(0.3), Inches(0.7), Inches(0.6),
             icon, size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Name & desc
    add_text(slide, x + Inches(1.1), y + Inches(0.2), w - Inches(1.4), Inches(0.35),
             role_name, size=18, color=color, bold=True)
    add_text(slide, x + Inches(1.1), y + Inches(0.55), w - Inches(1.4), Inches(0.3),
             role_desc, size=12, color=GRAY)
    # Examples
    for i, ex in enumerate(examples):
        add_text(slide, x + Inches(0.3), y + Inches(1.1) + i * Inches(0.3), w - Inches(0.6), Inches(0.3),
                 f"  {ex}", size=11, color=DARK)


def checklist_card(slide, x, y, w, h, title, items, color):
    """Render a checklist card"""
    add_shape(slide, x, y, w, h, WHITE, radius=0.04)
    add_shape(slide, x, y, Inches(0.06), h, color)
    add_text(slide, x + Inches(0.25), y + Inches(0.12), w - Inches(0.5), Inches(0.3),
             title, size=14, color=color, bold=True)
    for i, item in enumerate(items):
        add_text(slide, x + Inches(0.25), y + Inches(0.5) + i * Inches(0.28), w - Inches(0.5), Inches(0.28),
                 item, size=11, color=DARK)


def build_all(prs):
    W = prs.slide_width

    # ─── SLIDE 1: Title ───
    title_slide(prs,
        "QIP 교육 관리 시스템",
        "사용자 역할 구성 & 시스템 관리 가이드",
        "HWK Vietnam (화승비나) | QIP Training Management",
        "https://q-train-web.web.app/",
        "Q-TRAIN v2.0 | 2026.02")

    # ─── SLIDE 2: 목차 ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "목차", "사용자 역할 구성 & 시스템 관리 가이드")

    toc = [
        ("1", "사용자 역할 개요", "3가지 역할: ADMIN / TRAINER / VIEWER", BLUE),
        ("2", "역할별 권한 매트릭스", "각 역할이 할 수 있는 것과 할 수 없는 것", GREEN),
        ("3", "ADMIN — 시스템 마스터", "초기 셋업 + 일상 관리 + 데이터 무결성", ORANGE),
        ("4", "TRAINER — 교육 담당자", "교육 실시, 결과 입력, 일정 관리", ACCENT),
        ("5", "VIEWER — 조회자", "대시보드, 리포트, 교육 현황 조회", TEAL),
        ("6", "사용자 등록 & 보안 정책", "이메일 도메인, 역할 배정, 감사 로그", RED),
    ]
    for i, (num, t, d, c) in enumerate(toc):
        row = i // 2
        col = i % 2
        x = Inches(0.5) + col * Inches(6.2)
        y = Inches(2.3) + row * Inches(1.5)
        add_shape(slide, x, y, Inches(5.9), Inches(1.3), LIGHT_GRAY, radius=0.04)
        add_shape(slide, x + Inches(0.15), y + Inches(0.25), Inches(0.6), Inches(0.6), c, radius=0.1)
        add_text(slide, x + Inches(0.15), y + Inches(0.3), Inches(0.6), Inches(0.5),
                 num, size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.9), y + Inches(0.2), Inches(4.7), Inches(0.35),
                 t, size=16, color=DARK, bold=True)
        add_text(slide, x + Inches(0.9), y + Inches(0.6), Inches(4.7), Inches(0.5),
                 d, size=12, color=GRAY)

    # ═══════════════════════════════════════════════════
    # SECTION 1: 사용자 역할 개요
    # ═══════════════════════════════════════════════════
    section_divider(prs, "1", "사용자 역할 개요", "Q-TRAIN 시스템의 3가지 사용자 유형")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "3가지 사용자 역할", "Q-TRAIN은 역할 기반 접근 제어(RBAC)를 사용합니다")

    role_card(slide, Inches(0.5), Inches(2.3), Inches(3.9), Inches(4.5),
        "A", "ADMIN (관리자)", "시스템 마스터 · 전체 권한",
        BLUE,
        ["  QA 매니저, 시스템 관리자",
         "  프로그램/직원/사용자 관리",
         "  감사 로그, 데이터 동기화",
         "  경영진 대시보드 접근",
         "  모든 데이터 편집/삭제 권한",
         "  시스템 초기 셋업 담당",
         "  예: ksmoon@hsvina.com"])

    role_card(slide, Inches(4.7), Inches(2.3), Inches(3.9), Inches(4.5),
        "T", "TRAINER (교육 담당자)", "교육 실시 · 결과 입력",
        GREEN,
        ["  라인 리더, 현장 강사, QC 담당",
         "  교육 세션 생성/편집",
         "  출석 체크, 점수 입력",
         "  교육 일정 관리",
         "  재교육 세션 생성",
         "  프로그램 편집 불가",
         "  예: 현장 교육 강사"])

    role_card(slide, Inches(8.9), Inches(2.3), Inches(3.9), Inches(4.5),
        "V", "VIEWER (조회자)", "읽기 전용 · 조회만 가능",
        TEAL,
        ["  경영진, HR, 타 부서",
         "  대시보드 조회",
         "  교육 현황/결과 열람",
         "  리포트 다운로드",
         "  모든 데이터 편집 불가",
         "  사용자 관리 불가",
         "  예: 경영지원팀, HR"])

    # ═══════════════════════════════════════════════════
    # SECTION 2: 권한 매트릭스
    # ═══════════════════════════════════════════════════
    section_divider(prs, "2", "역할별 권한 매트릭스", "각 역할이 할 수 있는 것과 할 수 없는 것")

    # Permission matrix table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "기능별 권한 매트릭스", "O = 가능 / X = 불가 / △ = 제한적")

    rows = [
        ["기능 영역", "세부 기능", "ADMIN", "TRAINER", "VIEWER"],
        ["대시보드", "메인/경영진/부서별 대시보드 조회", "O", "O", "O"],
        ["교육 프로그램", "프로그램 생성/편집/비활성화", "O", "X", "X"],
        ["교육 프로그램", "프로그램 조회", "O", "O", "O"],
        ["교육 세션", "세션 생성/편집/삭제", "O", "O", "X"],
        ["출석 관리", "출석 체크/수정", "O", "O", "X"],
        ["교육 결과", "점수/등급 입력/수정", "O", "O", "X"],
        ["교육 결과", "결과 삭제", "X", "X", "X"],
        ["직원 관리", "직원 정보 생성/편집", "O", "X", "X"],
        ["재교육", "재교육 세션 생성", "O", "O", "X"],
        ["자격증", "자격증 발급/관리", "O", "O", "X"],
        ["신입 TQC", "교육생 등록/관리", "O", "O", "X"],
        ["CAPA", "CAPA 등록/관리", "O", "O", "X"],
        ["5PRS 분석", "분석 실행/교육 추천", "O", "O", "O"],
        ["리포트", "리포트 조회/다운로드", "O", "O", "O"],
        ["사용자 관리", "역할 배정/사용자 관리", "O", "X", "X"],
        ["감사 로그", "변경 이력 조회", "O", "X", "X"],
        ["데이터 동기화", "Google Sheets 동기화", "O", "X", "X"],
        ["시스템 설정", "TQC 팀/교육 단계 설정", "O", "X", "X"],
    ]

    add_table_slide(slide, Inches(0.5), Inches(2.2), Inches(12.3), rows,
                    [Inches(2.0), Inches(4.5), Inches(1.8), Inches(1.8), Inches(1.8)])

    # ═══════════════════════════════════════════════════
    # SECTION 3: ADMIN — 시스템 마스터
    # ═══════════════════════════════════════════════════
    section_divider(prs, "3", "ADMIN — 시스템 마스터", "초기 셋업, 일상 관리, 데이터 무결성 책임")

    # 3-1: Initial Setup
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "ADMIN 초기 셋업 체크리스트", "시스템 도입 시 반드시 완료해야 할 작업")

    setup_groups = [
        ("1단계: 기본 데이터 등록", BLUE, [
            "직원 DB 등록 (이름, 부서, 직급, 빌딩, 라인, 입사일)",
            "부서/빌딩/라인 구조 확인 및 정리",
            "강사 정보 등록 (자격, 전문 분야)",
            "허용 이메일 도메인 확인 (@hsvina.com 등)",
        ]),
        ("2단계: 교육 프로그램 설정", GREEN, [
            "교육 프로그램 코드 체계 수립 (예: SAF-001)",
            "각 프로그램별 평가 기준 설정 (합격 점수, 등급)",
            "유효기간 및 만료 알림 주기 설정",
            "대상 직급/부서별 필수 교육 매핑",
        ]),
        ("3단계: TQC & 역량 설정", ORANGE, [
            "TQC 팀 생성 및 리더 배정",
            "4단계 교육 과정 정의",
            "역량 항목 정의 (부서/직급별 필요 역량)",
            "역량 목표 레벨 설정",
        ]),
        ("4단계: 품질 연동 설정", ACCENT, [
            "5PRS 데이터 소스 연결 (Google Sheets)",
            "불량 유형 → 교육 프로그램 매핑 설정",
            "CAPA 프로세스 카테고리 정의",
            "데이터 동기화 스케줄 확인",
        ]),
    ]

    for i, (gtitle, gcolor, items) in enumerate(setup_groups):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.2)
        y = Inches(2.3) + row * Inches(2.45)
        checklist_card(slide, x, y, Inches(5.9), Inches(2.2), gtitle, items, gcolor)

    # 3-2: Daily Admin Tasks
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "ADMIN 일상 관리 업무", "시스템 운영 중 지속적으로 수행해야 할 작업")

    daily_groups = [
        ("일일 업무", BLUE, [
            "대시보드 확인 (이수율, 합격률, 재교육 현황)",
            "알림 확인 (자격증 만료, 미완료 교육)",
            "신규 등록 요청 처리 (직원, 프로그램)",
        ]),
        ("주간 업무", GREEN, [
            "교육 일정 확인 및 조정",
            "주간 교육 추천 결과 검토 (자동 분석)",
            "불합격자 재교육 현황 점검",
            "신입 TQC 교육 진행 상황 확인",
        ]),
        ("월간 업무", ORANGE, [
            "경영진 보고용 리포트 생성",
            "부서별 교육 완료율 분석",
            "역량 매트릭스 업데이트 확인",
            "5PRS 월간 분석 실행 & 교육 추천",
        ]),
        ("분기/연간 업무", ACCENT, [
            "교육 프로그램 유효성 검토 및 갱신",
            "역량 평가 실시 (스킬 갭 분석)",
            "감사 로그 검토 및 규정 준수 확인",
            "시스템 데이터 정합성 점검 (동기화)",
        ]),
    ]

    for i, (gtitle, gcolor, items) in enumerate(daily_groups):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.2)
        y = Inches(2.3) + row * Inches(2.45)
        checklist_card(slide, x, y, Inches(5.9), Inches(2.2), gtitle, items, gcolor)

    # 3-3: Admin-only pages
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "ADMIN 전용 페이지", "ADMIN만 접근할 수 있는 시스템 관리 메뉴")

    admin_pages = [
        ("경영진 대시보드", "전사 KPI, 부서별 비교, 추이 분석", "사이드바 → 대시보드 → 경영진", BLUE),
        ("감사 & 규정 준수", "규정 준수 현황, 감사 체크리스트", "사이드바 → 시스템 관리 → 감사", GREEN),
        ("감사 로그", "모든 데이터 변경 이력 (누가/언제/무엇)", "사이드바 → 시스템 관리 → 감사 로그", ORANGE),
        ("데이터 동기화", "Google Sheets ↔ Firestore 동기화", "사이드바 → 시스템 관리 → 데이터 동기화", ACCENT),
        ("사용자 관리", "사용자 역할 배정, 접근 권한 관리", "Firebase Console → Authentication", RED),
        ("TQC 설정", "TQC 팀 관리, 교육 단계 설정", "사이드바 → 신입 TQC → 설정", TEAL),
    ]

    for i, (ptitle, pdesc, ploc, pcolor) in enumerate(admin_pages):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col * Inches(4.15)
        y = Inches(2.3) + row * Inches(2.4)
        add_shape(slide, x, y, Inches(3.9), Inches(2.1), LIGHT_GRAY, radius=0.04)
        add_shape(slide, x, y, Inches(3.9), Inches(0.06), pcolor)
        add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.5), Inches(0.3),
                 ptitle, size=15, color=pcolor, bold=True)
        add_text(slide, x + Inches(0.2), y + Inches(0.55), Inches(3.5), Inches(0.6),
                 pdesc, size=12, color=DARK)
        add_shape(slide, x + Inches(0.2), y + Inches(1.25), Inches(3.5), Inches(0.55),
                  RGBColor(0xEF, 0xF6, 0xFF), radius=0.03)
        add_text(slide, x + Inches(0.3), y + Inches(1.35), Inches(3.3), Inches(0.35),
                 ploc, size=10, color=BLUE)

    # ═══════════════════════════════════════════════════
    # SECTION 4: TRAINER — 교육 담당자
    # ═══════════════════════════════════════════════════
    section_divider(prs, "4", "TRAINER — 교육 담당자", "교육 실시, 결과 입력, 일정 관리")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "TRAINER 역할 & 업무", "교육 현장에서 직접 교육을 운영하는 담당자")

    # Who is TRAINER
    add_shape(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.0),
              RGBColor(0xEC, 0xFD, 0xF5), radius=0.03)
    add_text(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.3),
             "TRAINER 대상: 라인 리더, 현장 강사, QC 교육 담당자, TQC 팀 리더",
             size=15, color=GREEN, bold=True)
    add_text(slide, Inches(0.8), Inches(2.75), Inches(11.5), Inches(0.3),
             "자동 배정: 허용 도메인(@hsvina.com, @hwaseung.com 등)으로 로그인하면 TRAINER 역할 자동 부여",
             size=12, color=DARK)

    trainer_tasks = [
        ("교육 세션 관리", GREEN, [
            "새 교육 세션 생성 (날짜, 장소, 강사)",
            "참석자 등록 (개별/일괄)",
            "세션 일정 변경/취소",
            "세션 상태 관리 (예정/진행/완료)",
        ]),
        ("출석 & 결과 입력", BLUE, [
            "교육 당일 출석 체크",
            "참석/불참/지각 기록",
            "점수 및 등급 입력",
            "합격/불합격 자동 판정 확인",
        ]),
        ("재교육 관리", ORANGE, [
            "불합격자 재교육 세션 생성",
            "재교육 출석 및 결과 입력",
            "재평가 후 합격 처리",
        ]),
        ("신입 TQC 교육", ACCENT, [
            "신입 교육생 등록",
            "4단계 교육 진행 기록",
            "면담 기록 입력",
            "최종 평가 입력",
        ]),
    ]

    for i, (ttitle, tcolor, items) in enumerate(trainer_tasks):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.2)
        y = Inches(3.6) + row * Inches(1.85)
        checklist_card(slide, x, y, Inches(5.9), Inches(1.65), ttitle, items, tcolor)

    # TRAINER cannot do
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "TRAINER 할 수 없는 것", "ADMIN에게 요청해야 하는 작업")

    cannot_items = [
        ("교육 프로그램 생성/편집", "새 프로그램이 필요하면 ADMIN에게 요청", "ADMIN이 프로그램 등록 후 TRAINER가 세션 생성"),
        ("직원 정보 편집", "직원 추가/수정이 필요하면 ADMIN에게 요청", "ADMIN이 직원 DB 관리"),
        ("교육 결과 삭제", "잘못된 결과는 수정만 가능 (삭제 불가)", "수정 시 편집 로그가 자동 기록됨"),
        ("감사 로그 조회", "변경 이력 확인이 필요하면 ADMIN에게 요청", "데이터 무결성 보호를 위한 정책"),
        ("사용자 관리", "새 사용자 역할 변경이 필요하면 ADMIN에게 요청", "ADMIN만 역할 배정 가능"),
        ("데이터 동기화", "외부 데이터 동기화 실행 권한 없음", "ADMIN이 스케줄에 따라 동기화 관리"),
    ]

    for i, (ctitle, cdesc, cnote) in enumerate(cannot_items):
        row = i // 2
        col = i % 2
        x = Inches(0.5) + col * Inches(6.2)
        y = Inches(2.3) + row * Inches(1.55)
        add_shape(slide, x, y, Inches(5.9), Inches(1.35), RGBColor(0xFE, 0xF2, 0xF2), radius=0.03)
        add_text(slide, x + Inches(0.2), y + Inches(0.12), Inches(5.5), Inches(0.3),
                 f"X  {ctitle}", size=14, color=RED, bold=True)
        add_text(slide, x + Inches(0.2), y + Inches(0.45), Inches(5.5), Inches(0.3),
                 cdesc, size=12, color=DARK)
        add_text(slide, x + Inches(0.2), y + Inches(0.8), Inches(5.5), Inches(0.4),
                 f"→ {cnote}", size=11, color=GRAY)

    # ═══════════════════════════════════════════════════
    # SECTION 5: VIEWER — 조회자
    # ═══════════════════════════════════════════════════
    section_divider(prs, "5", "VIEWER — 조회자", "읽기 전용 접근, 대시보드 & 리포트 조회")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "VIEWER 역할 & 활용법", "경영진, HR, 타 부서 담당자를 위한 조회 기능")

    add_shape(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.8),
              RGBColor(0xEC, 0xFE, 0xFF), radius=0.03)
    add_text(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.3),
             "VIEWER 대상: 공장장, 경영진, HR 담당, 품질 관리 부서, 타 부서 관리자",
             size=15, color=TEAL, bold=True)
    add_text(slide, Inches(0.8), Inches(2.7), Inches(11.5), Inches(0.3),
             "모든 데이터를 조회할 수 있지만, 생성/편집/삭제는 할 수 없습니다",
             size=12, color=DARK)

    viewer_tasks = [
        ("대시보드 조회", TEAL, [
            "메인 대시보드 (전체 현황)",
            "부서별 대시보드 (부서 비교)",
            "교육 진도 현황 (프로그램별)",
        ]),
        ("교육 현황 열람", BLUE, [
            "교육 프로그램 목록 조회",
            "교육 일정 확인",
            "교육 결과/합격률 조회",
            "재교육 현황 확인",
        ]),
        ("리포트 활용", GREEN, [
            "리포트 조회 및 다운로드 (Excel/PDF)",
            "5PRS 분석 결과 조회",
            "역량 매트릭스/스킬 갭 조회",
        ]),
        ("인력 현황 조회", ORANGE, [
            "직원 목록 및 프로필 조회",
            "강사 정보 조회",
            "자격증 현황 확인",
            "신입 TQC 교육 진행 확인",
        ]),
    ]

    for i, (vtitle, vcolor, items) in enumerate(viewer_tasks):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.2)
        y = Inches(3.4) + row * Inches(1.85)
        checklist_card(slide, x, y, Inches(5.9), Inches(1.65), vtitle, items, vcolor)

    # ═══════════════════════════════════════════════════
    # SECTION 6: 사용자 등록 & 보안
    # ═══════════════════════════════════════════════════
    section_divider(prs, "6", "사용자 등록 & 보안 정책", "이메일 도메인, 역할 배정, 감사 로그")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "사용자 등록 프로세스", "이메일 기반 자동 역할 배정")

    # Flow diagram using shapes
    flow_steps = [
        ("1. 접속", "q-train-web.web.app\n접속", BLUE),
        ("2. 로그인", "회사 이메일로\nGoogle 로그인", GREEN),
        ("3. 도메인 검증", "@hsvina.com\n@hwaseung.com\n허용 도메인 확인", ORANGE),
        ("4. 역할 자동 배정", "ADMIN 이메일 목록\n→ ADMIN\n그 외 → TRAINER", ACCENT),
        ("5. 시스템 사용", "역할에 맞는\n메뉴/기능 활성화", TEAL),
    ]

    for i, (ftitle, fdesc, fcolor) in enumerate(flow_steps):
        x = Inches(0.3) + i * Inches(2.55)
        y = Inches(2.3)
        add_shape(slide, x, y, Inches(2.3), Inches(2.0), LIGHT_GRAY, radius=0.04)
        add_shape(slide, x, y, Inches(2.3), Inches(0.06), fcolor)
        add_text(slide, x + Inches(0.1), y + Inches(0.2), Inches(2.1), Inches(0.3),
                 ftitle, size=13, color=fcolor, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.1), y + Inches(0.6), Inches(2.1), Inches(1.2),
                 fdesc, size=11, color=DARK, align=PP_ALIGN.CENTER)
        if i < 4:
            add_text(slide, x + Inches(2.3), y + Inches(0.8), Inches(0.25), Inches(0.3),
                     "→", size=18, color=GRAY, bold=True)

    # Security policies
    add_shape(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.6),
              RGBColor(0xFE, 0xF3, 0xC7), radius=0.03)
    add_text(slide, Inches(0.8), Inches(4.7), Inches(11.5), Inches(0.3),
             "보안 정책 & 데이터 보호", size=16, color=ORANGE, bold=True)

    policies = [
        ("허용 이메일 도메인", "@hwaseung.com, @hwaseungvina.com, @hsvina.com (+ @gmail.com 개발용)"),
        ("ADMIN 이메일", "admin@hwaseung.com, qip.admin@hwaseungvina.com, ksmoon@hsvina.com"),
        ("교육 결과 보호", "삭제 불가 (NO DELETE) — 수정 시 편집 로그 자동 기록"),
        ("감사 로그", "모든 변경 사항 자동 기록 (추가 전용, 수정/삭제 불가)"),
        ("자격증/이수 기록", "삭제 불가 — 규정 준수 및 감사 추적을 위해 영구 보존"),
        ("업데이트 제한", "1초 미만 간격의 연속 업데이트 차단 (Rate Limiting)"),
    ]

    for i, (ptitle, pdesc) in enumerate(policies):
        row = i // 2
        col = i % 2
        x = Inches(0.8) + col * Inches(6.0)
        y = Inches(5.15) + row * Inches(0.6)
        add_text(slide, x, y, Inches(5.5), Inches(0.3),
                 f"  {ptitle}: {pdesc}", size=11, color=DARK)

    # ─── SUMMARY SLIDE ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "역할별 요약 비교", "한눈에 보는 3가지 역할 차이")

    summary_rows = [
        ["구분", "ADMIN (관리자)", "TRAINER (교육 담당)", "VIEWER (조회자)"],
        ["대상", "QA 매니저, 시스템 관리자", "라인 리더, 현장 강사", "경영진, HR, 타 부서"],
        ["교육 프로그램", "생성/편집/비활성화", "조회만", "조회만"],
        ["교육 세션", "생성/편집/삭제", "생성/편집", "조회만"],
        ["출석 & 결과", "입력/수정", "입력/수정", "조회만"],
        ["직원 관리", "생성/편집/삭제", "조회만", "조회만"],
        ["사용자 관리", "역할 배정/관리", "불가", "불가"],
        ["감사 로그", "조회 가능", "불가", "불가"],
        ["데이터 동기화", "실행 가능", "불가", "불가"],
        ["경영진 대시보드", "접근 가능", "불가", "불가"],
        ["리포트 다운로드", "가능", "가능", "가능"],
        ["초기 셋업", "담당", "해당 없음", "해당 없음"],
    ]

    add_table_slide(slide, Inches(0.5), Inches(2.2), Inches(12.3), summary_rows,
                    [Inches(2.5), Inches(3.5), Inches(3.0), Inches(3.0)])

    # ─── Thank You ───
    thank_you_slide(prs, "https://q-train-web.web.app/", "ksmoon@hsvina.com")


def main():
    prs, W, H = create_presentation()
    build_all(prs)
    prs.save(OUT_FILE)
    count = len(prs.slides)
    size = os.path.getsize(OUT_FILE)
    print(f"Generated: {os.path.basename(OUT_FILE)} — {count} slides, {size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
