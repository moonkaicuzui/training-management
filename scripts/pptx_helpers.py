#!/usr/bin/env python3
"""
Q-TRAIN PPTX Helpers — Shared helper functions for all PPTX generators
Brand colors, shapes, text, tables, and standard slide layouts.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Colors ──
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
BLUE = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
GREEN = RGBColor(0x05, 0x96, 0x69)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT = RGBColor(0x7C, 0x3A, 0xED)


def add_bg(slide, color):
    """Set solid background color"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, radius=None):
    """Add a colored rectangle shape"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if hasattr(shape, 'adjustments') and radius:
        try:
            shape.adjustments[0] = radius
        except:
            pass
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font_name='Arial'):
    """Add a text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_para(text_frame, text, size=16, color=DARK, bold=False, space_before=Pt(4), space_after=Pt(4), align=PP_ALIGN.LEFT, bullet=False):
    """Add a paragraph to existing text frame"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Arial'
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    if bullet:
        p.level = 0
    return p


def add_card(slide, left, top, width, height, title, value, subtitle="", bg_color=WHITE, title_color=GRAY, value_color=DARK, icon=""):
    """Add a KPI card"""
    card = add_shape(slide, left, top, width, height, bg_color, radius=0.05)
    card.shadow.inherit = False

    if icon:
        add_text(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.5), Inches(0.4),
                 icon, size=22, color=title_color)

    add_text(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.3),
             title, size=11, color=title_color)
    add_text(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), Inches(0.5),
             value, size=28, color=value_color, bold=True)
    if subtitle:
        add_text(slide, left + Inches(0.2), top + height - Inches(0.35), width - Inches(0.4), Inches(0.3),
                 subtitle, size=10, color=GRAY)
    return card


def slide_header(slide, title, subtitle="", bg_color=WHITE):
    """Add consistent header bar"""
    add_bg(slide, bg_color)
    # Top bar
    try:
        W = slide.part.package.presentation_part.presentation.slide_width
    except (AttributeError, TypeError):
        W = Inches(13.333)
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), BLUE)
    # Logo area
    add_text(slide, Inches(0.5), Inches(0.25), Inches(3), Inches(0.4),
             "Q-TRAIN", size=14, color=BLUE, bold=True)
    add_text(slide, Inches(0.5), Inches(0.55), Inches(5), Inches(0.3),
             "HWK Vietnam QIP Training Management System", size=9, color=GRAY)
    # Title
    add_text(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
             title, size=28, color=NAVY, bold=True)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.65), Inches(12), Inches(0.4),
                 subtitle, size=14, color=GRAY)
    # Divider
    add_shape(slide, Inches(0.5), Inches(2.05), Inches(1.5), Inches(0.04), BLUE)


def add_table_slide(slide, left, top, width, rows_data, col_widths, header_color=NAVY):
    """Add a styled table"""
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.4 * rows))
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for r, row_data in enumerate(rows_data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = 'Arial'
                if r == 0:
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.bold = True
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    return table_shape


def create_presentation():
    """Create a standard Q-TRAIN presentation with 13.333x7.5 dimensions.
    Returns (prs, W, H) tuple."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height
    return prs, W, H


def title_slide(prs, title, subtitle, info_line, url, version):
    """Create standard title slide"""
    W = prs.slide_width
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide, NAVY)

    # Accent shapes
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), BLUE)
    add_shape(slide, Inches(0), Inches(7.0), W, Inches(0.5), RGBColor(0x15, 0x2A, 0x4A))

    # Main title
    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "Q-TRAIN", size=60, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
    add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6),
             title, size=28, color=RGBColor(0x93, 0xC5, 0xFD))

    # Divider
    add_shape(slide, Inches(1), Inches(3.3), Inches(2), Inches(0.05), BLUE)

    # Subtitle
    add_text(slide, Inches(1), Inches(3.6), Inches(10), Inches(0.5),
             subtitle, size=20, color=RGBColor(0xBF, 0xDB, 0xFE))

    # Info boxes
    add_text(slide, Inches(1), Inches(4.5), Inches(5), Inches(0.3),
             info_line, size=14, color=RGBColor(0x93, 0xC5, 0xFD))
    add_text(slide, Inches(1), Inches(4.9), Inches(5), Inches(0.3),
             url, size=13, color=RGBColor(0x60, 0xA5, 0xFA))

    # Bottom
    add_text(slide, Inches(1), Inches(6.5), Inches(5), Inches(0.3),
             version, size=11, color=RGBColor(0x64, 0x78, 0x96))

    return slide


def thank_you_slide(prs, url, contact_email):
    """Create thank you slide"""
    W = prs.slide_width
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), BLUE)

    add_text(slide, Inches(1), Inches(2), Inches(11), Inches(0.8),
             "Q-TRAIN", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.5),
             "QIP Training Management System", size=22, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)

    add_shape(slide, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.04), BLUE)

    add_text(slide, Inches(1), Inches(3.9), Inches(11), Inches(0.4),
             "Thank You", size=28, color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.3),
             url, size=14, color=RGBColor(0x60, 0xA5, 0xFA), align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(5.4), Inches(11), Inches(0.5),
             f"HWK Vietnam | QIP Training Management\nContact: {contact_email}",
             size=12, color=RGBColor(0x64, 0x78, 0x96), align=PP_ALIGN.CENTER)

    return slide
