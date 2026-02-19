#!/usr/bin/env python3
"""Generate Q-TRAIN Scenario-Based Training Guide in 3 languages (KO, EN, VI)"""

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx_helpers import (
    NAVY, BLUE, LIGHT_BLUE, WHITE, DARK, GRAY, LIGHT_GRAY, GREEN, ORANGE, RED, ACCENT,
    add_bg, add_shape, add_text, add_para, add_card, slide_header,
    add_table_slide, create_presentation, title_slide, thank_you_slide,
)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

# ═══════════════════════════════════════════════════════════════
# TEXT CONTENT — 3 languages
# ═══════════════════════════════════════════════════════════════

TEXTS = {
    "ko": {
        "file": "Q-TRAIN_Scenario_Training_KO.pptx",
        "title": "QIP 교육 관리 시스템",
        "subtitle": "시나리오 기반 실무 교육 가이드",
        "info": "HWK Vietnam (화승비나) | QIP Training Management",
        "version": "Q-TRAIN v2.0 | 2026.02",
        # ── Table of Contents ──
        "toc_title": "교육 과정 안내",
        "toc_sub": "8가지 실무 시나리오로 배우는 Q-TRAIN 시스템",
        "toc_items": [
            ("A", "출근 후 일일 업무", "대시보드 확인, 알림, 오늘 일정 점검"),
            ("B", "신규 교육 프로그램 등록", "프로그램 생성 → 평가 기준 → 유효기간 설정"),
            ("C", "교육 세션 실시 & 결과 입력", "세션 생성 → 출석 → 점수 → 합격/불합격"),
            ("D", "불합격자 재교육 처리", "재교육 대상 확인 → 재교육 → 재평가"),
            ("E", "신입 검사원 교육 (TQC)", "팀 배정 → 4단계 교육 → 면담 → 최종 평가"),
            ("F", "품질 문제 대응 (CAPA & 5PRS)", "불량 분석 → 교육 추천 → 자동 등록"),
            ("G", "역량 매트릭스 관리", "역량 정의 → 직원 할당 → 갭 분석 → 학습 경로"),
            ("H", "경영진 보고 & 리포트", "KPI 조회 → 리포트 생성 → 공유"),
        ],
        # ── System Overview ──
        "overview_title": "Q-TRAIN 시스템 소개",
        "overview_sub": "QIP 부서장을 위한 교육 관리 플랫폼",
        "overview_what": "Q-TRAIN이란?",
        "overview_desc": "HWK Vietnam QIP(Quality Improvement Program) 부서의\n교육 운영, 역량 관리, 품질 개선을 통합 관리하는 웹 시스템",
        "overview_features": [
            ("교육 운영", "프로그램, 일정, 출석, 결과를 한곳에서 관리"),
            ("인력 & 역량", "직원 DB, 역량 매트릭스, 스킬 갭 분석"),
            ("품질 개선", "CAPA 프로세스, 5PRS 불량 분석, 교육 추천"),
            ("신입 TQC", "4단계 교육, 면담 기록, 최종 평가"),
            ("분석 & 리포트", "대시보드, KPI, 경영진 보고서"),
            ("시스템 관리", "감사 로그, 자료실, 데이터 동기화"),
        ],
        "overview_sidebar": "사이드바 메뉴 구성 (9개 카테고리)",
        "sidebar_cats": [
            "대시보드", "교육 운영", "후속 관리", "인력 & 역량",
            "신입 TQC", "프로젝트", "품질 개선", "분석 & 리포트", "시스템 관리",
        ],
        "overview_url": "접속 주소: https://q-train-web.web.app/",
        "overview_lang": "지원 언어: 한국어 / 영어 / 베트남어",
        # ── Scenarios ──
        "scenarios": [
            {
                "id": "A",
                "title": "시나리오 A: 출근 후 일일 업무",
                "sub": "매일 아침 5분, Q-TRAIN으로 교육 현황 파악",
                "persona": "문경서 QA 매니저",
                "situation": "월요일 오전 8시, 문경서 매니저가 출근하여\n이번 주 교육 현황과 주요 알림을 확인합니다.",
                "steps": [
                    ("1. 로그인 & 대시보드", "Q-TRAIN 접속 → 대시보드 자동 표시\n이번 달 교육 이수율, 합격률, 재교육 필요 인원 확인"),
                    ("2. 알림 확인", "상단 알림 아이콘 클릭\n자격증 만료 예정, 미완료 교육, 재교육 대상자 알림"),
                    ("3. 오늘 일정 확인", "사이드바 → 교육 운영 → 일정\n오늘/이번 주 예정된 교육 세션 확인"),
                    ("4. 경영진 대시보드", "사이드바 → 대시보드 → 경영진\n부서별 교육 현황, KPI 추이 그래프 확인"),
                ],
            },
            {
                "id": "B",
                "title": "시나리오 B: 신규 교육 프로그램 등록",
                "sub": "새로운 품질 검사 교육을 시스템에 등록",
                "persona": "Nguyen Thi Mai (QC 검사원 교육 담당)",
                "situation": "새로운 '솔 접착 품질 검사' 교육 프로그램이 필요합니다.\nMai 담당자가 프로그램을 등록합니다.",
                "steps": [
                    ("1. 프로그램 생성", "사이드바 → 교육 운영 → 프로그램\n[+ 새 프로그램] 버튼 → 프로그램명, 코드, 카테고리 입력"),
                    ("2. 평가 기준 설정", "합격 점수: 80점 / 평가 방식: 필기 + 실기\n등급 기준: A(90+), B(80+), C(70+), F(<70)"),
                    ("3. 대상자 지정", "대상 직급: QC 검사원, 라인 리더\n대상 부서: Building A, B 품질팀"),
                    ("4. 유효기간 설정", "교육 유효기간: 12개월\n만료 전 30일/14일/7일 알림 자동 발송"),
                ],
            },
            {
                "id": "C",
                "title": "시나리오 C: 교육 세션 실시 & 결과 입력",
                "sub": "월간 안전교육 세션을 진행하고 결과를 기록",
                "persona": "Tran Van Duc (라인 리더, 현장 강사)",
                "situation": "이번 달 안전교육 세션을 생성하고,\n교육 실시 후 출석과 점수를 입력합니다.",
                "steps": [
                    ("1. 세션 생성", "사이드바 → 교육 운영 → 일정\n[+ 새 세션] → 프로그램 선택, 날짜, 장소, 강사 지정"),
                    ("2. 참석자 등록", "세션 상세 → [참석자 추가]\n대상자 명단에서 선택 또는 부서/직급별 일괄 등록"),
                    ("3. 출석 체크", "교육 당일: 사이드바 → 교육 운영 → 출석\n참석/불참/지각 체크 → 서명 확인"),
                    ("4. 점수 & 결과 입력", "사이드바 → 교육 운영 → 결과\n각 참석자의 점수 입력 → 합격/불합격 자동 판정"),
                ],
            },
            {
                "id": "D",
                "title": "시나리오 D: 불합격자 재교육 처리",
                "sub": "교육 불합격자를 재교육하고 재평가",
                "persona": "문경서 QA 매니저",
                "situation": "안전교육에서 3명이 불합격했습니다.\n재교육 세션을 생성하고 재평가를 진행합니다.",
                "steps": [
                    ("1. 재교육 대상 확인", "사이드바 → 후속 관리 → 재교육\n불합격자 자동 목록 표시 (이름, 점수, 불합격 사유)"),
                    ("2. 재교육 세션 생성", "[재교육 생성] 버튼 → 원본 프로그램 자동 연결\n날짜, 강사 지정 → 불합격자 자동 등록"),
                    ("3. 재교육 실시 & 재평가", "재교육 후 다시 출석/점수 입력\n합격 시: 상태 자동 업데이트"),
                    ("4. 이력 확인", "해당 직원 프로필 → 교육 이력 탭\n원래 교육 → 불합격 → 재교육 → 합격 전체 이력 표시"),
                ],
            },
            {
                "id": "E",
                "title": "시나리오 E: 신입 검사원 교육 (TQC)",
                "sub": "신입 TQC 직원의 4단계 교육 프로세스",
                "persona": "Le Thi Hoa (HR 담당)",
                "situation": "신입 QC 검사원 Pham Van An이 입사했습니다.\nTQC 팀에 배정하고 4단계 교육을 시작합니다.",
                "steps": [
                    ("1. 신입 등록", "사이드바 → 신입 TQC → 교육생\n[+ 신입 등록] → 이름, 부서, 입사일, 배정 팀 입력"),
                    ("2. 4단계 교육 진행", "1단계: 오리엔테이션 (1주)\n2단계: 이론 교육 (2주)\n3단계: 현장 실습 (2주)\n4단계: 독립 작업 (1주)"),
                    ("3. 면담 기록", "사이드바 → 신입 TQC → 면담\n각 단계별 면담 기록, 평가 코멘트 입력"),
                    ("4. 최종 평가 & 수료", "사이드바 → 신입 TQC → 최종 결과\n종합 점수, 합격 여부 판정 → 수료증 발급"),
                ],
            },
            {
                "id": "F",
                "title": "시나리오 F: 품질 문제 대응 (CAPA & 5PRS)",
                "sub": "불량 데이터 분석 → 교육 추천 → 자동 등록",
                "persona": "문경서 QA 매니저",
                "situation": "Building A에서 접착 불량이 급증했습니다.\n5PRS 데이터를 분석하고 관련 교육을 추천받습니다.",
                "steps": [
                    ("1. CAPA 등록", "사이드바 → 품질 개선 → CAPA\n[새 CAPA 등록] → 문제 유형, 발생 일시, 영향 범위 입력"),
                    ("2. 5PRS 분석 실행", "사이드바 → 품질 개선 → 5PRS\n월 선택 → [분석 실행] → TQC별 불량률 자동 계산"),
                    ("3. 교육 추천 확인", "사이드바 → 품질 개선 → 교육 추천\n불량 유형별 추천 교육 프로그램 자동 매핑\nIMMEDIATE / PREVENTIVE / SURGE 우선순위 표시"),
                    ("4. 자동 등록 & 추적", "[등록] 버튼 → 해당 직원 교육 세션 자동 생성\n교육 완료 후 불량률 변화 추적"),
                ],
            },
            {
                "id": "G",
                "title": "시나리오 G: 역량 매트릭스 관리",
                "sub": "직원 역량을 체계적으로 관리하고 갭을 분석",
                "persona": "Tran Van Duc (라인 리더)",
                "situation": "분기별 역량 평가를 실시하여\n팀원들의 역량 현황과 개선 필요 영역을 파악합니다.",
                "steps": [
                    ("1. 역량 항목 확인", "사이드바 → 인력 & 역량 → 역량 매트릭스\n부서/직급별 필요 역량 목록과 목표 레벨 확인"),
                    ("2. 직원별 레벨 평가", "역량 매트릭스에서 직원 선택\n각 역량 항목의 현재 레벨(1~5) 입력"),
                    ("3. 스킬 갭 분석", "사이드바 → 인력 & 역량 → 스킬 갭 분석\n목표 레벨 vs 현재 레벨 차이를 시각화\n갭이 큰 항목에 추천 교육 프로그램 표시"),
                    ("4. 학습 경로 생성", "갭 항목 클릭 → 추천 교육 바로 등록\n교육 완료 후 역량 레벨 자동 업데이트"),
                ],
            },
            {
                "id": "H",
                "title": "시나리오 H: 경영진 보고 & 리포트",
                "sub": "경영진에게 교육 현황을 보고하는 방법",
                "persona": "문경서 QA 매니저",
                "situation": "월간 경영회의에서 교육 현황을 보고해야 합니다.\nQ-TRAIN에서 데이터를 추출하여 보고서를 준비합니다.",
                "steps": [
                    ("1. 경영진 대시보드", "사이드바 → 대시보드 → 경영진\n전사 교육 이수율, 합격률, 재교육 현황 한눈에 확인"),
                    ("2. 부서별 분석", "사이드바 → 대시보드 → 부서별\n각 부서의 교육 완료율, 역량 달성률 비교"),
                    ("3. 리포트 생성", "사이드바 → 분석 & 리포트 → 리포트\n기간 선택 → 필터 적용 → Excel/PDF 다운로드"),
                    ("4. 교육 계획 확인", "사이드바 → 분석 & 리포트 → 교육 계획\n다음 달/분기 예정 교육, 예산, 강사 배정 현황"),
                ],
            },
        ],
        # ── Summary slide ──
        "summary_title": "핵심 정리",
        "summary_sub": "Q-TRAIN 시스템 활용 요약",
        "summary_items": [
            ("일일 업무", "매일 대시보드 확인 → 알림 체크 → 일정 점검"),
            ("교육 운영", "프로그램 등록 → 세션 생성 → 출석/결과 → 재교육"),
            ("인력 관리", "역량 매트릭스 → 갭 분석 → 학습 경로 → 레벨 업데이트"),
            ("품질 개선", "CAPA 등록 → 5PRS 분석 → 교육 추천 → 효과 추적"),
            ("보고", "경영진 대시보드 → 부서별 분석 → 리포트 다운로드"),
        ],
        "summary_rules": [
            "교육 결과는 삭제 불가 (수정 시 편집 로그 자동 기록)",
            "모든 변경 사항은 감사 로그에 자동 기록",
            "자격증 만료 전 자동 알림 (30/14/7일 전)",
            "3개 국어 지원 (한국어/영어/베트남어) — 우측 상단 언어 전환",
        ],
        "summary_tip": "단축키: Ctrl+K (빠른 검색) | 사이드바는 9개 카테고리로 구성",
        "thankyou": "감사합니다",
        "contact": "문의: ksmoon@hsvina.com",
    },

    "en": {
        "file": "Q-TRAIN_Scenario_Training_EN.pptx",
        "title": "QIP Training Management System",
        "subtitle": "Scenario-Based Practical Training Guide",
        "info": "HWK Vietnam | QIP Training Management",
        "version": "Q-TRAIN v2.0 | Feb 2026",
        "toc_title": "Training Course Overview",
        "toc_sub": "Learn Q-TRAIN through 8 real-world scenarios",
        "toc_items": [
            ("A", "Daily Morning Routine", "Dashboard check, notifications, today's schedule"),
            ("B", "Register New Training Program", "Create program → evaluation criteria → validity period"),
            ("C", "Conduct Session & Enter Results", "Create session → attendance → scores → pass/fail"),
            ("D", "Handle Failed Trainees", "Identify failures → retraining → re-evaluation"),
            ("E", "New Inspector Training (TQC)", "Team assignment → 4-stage training → interviews → final"),
            ("F", "Quality Issue Response (CAPA & 5PRS)", "Defect analysis → training recommendation → auto-enroll"),
            ("G", "Competency Matrix Management", "Define competencies → assign → gap analysis → learning path"),
            ("H", "Executive Reporting", "KPI review → report generation → share"),
        ],
        "overview_title": "Q-TRAIN System Overview",
        "overview_sub": "Integrated Training Management Platform for QIP Department Heads",
        "overview_what": "What is Q-TRAIN?",
        "overview_desc": "A web-based system that integrates training operations,\ncompetency management, and quality improvement for HWK Vietnam QIP.",
        "overview_features": [
            ("Training Ops", "Programs, schedules, attendance, results in one place"),
            ("People & Competency", "Employee DB, competency matrix, skill gap analysis"),
            ("Quality Improvement", "CAPA process, 5PRS defect analysis, training recommendations"),
            ("New TQC", "4-stage training, interview records, final evaluation"),
            ("Analytics & Reports", "Dashboards, KPIs, executive reports"),
            ("System Admin", "Audit logs, material library, data sync"),
        ],
        "overview_sidebar": "Sidebar Menu Structure (9 Categories)",
        "sidebar_cats": [
            "Dashboards", "Training Operations", "Follow-up", "People & Competency",
            "New TQC", "Projects", "Quality Improvement", "Analytics & Reports", "System Admin",
        ],
        "overview_url": "Access URL: https://q-train-web.web.app/",
        "overview_lang": "Supported Languages: Korean / English / Vietnamese",
        "scenarios": [
            {
                "id": "A",
                "title": "Scenario A: Daily Morning Routine",
                "sub": "5 minutes each morning with Q-TRAIN",
                "persona": "Mr. Moon — QA Manager",
                "situation": "Monday 8 AM. Mr. Moon arrives at the office\nand checks this week's training status and key alerts.",
                "steps": [
                    ("1. Login & Dashboard", "Open Q-TRAIN → Dashboard auto-displayed\nCheck monthly completion rate, pass rate, retraining count"),
                    ("2. Check Notifications", "Click notification icon in top bar\nCertificate expiry, incomplete training, retraining alerts"),
                    ("3. Today's Schedule", "Sidebar → Training Operations → Schedule\nView today's/this week's planned training sessions"),
                    ("4. Executive Dashboard", "Sidebar → Dashboards → Executive\nDepartment-wise training status, KPI trend charts"),
                ],
            },
            {
                "id": "B",
                "title": "Scenario B: Register New Training Program",
                "sub": "Register a new quality inspection training in the system",
                "persona": "Nguyen Thi Mai — QC Training Lead",
                "situation": "A new 'Sole Bonding Quality Inspection' training is needed.\nMai registers the program in Q-TRAIN.",
                "steps": [
                    ("1. Create Program", "Sidebar → Training Operations → Programs\n[+ New Program] → Enter name, code, category"),
                    ("2. Set Evaluation Criteria", "Passing score: 80 / Method: Written + Practical\nGrades: A(90+), B(80+), C(70+), F(<70)"),
                    ("3. Assign Target Audience", "Target positions: QC Inspectors, Line Leaders\nTarget departments: Building A, B Quality Teams"),
                    ("4. Set Validity Period", "Training validity: 12 months\nAuto-alerts at 30/14/7 days before expiry"),
                ],
            },
            {
                "id": "C",
                "title": "Scenario C: Conduct Session & Enter Results",
                "sub": "Run monthly safety training and record results",
                "persona": "Tran Van Duc — Line Leader / Trainer",
                "situation": "Create this month's safety training session,\nconduct training, and enter attendance and scores.",
                "steps": [
                    ("1. Create Session", "Sidebar → Training Operations → Schedule\n[+ New Session] → Select program, date, location, trainer"),
                    ("2. Register Attendees", "Session detail → [Add Attendees]\nSelect from target list or bulk-add by department/position"),
                    ("3. Mark Attendance", "Training day: Sidebar → Training Ops → Attendance\nCheck present/absent/late → Confirm signatures"),
                    ("4. Enter Scores & Results", "Sidebar → Training Operations → Results\nEnter each attendee's score → Auto pass/fail determination"),
                ],
            },
            {
                "id": "D",
                "title": "Scenario D: Handle Failed Trainees",
                "sub": "Retrain and re-evaluate failed trainees",
                "persona": "Mr. Moon — QA Manager",
                "situation": "3 trainees failed the safety training.\nCreate retraining sessions and re-evaluate.",
                "steps": [
                    ("1. Check Retraining List", "Sidebar → Follow-up → Retraining\nAuto-generated list of failures (name, score, reason)"),
                    ("2. Create Retraining Session", "[Create Retraining] → Original program auto-linked\nSet date, trainer → Failed trainees auto-enrolled"),
                    ("3. Retrain & Re-evaluate", "After retraining, enter attendance/scores again\nOn pass: Status automatically updated"),
                    ("4. Check History", "Employee profile → Training History tab\nFull trail: Original → Fail → Retrain → Pass"),
                ],
            },
            {
                "id": "E",
                "title": "Scenario E: New Inspector Training (TQC)",
                "sub": "4-stage onboarding process for new QC inspectors",
                "persona": "Le Thi Hoa — HR Representative",
                "situation": "New QC inspector Pham Van An has joined.\nAssign to TQC team and start 4-stage training.",
                "steps": [
                    ("1. Register New Trainee", "Sidebar → New TQC → Trainees\n[+ Register] → Name, department, hire date, assigned team"),
                    ("2. 4-Stage Training", "Stage 1: Orientation (1 week)\nStage 2: Theory Training (2 weeks)\nStage 3: Practical Training (2 weeks)\nStage 4: Independent Work (1 week)"),
                    ("3. Record Interviews", "Sidebar → New TQC → Meetings\nRecord interview notes and evaluation comments per stage"),
                    ("4. Final Evaluation", "Sidebar → New TQC → Final Result\nOverall score, pass/fail determination → Issue certificate"),
                ],
            },
            {
                "id": "F",
                "title": "Scenario F: Quality Issue Response (CAPA & 5PRS)",
                "sub": "Analyze defects → Recommend training → Auto-enroll",
                "persona": "Mr. Moon — QA Manager",
                "situation": "Bonding defects surging in Building A.\nAnalyze 5PRS data and get training recommendations.",
                "steps": [
                    ("1. Register CAPA", "Sidebar → Quality Improvement → CAPA\n[New CAPA] → Issue type, occurrence date, impact scope"),
                    ("2. Run 5PRS Analysis", "Sidebar → Quality Improvement → 5PRS\nSelect month → [Run Analysis] → Auto defect rate calculation"),
                    ("3. Review Recommendations", "Sidebar → Quality Improvement → Training Recommendations\nAuto-mapped training programs by defect type\nPriority: IMMEDIATE / PREVENTIVE / SURGE"),
                    ("4. Auto-Enroll & Track", "[Enroll] → Training sessions auto-created for target employees\nTrack defect rate changes after training completion"),
                ],
            },
            {
                "id": "G",
                "title": "Scenario G: Competency Matrix Management",
                "sub": "Systematically manage employee competencies and analyze gaps",
                "persona": "Tran Van Duc — Line Leader",
                "situation": "Conducting quarterly competency assessment\nto identify team skill levels and improvement areas.",
                "steps": [
                    ("1. Review Competencies", "Sidebar → People & Competency → Competency Matrix\nView required competencies and target levels by dept/position"),
                    ("2. Evaluate Employee Levels", "Select employee in competency matrix\nEnter current level (1~5) for each competency item"),
                    ("3. Skill Gap Analysis", "Sidebar → People & Competency → Skill Gap Analysis\nVisualize target vs. current level gaps\nRecommended training programs for large gaps"),
                    ("4. Create Learning Path", "Click gap item → Directly enroll in recommended training\nCompetency level auto-updates after training completion"),
                ],
            },
            {
                "id": "H",
                "title": "Scenario H: Executive Reporting",
                "sub": "How to report training status to management",
                "persona": "Mr. Moon — QA Manager",
                "situation": "Monthly management meeting requires training status report.\nExtract data from Q-TRAIN and prepare the report.",
                "steps": [
                    ("1. Executive Dashboard", "Sidebar → Dashboards → Executive\nCompany-wide completion rate, pass rate, retraining status"),
                    ("2. Department Analysis", "Sidebar → Dashboards → Department\nCompare each department's completion and competency rates"),
                    ("3. Generate Report", "Sidebar → Analytics & Reports → Reports\nSelect period → Apply filters → Download Excel/PDF"),
                    ("4. Training Plan Review", "Sidebar → Analytics & Reports → Training Plan\nNext month/quarter planned training, budget, trainer allocation"),
                ],
            },
        ],
        "summary_title": "Key Takeaways",
        "summary_sub": "Q-TRAIN System Usage Summary",
        "summary_items": [
            ("Daily Tasks", "Check dashboard → Review alerts → Verify schedule"),
            ("Training Ops", "Register program → Create session → Attendance/Results → Retraining"),
            ("People Mgmt", "Competency matrix → Gap analysis → Learning path → Level update"),
            ("Quality", "Register CAPA → 5PRS analysis → Training recommendations → Track effect"),
            ("Reporting", "Executive dashboard → Department analysis → Download reports"),
        ],
        "summary_rules": [
            "Training results cannot be deleted (edit logs auto-recorded)",
            "All changes are automatically recorded in audit logs",
            "Auto-alerts for certificate expiry (30/14/7 days before)",
            "3 languages supported (Korean/English/Vietnamese) — top-right language switch",
        ],
        "summary_tip": "Shortcut: Ctrl+K (Quick Search) | Sidebar has 9 organized categories",
        "thankyou": "Thank You",
        "contact": "Contact: ksmoon@hsvina.com",
    },

    "vi": {
        "file": "Q-TRAIN_Scenario_Training_VI.pptx",
        "title": "Hệ thống Quản lý Đào tạo QIP",
        "subtitle": "Hướng dẫn đào tạo thực tế theo kịch bản",
        "info": "HWK Vietnam | QIP Training Management",
        "version": "Q-TRAIN v2.0 | Tháng 02/2026",
        "toc_title": "Tổng quan khóa đào tạo",
        "toc_sub": "Học Q-TRAIN qua 8 kịch bản thực tế",
        "toc_items": [
            ("A", "Công việc hàng ngày buổi sáng", "Kiểm tra dashboard, thông báo, lịch hôm nay"),
            ("B", "Đăng ký chương trình đào tạo mới", "Tạo chương trình → tiêu chí → thời hạn hiệu lực"),
            ("C", "Thực hiện buổi đào tạo & nhập kết quả", "Tạo buổi → điểm danh → chấm điểm → đạt/không đạt"),
            ("D", "Xử lý học viên không đạt", "Xác định → đào tạo lại → đánh giá lại"),
            ("E", "Đào tạo kiểm tra viên mới (TQC)", "Phân nhóm → 4 giai đoạn → phỏng vấn → kết quả cuối"),
            ("F", "Ứng phó vấn đề chất lượng (CAPA & 5PRS)", "Phân tích lỗi → đề xuất đào tạo → tự động đăng ký"),
            ("G", "Quản lý ma trận năng lực", "Định nghĩa → gán → phân tích gap → lộ trình học tập"),
            ("H", "Báo cáo cho ban lãnh đạo", "Xem KPI → tạo báo cáo → chia sẻ"),
        ],
        "overview_title": "Giới thiệu hệ thống Q-TRAIN",
        "overview_sub": "Nền tảng quản lý đào tạo tích hợp cho Trưởng bộ phận QIP",
        "overview_what": "Q-TRAIN là gì?",
        "overview_desc": "Hệ thống web tích hợp quản lý vận hành đào tạo,\nquản lý năng lực và cải tiến chất lượng cho QIP HWK Vietnam.",
        "overview_features": [
            ("Vận hành ĐT", "Chương trình, lịch, điểm danh, kết quả tại một nơi"),
            ("Nhân sự & NL", "DB nhân viên, ma trận năng lực, phân tích gap kỹ năng"),
            ("Cải tiến CL", "Quy trình CAPA, phân tích lỗi 5PRS, đề xuất đào tạo"),
            ("TQC Mới", "Đào tạo 4 giai đoạn, ghi nhận phỏng vấn, đánh giá cuối"),
            ("PT & BC", "Dashboard, KPI, báo cáo cho ban lãnh đạo"),
            ("Quản trị HT", "Nhật ký kiểm toán, thư viện tài liệu, đồng bộ dữ liệu"),
        ],
        "overview_sidebar": "Cấu trúc menu sidebar (9 danh mục)",
        "sidebar_cats": [
            "Tổng quan", "Vận hành đào tạo", "Theo dõi sau", "Nhân sự & Năng lực",
            "TQC Mới", "Dự án", "Cải tiến chất lượng", "Phân tích & Báo cáo", "Quản trị hệ thống",
        ],
        "overview_url": "Địa chỉ truy cập: https://q-train-web.web.app/",
        "overview_lang": "Ngôn ngữ hỗ trợ: Tiếng Hàn / Tiếng Anh / Tiếng Việt",
        "scenarios": [
            {
                "id": "A",
                "title": "Kịch bản A: Công việc hàng ngày buổi sáng",
                "sub": "5 phút mỗi sáng với Q-TRAIN",
                "persona": "Ông Moon — Quản lý QA",
                "situation": "Thứ Hai 8 giờ sáng. Ông Moon đến văn phòng\nvà kiểm tra tình hình đào tạo tuần này và các cảnh báo.",
                "steps": [
                    ("1. Đăng nhập & Dashboard", "Mở Q-TRAIN → Dashboard tự động hiển thị\nKiểm tra tỷ lệ hoàn thành, tỷ lệ đạt, số cần đào tạo lại"),
                    ("2. Kiểm tra thông báo", "Click biểu tượng thông báo trên thanh trên\nHết hạn chứng chỉ, đào tạo chưa hoàn thành, cảnh báo"),
                    ("3. Lịch hôm nay", "Sidebar → Vận hành đào tạo → Lịch\nXem các buổi đào tạo hôm nay/tuần này"),
                    ("4. Dashboard lãnh đạo", "Sidebar → Tổng quan → Lãnh đạo\nTình hình đào tạo theo phòng ban, biểu đồ KPI"),
                ],
            },
            {
                "id": "B",
                "title": "Kịch bản B: Đăng ký chương trình đào tạo mới",
                "sub": "Đăng ký đào tạo kiểm tra chất lượng mới vào hệ thống",
                "persona": "Nguyen Thi Mai — Phụ trách đào tạo QC",
                "situation": "Cần chương trình 'Kiểm tra chất lượng dán đế' mới.\nMai đăng ký chương trình vào Q-TRAIN.",
                "steps": [
                    ("1. Tạo chương trình", "Sidebar → Vận hành đào tạo → Chương trình\n[+ CT mới] → Nhập tên, mã, danh mục"),
                    ("2. Thiết lập tiêu chí", "Điểm đạt: 80 / Phương pháp: Lý thuyết + Thực hành\nXếp loại: A(90+), B(80+), C(70+), F(<70)"),
                    ("3. Chỉ định đối tượng", "Vị trí: Kiểm tra viên QC, Trưởng chuyền\nPhòng ban: Nhóm CL Nhà máy A, B"),
                    ("4. Thiết lập thời hạn", "Hiệu lực đào tạo: 12 tháng\nCảnh báo tự động 30/14/7 ngày trước khi hết hạn"),
                ],
            },
            {
                "id": "C",
                "title": "Kịch bản C: Thực hiện buổi ĐT & nhập kết quả",
                "sub": "Tổ chức đào tạo an toàn hàng tháng và ghi nhận kết quả",
                "persona": "Tran Van Duc — Trưởng chuyền / Giảng viên",
                "situation": "Tạo buổi đào tạo an toàn tháng này,\nthực hiện và nhập điểm danh, điểm số.",
                "steps": [
                    ("1. Tạo buổi đào tạo", "Sidebar → Vận hành đào tạo → Lịch\n[+ Buổi mới] → Chọn CT, ngày, địa điểm, giảng viên"),
                    ("2. Đăng ký học viên", "Chi tiết buổi → [Thêm học viên]\nChọn từ danh sách hoặc thêm hàng loạt theo phòng/vị trí"),
                    ("3. Điểm danh", "Ngày ĐT: Sidebar → Vận hành ĐT → Điểm danh\nĐánh dấu có mặt/vắng/trễ → Xác nhận chữ ký"),
                    ("4. Nhập điểm & kết quả", "Sidebar → Vận hành đào tạo → Kết quả\nNhập điểm từng người → Tự động xác định đạt/không đạt"),
                ],
            },
            {
                "id": "D",
                "title": "Kịch bản D: Xử lý học viên không đạt",
                "sub": "Đào tạo lại và đánh giá lại học viên không đạt",
                "persona": "Ông Moon — Quản lý QA",
                "situation": "3 học viên không đạt đào tạo an toàn.\nTạo buổi đào tạo lại và đánh giá lại.",
                "steps": [
                    ("1. Kiểm tra DS đào tạo lại", "Sidebar → Theo dõi sau → Đào tạo lại\nDanh sách tự động (tên, điểm, lý do không đạt)"),
                    ("2. Tạo buổi đào tạo lại", "[Tạo ĐT lại] → CT gốc tự động liên kết\nĐặt ngày, GV → Học viên không đạt tự động đăng ký"),
                    ("3. ĐT lại & đánh giá lại", "Sau ĐT lại, nhập điểm danh/điểm lại\nKhi đạt: Trạng thái tự động cập nhật"),
                    ("4. Kiểm tra lịch sử", "Hồ sơ nhân viên → Tab lịch sử đào tạo\nToàn bộ: Gốc → Không đạt → ĐT lại → Đạt"),
                ],
            },
            {
                "id": "E",
                "title": "Kịch bản E: Đào tạo KTV mới (TQC)",
                "sub": "Quy trình 4 giai đoạn cho kiểm tra viên mới",
                "persona": "Le Thi Hoa — Đại diện HR",
                "situation": "KTV mới Pham Van An đã gia nhập.\nPhân nhóm TQC và bắt đầu đào tạo 4 giai đoạn.",
                "steps": [
                    ("1. Đăng ký học viên mới", "Sidebar → TQC Mới → Học viên\n[+ Đăng ký] → Tên, phòng ban, ngày nhận việc, nhóm"),
                    ("2. Đào tạo 4 giai đoạn", "GĐ 1: Định hướng (1 tuần)\nGĐ 2: Lý thuyết (2 tuần)\nGĐ 3: Thực hành (2 tuần)\nGĐ 4: Làm việc độc lập (1 tuần)"),
                    ("3. Ghi nhận phỏng vấn", "Sidebar → TQC Mới → Phỏng vấn\nGhi chú và nhận xét đánh giá theo từng giai đoạn"),
                    ("4. Đánh giá cuối cùng", "Sidebar → TQC Mới → Kết quả cuối\nĐiểm tổng, xác định đạt/không → Cấp chứng chỉ"),
                ],
            },
            {
                "id": "F",
                "title": "Kịch bản F: Ứng phó vấn đề CL (CAPA & 5PRS)",
                "sub": "Phân tích lỗi → Đề xuất đào tạo → Tự động đăng ký",
                "persona": "Ông Moon — Quản lý QA",
                "situation": "Lỗi dán đế tăng đột biến tại Nhà máy A.\nPhân tích dữ liệu 5PRS và nhận đề xuất đào tạo.",
                "steps": [
                    ("1. Đăng ký CAPA", "Sidebar → Cải tiến CL → CAPA\n[CAPA mới] → Loại vấn đề, ngày phát sinh, phạm vi"),
                    ("2. Chạy phân tích 5PRS", "Sidebar → Cải tiến CL → 5PRS\nChọn tháng → [Chạy PT] → Tự động tính tỷ lệ lỗi theo TQC"),
                    ("3. Xem đề xuất đào tạo", "Sidebar → Cải tiến CL → Đề xuất đào tạo\nCT đào tạo tự động gợi ý theo loại lỗi\nƯu tiên: NGAY LẬP TỨC / PHÒNG NGỪA / ĐỘT BIẾN"),
                    ("4. Tự động đăng ký & theo dõi", "[Đăng ký] → Buổi ĐT tự động tạo cho NV liên quan\nTheo dõi thay đổi tỷ lệ lỗi sau đào tạo"),
                ],
            },
            {
                "id": "G",
                "title": "Kịch bản G: Quản lý ma trận năng lực",
                "sub": "Quản lý năng lực nhân viên và phân tích khoảng trống",
                "persona": "Tran Van Duc — Trưởng chuyền",
                "situation": "Đánh giá năng lực hàng quý\nđể xác định trình độ và lĩnh vực cần cải thiện.",
                "steps": [
                    ("1. Xem năng lực yêu cầu", "Sidebar → Nhân sự & NL → Ma trận năng lực\nXem danh sách năng lực và mức mục tiêu theo phòng/vị trí"),
                    ("2. Đánh giá cấp độ NV", "Chọn nhân viên trong ma trận\nNhập cấp độ hiện tại (1~5) cho từng mục năng lực"),
                    ("3. Phân tích gap kỹ năng", "Sidebar → Nhân sự & NL → Phân tích gap\nTrực quan hóa chênh lệch mục tiêu vs hiện tại\nCT đào tạo gợi ý cho gap lớn"),
                    ("4. Tạo lộ trình học tập", "Click mục gap → Đăng ký ĐT gợi ý ngay\nCấp độ năng lực tự động cập nhật sau hoàn thành ĐT"),
                ],
            },
            {
                "id": "H",
                "title": "Kịch bản H: Báo cáo cho ban lãnh đạo",
                "sub": "Cách báo cáo tình hình đào tạo cho ban lãnh đạo",
                "persona": "Ông Moon — Quản lý QA",
                "situation": "Cuộc họp ban lãnh đạo hàng tháng cần báo cáo đào tạo.\nTrích xuất dữ liệu từ Q-TRAIN và chuẩn bị báo cáo.",
                "steps": [
                    ("1. Dashboard lãnh đạo", "Sidebar → Tổng quan → Lãnh đạo\nTỷ lệ hoàn thành, đạt, đào tạo lại toàn công ty"),
                    ("2. Phân tích theo phòng ban", "Sidebar → Tổng quan → Phòng ban\nSo sánh tỷ lệ hoàn thành và năng lực từng phòng"),
                    ("3. Tạo báo cáo", "Sidebar → Phân tích & BC → Báo cáo\nChọn kỳ → Lọc → Tải Excel/PDF"),
                    ("4. Xem kế hoạch đào tạo", "Sidebar → Phân tích & BC → Kế hoạch ĐT\nĐT dự kiến tháng/quý tới, ngân sách, phân công GV"),
                ],
            },
        ],
        "summary_title": "Tóm tắt chính",
        "summary_sub": "Tổng kết sử dụng hệ thống Q-TRAIN",
        "summary_items": [
            ("Hàng ngày", "Kiểm tra dashboard → Xem cảnh báo → Xác nhận lịch"),
            ("Vận hành ĐT", "Đăng ký CT → Tạo buổi → Điểm danh/KQ → Đào tạo lại"),
            ("Quản lý NS", "Ma trận NL → Phân tích gap → Lộ trình → Cập nhật cấp độ"),
            ("Chất lượng", "Đăng ký CAPA → PT 5PRS → Đề xuất ĐT → Theo dõi hiệu quả"),
            ("Báo cáo", "Dashboard LĐ → PT phòng ban → Tải báo cáo"),
        ],
        "summary_rules": [
            "Kết quả đào tạo không thể xóa (nhật ký sửa tự động ghi nhận)",
            "Tất cả thay đổi được ghi nhận trong nhật ký kiểm toán",
            "Cảnh báo tự động khi chứng chỉ sắp hết hạn (30/14/7 ngày)",
            "Hỗ trợ 3 ngôn ngữ (Hàn/Anh/Việt) — chuyển ở góc trên bên phải",
        ],
        "summary_tip": "Phím tắt: Ctrl+K (Tìm nhanh) | Sidebar có 9 danh mục tổ chức rõ ràng",
        "thankyou": "Cảm ơn",
        "contact": "Liên hệ: ksmoon@hsvina.com",
    },
}


# ═══════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════

def build_toc_slide(prs, tx):
    """Table of Contents"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, tx["toc_title"], tx["toc_sub"])

    for i, (letter, name, desc) in enumerate(tx["toc_items"]):
        row = i // 2
        col = i % 2
        x = Inches(0.5) + col * Inches(6.2)
        y = Inches(2.3) + row * Inches(1.15)

        # Card bg
        card_color = [LIGHT_BLUE, RGBColor(0xEC, 0xFD, 0xF5), RGBColor(0xFE, 0xF3, 0xC7), RGBColor(0xF3, 0xE8, 0xFF),
                      LIGHT_BLUE, RGBColor(0xEC, 0xFD, 0xF5), RGBColor(0xFE, 0xF3, 0xC7), RGBColor(0xF3, 0xE8, 0xFF)][i]
        add_shape(slide, x, y, Inches(5.9), Inches(1.0), card_color, radius=0.04)

        # Letter badge
        badge_color = [BLUE, GREEN, ORANGE, ACCENT, BLUE, GREEN, ORANGE, ACCENT][i]
        badge = add_shape(slide, x + Inches(0.15), y + Inches(0.2), Inches(0.55), Inches(0.55), badge_color, radius=0.1)
        add_text(slide, x + Inches(0.15), y + Inches(0.25), Inches(0.55), Inches(0.5),
                 letter, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.85), y + Inches(0.15), Inches(4.8), Inches(0.35),
                 name, size=15, color=DARK, bold=True)
        add_text(slide, x + Inches(0.85), y + Inches(0.52), Inches(4.8), Inches(0.4),
                 desc, size=11, color=GRAY)


def build_overview_slide(prs, tx):
    """System overview — 2 slides"""
    # Slide 1: What is Q-TRAIN
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, tx["overview_title"], tx["overview_sub"])

    # What is Q-TRAIN box
    add_shape(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.4), RGBColor(0xEF, 0xF6, 0xFF), radius=0.03)
    add_text(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.35),
             tx["overview_what"], size=18, color=BLUE, bold=True)
    add_text(slide, Inches(0.8), Inches(2.8), Inches(11.5), Inches(0.8),
             tx["overview_desc"], size=14, color=DARK)

    # Feature cards (2x3 grid)
    for i, (ftitle, fdesc) in enumerate(tx["overview_features"]):
        row = i // 3
        col = i % 3
        x = Inches(0.5) + col * Inches(4.15)
        y = Inches(4.0) + row * Inches(1.5)
        colors = [BLUE, GREEN, ORANGE, ACCENT, RGBColor(0xDB, 0x27, 0x77), NAVY]
        add_shape(slide, x, y, Inches(3.9), Inches(1.3), LIGHT_GRAY, radius=0.03)
        add_shape(slide, x, y, Inches(0.06), Inches(1.3), colors[i])
        add_text(slide, x + Inches(0.25), y + Inches(0.15), Inches(3.5), Inches(0.3),
                 ftitle, size=14, color=colors[i], bold=True)
        add_text(slide, x + Inches(0.25), y + Inches(0.5), Inches(3.5), Inches(0.7),
                 fdesc, size=11, color=GRAY)

    # Slide 2: Sidebar structure
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide2, tx["overview_sidebar"])

    # Sidebar categories as cards
    cat_colors = [BLUE, GREEN, ORANGE, ACCENT, RGBColor(0xDB, 0x27, 0x77),
                  NAVY, RED, RGBColor(0x06, 0x94, 0xA2), GRAY]
    cat_icons = ["1", "2", "3", "4", "5", "6", "7", "8", "9"]

    for i, cat in enumerate(tx["sidebar_cats"]):
        row = i // 3
        col = i % 3
        x = Inches(0.5) + col * Inches(4.15)
        y = Inches(2.3) + row * Inches(1.4)
        add_shape(slide2, x, y, Inches(3.9), Inches(1.15), LIGHT_GRAY, radius=0.04)
        # Number badge
        add_shape(slide2, x + Inches(0.15), y + Inches(0.25), Inches(0.55), Inches(0.55), cat_colors[i], radius=0.1)
        add_text(slide2, x + Inches(0.15), y + Inches(0.3), Inches(0.55), Inches(0.5),
                 cat_icons[i], size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide2, x + Inches(0.85), y + Inches(0.3), Inches(2.8), Inches(0.5),
                 cat, size=16, color=DARK, bold=True)

    add_text(slide2, Inches(0.5), Inches(6.6), Inches(12), Inches(0.3),
             tx["overview_url"], size=13, color=BLUE, bold=True)
    add_text(slide2, Inches(0.5), Inches(6.95), Inches(12), Inches(0.3),
             tx["overview_lang"], size=12, color=GRAY)


def build_scenario_slides(prs, scenario):
    """Each scenario: 1 divider slide + 1 detail slide"""
    # Divider slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    W = prs.slide_width
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), BLUE)

    add_text(slide, Inches(1), Inches(1.5), Inches(1.2), Inches(1.2),
             scenario["id"], size=72, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(2.5), Inches(1.5), Inches(9.5), Inches(0.7),
             scenario["title"], size=32, color=WHITE, bold=True)
    add_text(slide, Inches(2.5), Inches(2.3), Inches(9.5), Inches(0.4),
             scenario["sub"], size=16, color=RGBColor(0x93, 0xC5, 0xFD))

    add_shape(slide, Inches(2.5), Inches(3.1), Inches(2), Inches(0.04), BLUE)

    # Persona & Situation
    add_shape(slide, Inches(2.5), Inches(3.5), Inches(9.5), Inches(1.6), RGBColor(0x15, 0x2A, 0x4A), radius=0.03)
    add_text(slide, Inches(2.8), Inches(3.65), Inches(8.8), Inches(0.3),
             scenario["persona"], size=14, color=RGBColor(0x60, 0xA5, 0xFA), bold=True)
    add_text(slide, Inches(2.8), Inches(4.1), Inches(8.8), Inches(0.8),
             scenario["situation"], size=13, color=RGBColor(0xBF, 0xDB, 0xFE))

    # Detail slide with steps
    detail = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(detail, scenario["title"], scenario["sub"])

    step_colors = [BLUE, GREEN, ORANGE, ACCENT]

    for i, (step_title, step_desc) in enumerate(scenario["steps"]):
        row = i // 2
        col = i % 2
        x = Inches(0.5) + col * Inches(6.2)
        y = Inches(2.3) + row * Inches(2.4)

        # Step card
        add_shape(detail, x, y, Inches(5.9), Inches(2.2), LIGHT_GRAY, radius=0.03)
        add_shape(detail, x, y, Inches(5.9), Inches(0.06), step_colors[i])

        add_text(detail, x + Inches(0.25), y + Inches(0.2), Inches(5.4), Inches(0.3),
                 step_title, size=15, color=step_colors[i], bold=True)

        # Description lines
        lines = step_desc.split('\n')
        for j, line in enumerate(lines):
            add_text(detail, x + Inches(0.25), y + Inches(0.6) + j * Inches(0.35), Inches(5.4), Inches(0.35),
                     line, size=12, color=DARK)


def build_summary_slide(prs, tx):
    """Summary slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, tx["summary_title"], tx["summary_sub"])

    # Workflow summary cards
    for i, (stitle, sdesc) in enumerate(tx["summary_items"]):
        x = Inches(0.5) + i * Inches(2.45)
        y = Inches(2.3)
        colors = [BLUE, GREEN, ORANGE, RED, ACCENT]
        add_shape(slide, x, y, Inches(2.25), Inches(2.3), LIGHT_GRAY, radius=0.04)
        add_shape(slide, x, y, Inches(2.25), Inches(0.06), colors[i])
        add_text(slide, x + Inches(0.15), y + Inches(0.2), Inches(1.95), Inches(0.3),
                 stitle, size=13, color=colors[i], bold=True)
        # Wrap description
        desc_lines = sdesc.split(' → ')
        for j, dl in enumerate(desc_lines):
            prefix = "→ " if j > 0 else ""
            add_text(slide, x + Inches(0.15), y + Inches(0.6) + j * Inches(0.35), Inches(1.95), Inches(0.35),
                     prefix + dl, size=10, color=DARK)

    # Rules
    add_shape(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(2.2), RGBColor(0xFE, 0xF3, 0xC7), radius=0.03)
    rules_box = add_text(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.3),
                         "RULES", size=14, color=ORANGE, bold=True)

    for i, rule in enumerate(tx["summary_rules"]):
        add_text(slide, Inches(0.8), Inches(5.4) + i * Inches(0.35), Inches(11.5), Inches(0.35),
                 f"  {rule}", size=12, color=DARK)

    add_text(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.3),
             tx["summary_tip"], size=11, color=BLUE, bold=True)


def generate_pptx(lang):
    """Generate PPTX for given language"""
    tx = TEXTS[lang]
    prs, W, H = create_presentation()

    # 1. Title slide
    title_slide(prs, tx["title"], tx["subtitle"], tx["info"],
                "https://q-train-web.web.app/", tx["version"])

    # 2. Table of Contents
    build_toc_slide(prs, tx)

    # 3. System Overview (2 slides)
    build_overview_slide(prs, tx)

    # 4. Scenarios (2 slides each = 16 slides)
    for scenario in tx["scenarios"]:
        build_scenario_slides(prs, scenario)

    # 5. Summary
    build_summary_slide(prs, tx)

    # 6. Thank You
    thank_you_slide(prs, "https://q-train-web.web.app/", "ksmoon@hsvina.com")

    # Save
    out_path = os.path.join(SCRIPT_DIR, tx["file"])
    prs.save(out_path)
    slide_count = len(prs.slides)
    file_size = os.path.getsize(out_path)
    print(f"  [{lang.upper()}] {tx['file']} — {slide_count} slides, {file_size / 1024:.1f} KB")
    return out_path


# ═══════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════

if __name__ == "__main__":
    print("Generating Q-TRAIN Scenario Training Guides...")
    files = []
    for lang in ["ko", "en", "vi"]:
        path = generate_pptx(lang)
        files.append(path)
    print(f"\nDone! Generated {len(files)} files.")
