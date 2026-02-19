#!/usr/bin/env python3
"""
Q-TRAIN QIP 부서장 업무 가이드
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
업무 유형별로 구성된 종합 교육 자료 (~35 슬라이드)
대상: QIP 부서장 (Department Head)
"""

import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from pptx_helpers import (
    NAVY, BLUE, LIGHT_BLUE, WHITE, DARK, GRAY, LIGHT_GRAY, GREEN, ORANGE, RED, ACCENT,
    add_bg, add_shape, add_text, add_para, add_card, slide_header, add_table_slide,
    create_presentation, title_slide, thank_you_slide,
)

# ── Additional Colors ──
AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)
RED_BG = RGBColor(0xFE, 0xF2, 0xF2)
GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)
BLUE_BG = RGBColor(0xEF, 0xF6, 0xFF)
PURPLE_BG = RGBColor(0xF5, 0xF3, 0xFF)
TEAL = RGBColor(0x06, 0x69, 0xB4)
PINK = RGBColor(0xDB, 0x27, 0x77)
INDIGO = RGBColor(0x49, 0x46, 0xB8)
EMERALD = RGBColor(0x04, 0x7A, 0x57)

prs, W, H = create_presentation()


# ═══════════════════════════════════════════════════════════════
# HELPER: Section divider slide
# ═══════════════════════════════════════════════════════════════
def section_divider(prs, section_letter, section_title, section_desc, color, items):
    """Create a section divider with numbered items."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, color)
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), WHITE)

    # Section label
    add_text(slide, Inches(1), Inches(1.5), Inches(1.2), Inches(1.2),
             section_letter, size=60, color=WHITE, bold=True)
    add_shape(slide, Inches(2.3), Inches(1.8), Inches(0.04), Inches(0.8), RGBColor(0xFF, 0xFF, 0xFF))
    add_text(slide, Inches(2.6), Inches(1.5), Inches(9), Inches(0.7),
             section_title, size=36, color=WHITE, bold=True)
    add_text(slide, Inches(2.6), Inches(2.2), Inches(9), Inches(0.4),
             section_desc, size=16, color=RGBColor(0xBF, 0xDB, 0xFE))

    # Items
    y = Inches(3.5)
    for i, item in enumerate(items):
        add_text(slide, Inches(2.6), y + i * Inches(0.5), Inches(0.4), Inches(0.4),
                 f"{i+1}.", size=18, color=RGBColor(0x93, 0xC5, 0xFD), bold=True)
        add_text(slide, Inches(3.1), y + i * Inches(0.5), Inches(8), Inches(0.4),
                 item, size=16, color=WHITE)
    return slide


# ═══════════════════════════════════════════════════════════════
# HELPER: Feature page with left info + right table/cards
# ═══════════════════════════════════════════════════════════════
def feature_page(prs, title, subtitle, left_title, left_items, right_title, right_content_fn):
    """Create a feature page with left description and right custom content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, title, subtitle)

    # Left panel
    add_shape(slide, Inches(0.5), Inches(2.3), Inches(5.8), Inches(4.8), WHITE, radius=0.04)
    add_text(slide, Inches(0.7), Inches(2.45), Inches(5.4), Inches(0.3),
             left_title, size=15, color=NAVY, bold=True)

    txBox = add_text(slide, Inches(0.7), Inches(2.9), Inches(5.2), Inches(4.0), "", size=11, color=DARK)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for item in left_items:
        if isinstance(item, tuple):
            text, clr, bld = item
        else:
            text, clr, bld = item, DARK, False
        add_para(tf, "  " + text, size=11, color=clr, bold=bld, space_before=Pt(5))

    # Right panel
    right_content_fn(slide)
    return slide


# ═══════════════════════════════════════════════════════════════
# HELPER: Numbered step flow (horizontal)
# ═══════════════════════════════════════════════════════════════
def add_step_flow(slide, x, y, steps, colors=None):
    """Add a horizontal numbered step flow."""
    cols = len(steps)
    step_w = min(Inches(2.0), (Inches(12.0)) / cols)
    for i, step_text in enumerate(steps):
        sx = x + i * step_w
        sc = colors[i] if colors and i < len(colors) else BLUE
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, sx, y, Inches(0.32), Inches(0.32))
        circle.fill.solid()
        circle.fill.fore_color.rgb = sc
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.paragraphs[0].text = str(i + 1)
        ctf.paragraphs[0].font.size = Pt(10)
        ctf.paragraphs[0].font.color.rgb = WHITE
        ctf.paragraphs[0].font.bold = True
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
        add_text(slide, sx + Inches(0.38), y - Inches(0.02), step_w - Inches(0.55), Inches(0.4),
                 step_text, size=10, color=DARK)
        if i < cols - 1:
            add_text(slide, sx + step_w - Inches(0.18), y, Inches(0.2), Inches(0.3),
                     "→", size=13, color=GRAY, bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════
title_slide(prs,
            "QIP 부서장 업무 가이드",
            "업무 유형별로 배우는 Q-TRAIN 시스템 활용법",
            "HWK Vietnam (화승비나) | QIP Training Management",
            "https://q-train-web.web.app",
            "2026.02 | Comprehensive Guide v2.0")


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: Table of Contents
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "목차", "업무 유형별 9개 카테고리")

toc_items = [
    ("A", "시스템 접속 & 기본 화면", "로그인, 메뉴 구조, 언어 전환", BLUE),
    ("B", "일일 업무 — 교육 현황 모니터링", "대시보드, 알림, 일일 점검", NAVY),
    ("C", "교육 운영 관리", "프로그램/일정/출석/결과 관리", GREEN),
    ("D", "직원 역량 & 진도 관리", "직원DB, 진도 매트릭스, 역량/갭 분석", TEAL),
    ("E", "재교육 & 자격증 관리", "불합격 처리, 만료 관리, 이수증 발급", ORANGE),
    ("F", "신입사원 교육 (New TQC)", "4단계 교육, 면담, 퇴사 분석", ACCENT),
    ("G", "품질 개선 활동", "CAPA, 5PRS 불량 분석, 교육 추천", RED),
    ("H", "프로젝트 & 협업", "프로젝트 관리, 과제, 캘린더", INDIGO),
    ("I", "분석/리포트 & 시스템 관리", "경영진 보고, 감사, 데이터 동기화", EMERALD),
]

for i, (letter, title, desc, color) in enumerate(toc_items):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + col * Inches(4.2)
    y = Inches(2.3) + row * Inches(1.55)

    card = add_shape(slide, x, y, Inches(3.95), Inches(1.35), WHITE, radius=0.04)
    # Letter badge
    badge = add_shape(slide, x + Inches(0.1), y + Inches(0.1), Inches(0.55), Inches(0.55), color, radius=0.08)
    btf = badge.text_frame
    btf.paragraphs[0].text = letter
    btf.paragraphs[0].font.size = Pt(18)
    btf.paragraphs[0].font.color.rgb = WHITE
    btf.paragraphs[0].font.bold = True
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, x + Inches(0.8), y + Inches(0.1), Inches(3), Inches(0.35),
             title, size=13, color=NAVY, bold=True)
    add_text(slide, x + Inches(0.8), y + Inches(0.45), Inches(3), Inches(0.25),
             desc, size=10, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SECTION A: 시스템 접속 & 기본 화면
# ═══════════════════════════════════════════════════════════════
section_divider(prs, "A", "시스템 접속 & 기본 화면",
    "Q-TRAIN에 접속하고 기본 화면을 이해합니다",
    BLUE, ["접속 방법 & 로그인", "메뉴 구조 이해", "언어 전환 기능"])

# A-1: System Access
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "시스템 접속 방법", "웹 브라우저에서 바로 접속 가능 (앱 설치 불필요)")

# Access info cards
add_card(slide, Inches(0.5), Inches(2.3), Inches(3.8), Inches(1.5),
         "접속 주소 (URL)", "q-train-web.web.app", "Chrome / Edge 권장", LIGHT_BLUE, BLUE, NAVY)
add_card(slide, Inches(4.6), Inches(2.3), Inches(3.8), Inches(1.5),
         "로그인 방식", "Google 계정", "@hsvina.com 회사 메일로 로그인", LIGHT_BLUE, BLUE, NAVY)
add_card(slide, Inches(8.7), Inches(2.3), Inches(4.1), Inches(1.5),
         "모바일 접속", "PWA 지원", "홈 화면에 추가 가능 (앱처럼 사용)", LIGHT_BLUE, BLUE, NAVY)

# Login steps
add_shape(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(3.1), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(4.2), Inches(11.5), Inches(0.3),
         "로그인 절차", size=15, color=NAVY, bold=True)

login_steps = [
    ("Step 1", "브라우저에서 https://q-train-web.web.app 접속", BLUE),
    ("Step 2", "Google 로그인 버튼 클릭 → @hsvina.com 계정 선택", GREEN),
    ("Step 3", "권한에 따라 자동으로 대시보드 이동", ACCENT),
]

for i, (step, desc, color) in enumerate(login_steps):
    y = Inches(4.7) + i * Inches(0.65)
    badge = add_shape(slide, Inches(0.9), y, Inches(1.0), Inches(0.35), color, radius=0.1)
    btf = badge.text_frame
    btf.paragraphs[0].text = step
    btf.paragraphs[0].font.size = Pt(10)
    btf.paragraphs[0].font.color.rgb = WHITE
    btf.paragraphs[0].font.bold = True
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(slide, Inches(2.1), y, Inches(10), Inches(0.35),
             desc, size=12, color=DARK)

# Language & features
add_text(slide, Inches(0.7), Inches(6.4), Inches(11.5), Inches(0.4),
         "기능:  우측 상단 언어 전환 (🇻🇳 베트남어 / 🇰🇷 한국어 / 🇺🇸 영어)  |  Ctrl+K 글로벌 검색  |  알림 벨 아이콘",
         size=11, color=GRAY)


# A-2: Menu Structure
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "전체 메뉴 구조", "좌측 사이드바 기반 7개 섹션, 30+ 메뉴 항목")

menu_sections = [
    ("메인 메뉴\n(7개)", [
        "대시보드", "교육 프로그램", "진도 현황",
        "교육 일정", "출석 체크", "결과 입력", "직원 관리"
    ], NAVY),
    ("신입 TQC\n(7개)", [
        "대시보드", "교육생 관리", "미팅 관리",
        "최종 결과", "수료증", "퇴사 분석", "설정"
    ], ACCENT),
    ("프로젝트\n(5개)", [
        "대시보드", "팀원", "과제",
        "캘린더", "설정"
    ], INDIGO),
    ("CAPA\n(2개)", [
        "관리", "새 등록"
    ], RED),
    ("5PRS 검사\n(2개)", [
        "대시보드", "교육 추천"
    ], TEAL),
    ("관리\n(9개)", [
        "재교육", "이수증", "리포트",
        "부서별", "알림", "평가",
        "교육자료", "역량", "스킬갭"
    ], ORANGE),
    ("관리자\n(6개)", [
        "경영진", "감사대응", "강사관리",
        "연간계획", "감사로그", "동기화"
    ], EMERALD),
]

for i, (section_name, items, color) in enumerate(menu_sections):
    x = Inches(0.3) + i * Inches(1.85)
    y = Inches(2.3)

    # Header
    add_shape(slide, x, y, Inches(1.72), Inches(0.65), color, radius=0.04)
    add_text(slide, x + Inches(0.05), y + Inches(0.05), Inches(1.62), Inches(0.55),
             section_name, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Items
    for j, item in enumerate(items):
        iy = y + Inches(0.75) + j * Inches(0.35)
        add_shape(slide, x, iy, Inches(1.72), Inches(0.3), WHITE if j % 2 == 0 else LIGHT_GRAY, radius=0.02)
        add_text(slide, x + Inches(0.08), iy, Inches(1.56), Inches(0.3),
                 item, size=9, color=DARK)

# Note
add_text(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
         "Tip: 사이드바 메뉴는 접기/펼치기 가능 | 현재 위치는 상단 브레드크럼에 표시됩니다",
         size=11, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SECTION B: 일일 업무 — 교육 현황 모니터링
# ═══════════════════════════════════════════════════════════════
section_divider(prs, "B", "일일 업무 — 교육 현황 모니터링",
    "매일 확인해야 할 핵심 KPI와 알림",
    NAVY, ["대시보드에서 KPI 한눈에 확인", "알림 확인 및 긴급 조치", "재교육/만료 경고 처리"])

# B-1: Dashboard
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "메인 대시보드", "경로: 사이드바 → 대시보드  |  매일 아침 가장 먼저 확인")

# KPI cards
add_card(slide, Inches(0.5), Inches(2.3), Inches(2.9), Inches(1.2),
         "전체 직원", "1,200명", "Active 기준", LIGHT_BLUE, BLUE, NAVY)
add_card(slide, Inches(3.6), Inches(2.3), Inches(2.9), Inches(1.2),
         "이번달 교육", "8건 / 12건", "완료 / 예정", LIGHT_BLUE, BLUE, NAVY)
add_card(slide, Inches(6.7), Inches(2.3), Inches(2.9), Inches(1.2),
         "전체 이수율", "67.2%", "목표: 90%", GREEN_BG, GREEN, GREEN)
add_card(slide, Inches(9.8), Inches(2.3), Inches(2.9), Inches(1.2),
         "재교육 대상", "45명", "즉시 조치 필요", RED_BG, RED, RED)

# Dashboard components
add_shape(slide, Inches(0.5), Inches(3.8), Inches(6), Inches(3.4), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(3.9), Inches(5.5), Inches(0.3),
         "대시보드 구성 요소", size=14, color=NAVY, bold=True)
dash_items = [
    ("등급 분포 차트", "AA/A/B/C 등급별 인원 파이차트", BLUE),
    ("월별 교육 추세", "계획 vs 완료 교육 건수 막대 그래프", GREEN),
    ("동별 현황", "Building별 재교육/만료 인원 현황", ORANGE),
    ("만료 예정", "30일 내 자격 만료 예정 목록", RED),
    ("다가오는 세션", "예정된 교육 세션 목록", ACCENT),
    ("최근 재교육 대상", "FAIL/ABSENT 처리 필요 인원", PINK),
]
for i, (title, desc, color) in enumerate(dash_items):
    y = Inches(4.3) + i * Inches(0.45)
    add_shape(slide, Inches(0.9), y + Inches(0.05), Inches(0.15), Inches(0.15), color)
    add_text(slide, Inches(1.15), y, Inches(2), Inches(0.3),
             title, size=11, color=DARK, bold=True)
    add_text(slide, Inches(3.1), y, Inches(3.2), Inches(0.3),
             desc, size=10, color=GRAY)

# Action tips
add_shape(slide, Inches(6.8), Inches(3.8), Inches(6), Inches(3.4), AMBER_BG, radius=0.04)
add_text(slide, Inches(7), Inches(3.9), Inches(5.5), Inches(0.3),
         "부서장 일일 점검 포인트", size=14, color=RGBColor(0x92, 0x40, 0x0E), bold=True)

check_items = [
    "1. 재교육 대상 확인 → 45명 이상이면 즉시 세션 배정",
    "2. 이수율 확인 → 70% 미만이면 부서별 원인 분석",
    "3. 만료 예정 30일 이내 직원 → 재교육 세션 스케줄링",
    "4. 이번달 교육 완료/예정 비교 → 일정 지연 확인",
    "5. 등급 분포 확인 → C등급 비율 10% 초과 시 점검",
    "",
    "Tip: 경영진 대시보드(/executive)에서",
    "더 상세한 KPI 분석이 가능합니다.",
]
txBox = add_text(slide, Inches(7), Inches(4.4), Inches(5.5), Inches(2.5), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in check_items:
    c = RGBColor(0x78, 0x35, 0x0F)
    add_para(tf, "  " + item, size=11, color=c, space_before=Pt(4))


# ═══════════════════════════════════════════════════════════════
# SECTION C: 교육 운영 관리
# ═══════════════════════════════════════════════════════════════
section_divider(prs, "C", "교육 운영 관리",
    "교육의 전체 라이프사이클을 관리합니다",
    GREEN, ["교육 프로그램 등록/수정", "교육 세션 스케줄링", "출석 관리", "결과 입력 및 등급 관리"])

# C-1: Programs
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "교육 프로그램 관리", "경로: 사이드바 → 교육 프로그램  |  모든 교육의 기준을 정의")

# Left: What is a program
add_shape(slide, Inches(0.5), Inches(2.3), Inches(6), Inches(4.8), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(2.4), Inches(5.5), Inches(0.3),
         "프로그램이란?", size=15, color=NAVY, bold=True)

prog_items = [
    ("모든 교육 세션의 기준이 되는 '교육 과정 정의'", DARK, False),
    ("", DARK, False),
    ("필수 설정 항목:", NAVY, True),
    ("  프로그램 코드 (예: QIP-001)", DARK, False),
    ("  교육명 — 한국어/영어/베트남어 3개 언어", DARK, False),
    ("  카테고리: QIP / 생산 / 재교육 / 신입 / 승진", DARK, False),
    ("  평가 방식: 점수제 or 합격/불합격", DARK, False),
    ("  합격 기준 점수 (예: 70점)", DARK, False),
    ("  등급 기준: AA≥95 / A≥85 / B≥70 / C<70", DARK, False),
    ("  교육 시간 (예: 4시간)", DARK, False),
    ("  유효기간 (예: 12개월, 없으면 영구)", DARK, False),
    ("  대상 직급 (복수 선택 가능)", DARK, False),
    ("", DARK, False),
    ("핵심 규칙:", RED, True),
    ("  프로그램 삭제 = 비활성화 (soft delete)", RED, False),
    ("  모든 변경사항 자동 기록 (change_logs)", RED, False),
]
txBox = add_text(slide, Inches(0.7), Inches(2.8), Inches(5.5), Inches(4.0), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for text, clr, bld in prog_items:
    add_para(tf, text, size=11, color=clr, bold=bld, space_before=Pt(3))

# Right: Example
add_text(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(0.3),
         "프로그램 등록 예시", size=14, color=NAVY, bold=True)

prog_data = [
    ["항목", "설정 값"],
    ["프로그램 코드", "QIP-001"],
    ["프로그램명 (KR)", "외관검사 기초"],
    ["프로그램명 (EN)", "Visual Inspection Basic"],
    ["프로그램명 (VI)", "Kiem tra truc quan co ban"],
    ["카테고리", "QIP"],
    ["평가 방식", "SCORE (점수제)"],
    ["합격 점수", "70점"],
    ["등급 기준", "AA≥95 / A≥85 / B≥70"],
    ["교육 시간", "4시간"],
    ["유효기간", "12개월"],
    ["대상 직급", "Worker, Line Leader"],
]
add_table_slide(slide, Inches(6.8), Inches(2.7), Inches(6), prog_data,
                [Inches(2.5), Inches(3.5)])

# Operation flow
add_text(slide, Inches(6.8), Inches(6.6), Inches(6), Inches(0.4),
         "조작: + 추가 버튼 → 양식 입력 → 저장 | 수정/비활성화 가능",
         size=10, color=GRAY)


# C-2: Schedule & Attendance & Results
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "교육 세션 → 출석 → 결과 입력", "하나의 교육이 진행되는 전체 흐름")

# Step flow at top
add_shape(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.8), LIGHT_BLUE, radius=0.04)
add_step_flow(slide, Inches(0.8), Inches(2.35),
    ["세션 생성\n(/schedule)", "참석자 등록", "출석 체크\n(/attendance)", "점수 입력\n(/results)", "등급 자동계산", "자격증 발급"],
    [NAVY, BLUE, GREEN, ORANGE, ACCENT, EMERALD])

# Three panels
panels = [
    ("세션 관리 (/schedule)", [
        "프로그램 선택 → 날짜/시간 설정",
        "강사 배정 (강사 관리에서 등록된 강사)",
        "장소 입력 (교육장 / 라인 / 회의실)",
        "최대 인원 설정",
        "참석 대상자 선택 (부서/직급 필터)",
        "상태: 계획 → 완료 → 취소",
    ], BLUE),
    ("출석 관리 (/attendance)", [
        "세션 선택 → 출석 체크 화면",
        "출석: PRESENT / LATE / ABSENT",
        "체크인/체크아웃 시간 기록",
        "일괄 출석 처리 가능",
        "ABSENT → 자동으로 재교육 대상",
        "",
    ], GREEN),
    ("결과 입력 (/results)", [
        "세션 선택 → 참석자별 점수 입력",
        "등급 자동 계산 (AA/A/B/C)",
        "FAIL 시 자동 재교육 플래그",
        "중복 입력 감지 (동일 세션+직원)",
        "수정 시 반드시 사유 입력 필수",
        "삭제 불가 (No Delete Policy)",
    ], ORANGE),
]

for i, (title, items, color) in enumerate(panels):
    x = Inches(0.5) + i * Inches(4.2)
    y = Inches(3.2)
    add_shape(slide, x, y, Inches(3.95), Inches(3.2), WHITE, radius=0.04)
    add_shape(slide, x, y, Inches(3.95), Inches(0.4), color, radius=0.04)
    add_shape(slide, x, y + Inches(0.25), Inches(3.95), Inches(0.15), color)
    add_text(slide, x + Inches(0.1), y + Inches(0.05), Inches(3.75), Inches(0.3),
             title, size=12, color=WHITE, bold=True)
    txBox = add_text(slide, x + Inches(0.15), y + Inches(0.5), Inches(3.65), Inches(2.5), "", size=10, color=DARK)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for item in items:
        c = RED if "삭제 불가" in item or "ABSENT" in item else DARK
        add_para(tf, "  " + item, size=10, color=c, bold=("불가" in item), space_before=Pt(5))

# Bottom rule
add_shape(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.7), RED_BG, radius=0.04)
add_text(slide, Inches(0.7), Inches(6.65), Inches(11.5), Inches(0.5),
         "핵심 규칙:  결과 수정 시 사유 필수 → result_edit_logs에 영구 기록  |  결과 삭제 절대 불가  |  등급 자동 계산 (AA≥95 / A≥85 / B≥70 / C<70)",
         size=11, color=RED, bold=True)


# ═══════════════════════════════════════════════════════════════
# SECTION D: 직원 역량 & 진도 관리
# ═══════════════════════════════════════════════════════════════
section_divider(prs, "D", "직원 역량 & 진도 관리",
    "직원 정보와 교육 진도, 역량 수준을 종합 관리합니다",
    TEAL, ["직원 데이터베이스 관리", "진도 현황 매트릭스", "역량 매트릭스 & 스킬 갭 분석"])

# D-1: Employees & Progress
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "직원 관리 & 진도 현황", "직원 DB와 교육 이수 매트릭스를 한눈에 확인")

# Left: Employees
add_shape(slide, Inches(0.5), Inches(2.3), Inches(6), Inches(4.8), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(2.4), Inches(5.5), Inches(0.3),
         "직원 관리 (/employees)", size=14, color=NAVY, bold=True)

emp_items = [
    ("직원 마스터 데이터 관리", NAVY, True),
    ("  사번, 이름, 부서, 직급, 동, 라인, 입사일, 상태", DARK, False),
    ("", DARK, False),
    ("필터링:", NAVY, True),
    ("  부서: ASSEMBLY, STITCHING, CUTTING 등 11개", DARK, False),
    ("  동: A, A1, B, B1, C, D 등 17개", DARK, False),
    ("  직급: TQC, RQC, CFA, Worker, Leader 등", DARK, False),
    ("  상태: Active / Inactive", DARK, False),
    ("", DARK, False),
    ("클릭하면 개인 상세 페이지:", NAVY, True),
    ("  교육 이력 전체 조회", DARK, False),
    ("  자격증 보유 현황", DARK, False),
    ("  역량 수준 확인", DARK, False),
    ("  재교육 이력", DARK, False),
]
txBox = add_text(slide, Inches(0.7), Inches(2.8), Inches(5.5), Inches(4.0), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for text, clr, bld in emp_items:
    add_para(tf, text, size=11, color=clr, bold=bld, space_before=Pt(3))

# Right: Progress Matrix
add_shape(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(4.8), LIGHT_BLUE, radius=0.04)
add_text(slide, Inches(7), Inches(2.4), Inches(5.5), Inches(0.3),
         "진도 현황 매트릭스 (/progress)", size=14, color=NAVY, bold=True)

progress_items = [
    ("직원 × 프로그램 교차 매트릭스", NAVY, True),
    ("", DARK, False),
    ("각 셀의 상태:", NAVY, True),
    ("  PASS — 유효한 자격 보유 (초록)", GREEN, False),
    ("  EXPIRING — 30일 내 만료 예정 (주황)", ORANGE, False),
    ("  EXPIRED — 자격 만료 (빨강)", RED, False),
    ("  FAIL — 불합격 (빨강)", RED, False),
    ("  NOT_TAKEN — 미이수 (회색)", GRAY, False),
    ("", DARK, False),
    ("4가지 필터:", NAVY, True),
    ("  동 (Building) / 부서 / 직급 / 카테고리", DARK, False),
    ("", DARK, False),
    ("셀 클릭 시:", NAVY, True),
    ("  해당 직원의 해당 교육 상세 정보 팝업", DARK, False),
    ("  점수, 등급, 교육일, 만료일 표시", DARK, False),
]
txBox = add_text(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(4.0), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for text, clr, bld in progress_items:
    add_para(tf, text, size=11, color=clr, bold=bld, space_before=Pt(3))


# D-2: Competency & Skill Gap
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "역량 매트릭스 & 스킬 갭 분석", "역량을 정의하고, 갭을 분석하고, 학습 경로를 배정합니다")

# Competency
add_shape(slide, Inches(0.5), Inches(2.3), Inches(6), Inches(2.2), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(2.4), Inches(5.5), Inches(0.3),
         "역량 매트릭스 (/competency)", size=14, color=NAVY, bold=True)

comp_table = [
    ["탭", "기능", "주요 조작"],
    ["역량 목록", "역량 정의 및 카테고리 관리", "+ 추가 / 수정 / 삭제"],
    ["스킬 매트릭스", "직원×역량 교차 매트릭스 (히트맵)", "셀 클릭 → 레벨 변경"],
    ["학습 경로", "순차 교육 과정 설계", "경로 생성 → 직원 할당"],
]
add_table_slide(slide, Inches(0.7), Inches(2.8), Inches(5.5), comp_table,
                [Inches(1.3), Inches(2.2), Inches(2)])

# Skill Gap
add_shape(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(2.2), WHITE, radius=0.04)
add_text(slide, Inches(7), Inches(2.4), Inches(5.5), Inches(0.3),
         "스킬 갭 분석 (/skill-gap)", size=14, color=NAVY, bold=True)

gap_table = [
    ["기능", "설명"],
    ["부서별 갭 분석", "부서 선택 → 역량별 현재 vs 목표 비교"],
    ["갭 히트맵", "부서×역량 교차표, 색상으로 갭 정도 표시"],
    ["교육 연계", "갭 있는 역량 → 추천 교육 프로그램 표시"],
    ["트렌드", "월별 갭 변화 추이 차트"],
]
add_table_slide(slide, Inches(7), Inches(2.8), Inches(5.5), gap_table,
                [Inches(2), Inches(3.5)])

# Competency levels
add_shape(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.5), GREEN_BG, radius=0.04)
add_text(slide, Inches(0.7), Inches(4.8), Inches(11.5), Inches(0.3),
         "역량 수준 체계 (5단계)", size=14, color=EMERALD, bold=True)

levels = [
    ("NOVICE (1)", "기초 지식만 보유, 지도 필요", RGBColor(0xEF, 0x44, 0x44)),
    ("BEGINNER (2)", "기본 작업 가능, 감독 필요", ORANGE),
    ("COMPETENT (3)", "독립 작업 가능", RGBColor(0xEA, 0xB3, 0x08)),
    ("PROFICIENT (4)", "복잡한 문제 해결 가능", GREEN),
    ("EXPERT (5)", "다른 직원 교육 가능, 개선 주도", BLUE),
]

for i, (level, desc, color) in enumerate(levels):
    x = Inches(0.7) + i * Inches(2.45)
    badge = add_shape(slide, x, Inches(5.25), Inches(2.2), Inches(0.35), color, radius=0.08)
    btf = badge.text_frame
    btf.paragraphs[0].text = level
    btf.paragraphs[0].font.size = Pt(10)
    btf.paragraphs[0].font.color.rgb = WHITE
    btf.paragraphs[0].font.bold = True
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(slide, x, Inches(5.7), Inches(2.2), Inches(0.3),
             desc, size=9, color=DARK, align=PP_ALIGN.CENTER)

# Category info
add_text(slide, Inches(0.7), Inches(6.3), Inches(11.5), Inches(0.6),
         "역량 카테고리:  TECHNICAL (기술) | QUALITY (품질) | SAFETY (안전) | LEADERSHIP (리더십) | COMMUNICATION (소통) | PROCESS (공정)\n"
         "학습 경로 유형:  ONBOARDING (입사) | POSITION (직무) | PROMOTION (승진) | SPECIALIZATION (전문) | REMEDIAL (보충)",
         size=10, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SECTION E: 재교육 & 자격증 관리
# ═══════════════════════════════════════════════════════════════
section_divider(prs, "E", "재교육 & 자격증 관리",
    "불합격/만료 대상자를 자동 감지하고, 자격증을 발급합니다",
    ORANGE, ["재교육 대상 자동 감지 & 관리", "자격 유효기간 만료 추적", "이수증/자격증 발급"])

# E-1: Retraining & Certificates
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "재교육 대상 & 이수증 관리", "시스템이 자동으로 대상자를 식별하고, 자격증을 발급합니다")

# Left: Retraining
add_shape(slide, Inches(0.5), Inches(2.3), Inches(6), Inches(4.8), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(2.4), Inches(5.5), Inches(0.3),
         "재교육 대상 (/retraining)", size=14, color=NAVY, bold=True)

retrain_items = [
    ("자동 식별되는 4가지 대상:", NAVY, True),
    ("", DARK, False),
    ("  1. FAILED — 시험 불합격자", RED, True),
    ("     점수가 합격 기준 미달 → 자동 감지", DARK, False),
    ("", DARK, False),
    ("  2. ABSENT — 교육 결석자", RED, True),
    ("     출석 체크에서 결석 처리된 직원", DARK, False),
    ("", DARK, False),
    ("  3. EXPIRED — 자격 유효기간 만료", ORANGE, True),
    ("     유효기간이 지난 자격 보유자", DARK, False),
    ("", DARK, False),
    ("  4. EXPIRING — 30일 내 만료 예정", ORANGE, True),
    ("     사전 경고로 미리 재교육 계획 가능", DARK, False),
    ("", DARK, False),
    ("부서장 액션:", NAVY, True),
    ("  대상자 확인 → 재교육 세션 생성 → 결과 추적", BLUE, False),
]
txBox = add_text(slide, Inches(0.7), Inches(2.8), Inches(5.5), Inches(4.0), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for text, clr, bld in retrain_items:
    add_para(tf, text, size=10, color=clr, bold=bld, space_before=Pt(2))

# Right: Certificates
add_shape(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(4.8), WHITE, radius=0.04)
add_text(slide, Inches(7), Inches(2.4), Inches(5.5), Inches(0.3),
         "이수증/자격증 발급 (/certificates)", size=14, color=NAVY, bold=True)

cert_items = [
    ("이수증 발급 기능:", NAVY, True),
    ("", DARK, False),
    ("  개별 발급 — 직원 선택 → 프로그램 선택 → 발급", DARK, False),
    ("  일괄 발급 — 세션 완료 시 PASS 전원 일괄 발급", DARK, False),
    ("  템플릿 관리 — 이수증 디자인 템플릿 관리", DARK, False),
    ("  자격 취소 — 사유 입력 후 자격 취소 (감사 기록)", DARK, False),
    ("", DARK, False),
    ("이수증에 포함되는 정보:", NAVY, True),
    ("  직원 이름 / 교육 프로그램명", DARK, False),
    ("  교육 완료일 / 점수 / 등급", DARK, False),
    ("  유효기간 / 디지털 인증 ID", DARK, False),
    ("", DARK, False),
    ("자격증 유효기간 자동 추적:", NAVY, True),
    ("  발급 시 만료일 자동 계산", DARK, False),
    ("  30일 전 EXPIRING 알림 발송", ORANGE, False),
    ("  만료 시 EXPIRED 상태 자동 변경", RED, False),
]
txBox = add_text(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(4.0), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for text, clr, bld in cert_items:
    add_para(tf, text, size=10, color=clr, bold=bld, space_before=Pt(2))


# ═══════════════════════════════════════════════════════════════
# SECTION F: 신입사원 교육 (New TQC)
# ═══════════════════════════════════════════════════════════════
section_divider(prs, "F", "신입사원 교육 (New TQC)",
    "신입 검사원의 4단계 온보딩 교육을 체계적으로 관리합니다",
    ACCENT, ["TQC 대시보드로 현황 파악", "교육생 등록 & 4단계 교육", "면담 관리 & 퇴사 분석"])

# F-1: New TQC Overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "신입 TQC 교육 시스템", "경로: 사이드바 → 신입 TQC 교육  |  완전한 온보딩 서브시스템")

# 4 Stage cards
stages = [
    ("Stage 1", "오리엔테이션", "공장 안전 교육\n사내 규정 안내\n기본 품질 개념", GREEN),
    ("Stage 2", "기초 교육", "QIP 교육 수강\n시험 응시\n합격 기준 충족", BLUE),
    ("Stage 3", "라인 배치", "실제 라인 OJT\n선임자 지도\n실무 능력 습득", ORANGE),
    ("Stage 4", "현장 평가", "독립 작업 평가\n최종 합격 판정\n수료증 발급", ACCENT),
]

for i, (label, title, desc, color) in enumerate(stages):
    x = Inches(0.5) + i * Inches(3.15)
    add_shape(slide, x, Inches(2.3), Inches(2.95), Inches(2.3), WHITE, radius=0.04)
    add_shape(slide, x, Inches(2.3), Inches(2.95), Inches(0.55), color, radius=0.04)
    add_shape(slide, x, Inches(2.65), Inches(2.95), Inches(0.2), color)
    add_text(slide, x + Inches(0.1), Inches(2.35), Inches(2.75), Inches(0.2),
             label, size=9, color=RGBColor(0xBF, 0xDB, 0xFE))
    add_text(slide, x + Inches(0.1), Inches(2.52), Inches(2.75), Inches(0.3),
             title, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txBox = add_text(slide, x + Inches(0.15), Inches(3.0), Inches(2.65), Inches(1.4), "", size=10, color=DARK)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for line in desc.split('\n'):
        add_para(tf, "  " + line, size=10, color=DARK, space_before=Pt(6))
    if i < 3:
        add_text(slide, x + Inches(2.95), Inches(3.2), Inches(0.2), Inches(0.3),
                 "→", size=18, color=GRAY, bold=True)

# TQC Sub-features
add_shape(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.4), PURPLE_BG, radius=0.04)
add_text(slide, Inches(0.7), Inches(4.9), Inches(11.5), Inches(0.3),
         "TQC 서브 메뉴 안내", size=14, color=ACCENT, bold=True)

tqc_menus = [
    ["메뉴", "기능", "주요 조작"],
    ["대시보드", "전체 현황: 교육중/완료/퇴사 인원, 이직률, 평균 진도", "현황 모니터링"],
    ["교육생 관리", "신입 등록 → 자동으로 4단계+3회 면담 생성", "등록/진도 관리"],
    ["미팅 관리", "1주/1개월/3개월 면담 일정 관리", "면담 실시/기록"],
    ["최종 결과", "4단계 완료 후 최종 합격/불합격 판정", "합격 처리"],
    ["수료증", "합격자 수료증 발급", "발급/출력"],
    ["퇴사 분석", "퇴사 사유 8가지 카테고리 분류, 분석", "퇴사 기록/분석"],
    ["설정", "TQC 팀 관리, 교육 기간 설정", "팀 생성/수정"],
]
add_table_slide(slide, Inches(0.7), Inches(5.3), Inches(11.5), tqc_menus,
                [Inches(2), Inches(6), Inches(3.5)], header_color=ACCENT)


# ═══════════════════════════════════════════════════════════════
# SECTION G: 품질 개선 활동
# ═══════════════════════════════════════════════════════════════
section_divider(prs, "G", "품질 개선 활동",
    "CAPA와 5PRS를 활용한 품질 기반 교육 관리",
    RED, ["CAPA 관리 (시정/예방 조치)", "5PRS 불량 데이터 분석", "AI 기반 교육 추천 시스템"])

# G-1: CAPA
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "CAPA 관리", "경로: 사이드바 → CAPA  |  품질 문제의 체계적 시정/예방 조치")

# CAPA 5 stages
capa_stages = [
    ("1", "발견", RED),
    ("2", "조사", ORANGE),
    ("3", "조치", BLUE),
    ("4", "검증", ACCENT),
    ("5", "종료", GREEN),
]

add_shape(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.2), WHITE, radius=0.04)
for i, (num, name, color) in enumerate(capa_stages):
    x = Inches(1.0) + i * Inches(2.4)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(2.5), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.paragraphs[0].text = num
    ctf.paragraphs[0].font.size = Pt(16)
    ctf.paragraphs[0].font.color.rgb = WHITE
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(slide, x - Inches(0.2), Inches(3.05), Inches(0.9), Inches(0.25),
             name, size=11, color=color, bold=True, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(slide, x + Inches(0.65), Inches(2.55), Inches(1.5), Inches(0.3),
                 "→→→", size=12, color=LIGHT_GRAY)

# CAPA details
add_shape(slide, Inches(0.5), Inches(3.6), Inches(6), Inches(3.5), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(3.7), Inches(5.5), Inches(0.3),
         "CAPA 등록 항목", size=14, color=NAVY, bold=True)

capa_items = [
    "제목 / 설명 (문제 상세)",
    "유형: CORRECTIVE (시정) 또는 PREVENTIVE (예방)",
    "심각도: LOW / MEDIUM / HIGH / CRITICAL",
    "출처: AUDIT / CUSTOMER / INTERNAL / 5PRS",
    "근본 원인 분석 (5-Why)",
    "시정/예방 조치 계획",
    "담당자 / 목표 완료일",
    "교육 연계 — 관련 교육 프로그램 링크",
    "효과성 검증 결과",
]
txBox = add_text(slide, Inches(0.7), Inches(4.1), Inches(5.5), Inches(2.8), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in capa_items:
    add_para(tf, "  " + item, size=11, color=DARK, space_before=Pt(5))

# Example
add_shape(slide, Inches(6.8), Inches(3.6), Inches(6), Inches(3.5), RED_BG, radius=0.04)
add_text(slide, Inches(7), Inches(3.7), Inches(5.5), Inches(0.3),
         "실전 예시: CAPA-2026-003", size=14, color=RED, bold=True)

example_items = [
    ("문제:", "A동 L3 외관검사 불량률 3.2% (기준 1.5%)", DARK),
    ("유형:", "CORRECTIVE (시정)", DARK),
    ("심각도:", "HIGH", RED),
    ("원인:", "QIP-001 교육이 이론 위주 → 실무 적용 부족", DARK),
    ("", "(5-Why 분석 결과)", GRAY),
    ("조치:", "보충교육 + 미세 스크래치 판별 가이드 제작", DARK),
    ("", "+ 커리큘럼에 실습 2시간 추가", DARK),
    ("교육:", "QIP-001 재교육 → 해당 라인 30명 등록", BLUE),
    ("결과:", "불량률 3.2% → 1.1% (기준 이내)", GREEN),
]
txBox = add_text(slide, Inches(7), Inches(4.1), Inches(5.5), Inches(2.8), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for label, text, clr in example_items:
    line = f"  {label} {text}" if label else f"       {text}"
    add_para(tf, line, size=10, color=clr, bold=(label != ""), space_before=Pt(3))


# G-2: 5PRS & Training Recommendations
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "5PRS 불량 분석 & 교육 추천", "경로: 사이드바 → 5PRS  |  불량 데이터 기반 자동 교육 추천")

# Left: 5PRS Dashboard
add_shape(slide, Inches(0.5), Inches(2.3), Inches(6), Inches(2.3), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(2.4), Inches(5.5), Inches(0.3),
         "5PRS 대시보드 (/five-prs)", size=14, color=NAVY, bold=True)

prs_items = [
    "월별 불량 데이터 조회 (Google Apps Script API)",
    "동/라인별 불량률 KPI 카드",
    "불량 유형별 추세 차트 (월별 트렌드)",
    "TQC별 불량률 히트맵",
    "불량 유형 분포 파이차트",
]
txBox = add_text(slide, Inches(0.7), Inches(2.8), Inches(5.5), Inches(1.5), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in prs_items:
    add_para(tf, "  " + item, size=11, color=DARK, space_before=Pt(5))

# Right: Recommendations
add_shape(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(2.3), GREEN_BG, radius=0.04)
add_text(slide, Inches(7), Inches(2.4), Inches(5.5), Inches(0.3),
         "교육 추천 (/five-prs/training-recommendations)", size=14, color=EMERALD, bold=True)

rec_items = [
    "불량률 기준 초과 → 교육 자동 추천",
    "우선순위: IMMEDIATE / PREVENTIVE / SURGE",
    "불량유형 → 교육 프로그램 자동 매핑",
    "대상 직원 자동 선별 (TQC-직원 연결)",
    "원클릭 교육 등록 → 결과까지 추적",
]
txBox = add_text(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(1.5), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in rec_items:
    add_para(tf, "  " + item, size=11, color=DARK, space_before=Pt(5))

# Bottom: Data flow
add_shape(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.4), LIGHT_BLUE, radius=0.04)
add_text(slide, Inches(0.7), Inches(4.9), Inches(11.5), Inches(0.3),
         "5PRS → 교육 추천 → 교육 등록 흐름", size=14, color=NAVY, bold=True)

add_step_flow(slide, Inches(0.8), Inches(5.4),
    ["5PRS 데이터\n수집", "불량률\n기준 비교", "초과 항목\n식별", "교육 프로그램\n매핑", "대상 직원\n선별", "교육 세션\n생성", "결과\n추적"],
    [NAVY, BLUE, RED, ORANGE, ACCENT, GREEN, EMERALD])

# Auto schedule note
add_shape(slide, Inches(0.7), Inches(6.0), Inches(11.5), Inches(1.0), AMBER_BG, radius=0.04)
add_text(slide, Inches(0.9), Inches(6.1), Inches(11), Inches(0.8),
         "자동화:  매주 월요일 06:00에 Cloud Function이 자동으로 5PRS 데이터를 분석하고,\n"
         "IMMEDIATE 우선순위 대상자를 교육 프로그램에 자동 등록합니다. 결과는 이메일로 발송됩니다.",
         size=11, color=RGBColor(0x92, 0x40, 0x0E), bold=True)


# ═══════════════════════════════════════════════════════════════
# SECTION H: 프로젝트 & 협업
# ═══════════════════════════════════════════════════════════════
section_divider(prs, "H", "프로젝트 & 협업",
    "품질 부서의 프로젝트와 과제를 체계적으로 관리합니다",
    INDIGO, ["프로젝트 관리 & 팀원 배정", "과제 관리 (칸반 보드)", "캘린더 & 일정 관리"])

# H-1: Projects
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "프로젝트 관리", "경로: 사이드바 → 프로젝트  |  품질 부서 내 프로젝트 협업 도구")

proj_menus = [
    ["메뉴", "기능", "주요 특징"],
    ["대시보드", "프로젝트 현황 종합 보기", "과제 상태 분포, 팀 통계, 긴급 과제"],
    ["팀원 관리", "프로젝트 참여 인원 관리", "역할 배정, 활성/비활성, 성과 지표"],
    ["과제 관리", "과제 생성/할당/추적", "칸반 보드, 우선순위, 담당자 배정"],
    ["캘린더", "프로젝트 일정 시각화", "타임라인 뷰, 마감일 관리"],
    ["설정", "프로젝트 기본 설정", "프로젝트 생성/수정, 일반 설정"],
]
add_table_slide(slide, Inches(0.5), Inches(2.3), Inches(12.3), proj_menus,
                [Inches(2), Inches(4.5), Inches(5.8)], header_color=INDIGO)

add_text(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(0.4),
         "활용 예시: 'ISO 인증 준비 프로젝트', '연간 교육 계획 수립', '5S 개선 활동' 등 품질 부서의 프로젝트를 체계적으로 관리",
         size=11, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SECTION I: 분석/리포트 & 시스템 관리
# ═══════════════════════════════════════════════════════════════
section_divider(prs, "I", "분석/리포트 & 시스템 관리",
    "경영진 보고, 감사 대응, 데이터 관리를 수행합니다",
    EMERALD, ["경영진/부서별 대시보드", "리포트 생성 & 내보내기", "감사 대응 & 시스템 관리"])

# I-1: Executive & Reports
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "경영진 대시보드 & 리포트", "경영진 보고용 KPI와 다양한 리포트를 생성합니다")

# Executive Dashboard
add_shape(slide, Inches(0.5), Inches(2.3), Inches(6), Inches(2.5), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(2.4), Inches(5.5), Inches(0.3),
         "경영진 대시보드 (/executive)", size=14, color=NAVY, bold=True)

exec_items = [
    ("5개 탭:", NAVY, True),
    ("  Overview — 전체 교육 현황 요약", DARK, False),
    ("  KPI — 핵심 성과 지표 상세", DARK, False),
    ("  Trends — 월별/분기별 추세 분석", DARK, False),
    ("  Benchmarks — 목표 대비 달성률", DARK, False),
    ("  ROI — 교육 투자 대비 효과 분석", DARK, False),
]
txBox = add_text(slide, Inches(0.7), Inches(2.8), Inches(5.5), Inches(1.8), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for text, clr, bld in exec_items:
    add_para(tf, text, size=11, color=clr, bold=bld, space_before=Pt(4))

# Reports
add_shape(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(2.5), WHITE, radius=0.04)
add_text(slide, Inches(7), Inches(2.4), Inches(5.5), Inches(0.3),
         "리포트 생성 (/reports)", size=14, color=NAVY, bold=True)

report_items = [
    ("리포트 종류:", NAVY, True),
    ("  프로그램별 이수 현황", DARK, False),
    ("  직원별 교육 이력", DARK, False),
    ("  부서별 이수율 비교", DARK, False),
    ("  재교육 현황 분석", DARK, False),
    ("  내보내기: CSV / PDF / 인쇄", BLUE, True),
]
txBox = add_text(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(1.8), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for text, clr, bld in report_items:
    add_para(tf, text, size=11, color=clr, bold=bld, space_before=Pt(4))

# Additional management tools
add_shape(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.2), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(5.1), Inches(11.5), Inches(0.3),
         "기타 관리 도구", size=14, color=NAVY, bold=True)

admin_table = [
    ["메뉴", "경로", "기능", "사용 빈도"],
    ["강사 관리", "/trainers", "사내/외부 강사 등록, 전문 분야, 배정 이력", "수시"],
    ["교육 자료", "/materials", "교육 자료 파일 업로드/관리 (폴더 구조)", "수시"],
    ["연간 계획", "/training-plan", "연간 교육 일정 계획 수립", "연 1회"],
    ["효과성 평가", "/evaluation", "교육 후 만족도/효과성 설문 관리", "교육 후"],
    ["감사 대응", "/audit", "감사 규정 준수 현황 및 지적 사항 관리", "분기"],
    ["감사 로그", "/audit-log", "시스템 전체 변경 이력 조회 (읽기 전용)", "감사 시"],
    ["데이터 동기화", "/data-sync", "Google Sheets ↔ Firestore 동기화", "초기/수시"],
    ["알림", "/notifications", "시스템 알림 확인 및 설정", "매일"],
]
add_table_slide(slide, Inches(0.7), Inches(5.5), Inches(11.5), admin_table,
                [Inches(1.8), Inches(2.0), Inches(5.5), Inches(1.5)])


# ═══════════════════════════════════════════════════════════════
# SLIDE: 사이드바 개선 제안
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "사이드바 메뉴 개선 제안", "현재 7개 섹션 → 9개 카테고리로 재구성하여 직관성 향상")

# Current vs Proposed
add_shape(slide, Inches(0.5), Inches(2.3), Inches(5.5), Inches(4.8), RED_BG, radius=0.04)
add_text(slide, Inches(0.7), Inches(2.4), Inches(5), Inches(0.3),
         "현재 구조 (문제점)", size=14, color=RED, bold=True)

current_issues = [
    ("'메인 메뉴' — 너무 포괄적", RED, True),
    ("  대시보드와 운영 페이지가 혼재", DARK, False),
    ("", DARK, False),
    ("'관리' 섹션 — 9개 항목이 뒤섞임", RED, True),
    ("  재교육, 이수증, 리포트, 역량, 평가 등", DARK, False),
    ("  성격이 다른 기능이 한 그룹에 혼재", DARK, False),
    ("", DARK, False),
    ("'관리자' 섹션 — 역할이 불명확", RED, True),
    ("  경영진 대시보드 + 기술 도구 혼재", DARK, False),
    ("", DARK, False),
    ("역량/스킬갭 — '관리'에 묻혀 있음", ORANGE, True),
    ("  QIP 부서장에게 핵심 기능인데 찾기 어려움", DARK, False),
]
txBox = add_text(slide, Inches(0.7), Inches(2.8), Inches(5), Inches(4.0), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for text, clr, bld in current_issues:
    add_para(tf, text, size=10, color=clr, bold=bld, space_before=Pt(3))

# Proposed
add_shape(slide, Inches(6.3), Inches(2.3), Inches(6.5), Inches(4.8), GREEN_BG, radius=0.04)
add_text(slide, Inches(6.5), Inches(2.4), Inches(6), Inches(0.3),
         "제안 구조 (9개 카테고리)", size=14, color=EMERALD, bold=True)

proposed = [
    ("📊 대시보드", "메인 / 경영진 / 부서별", NAVY),
    ("📚 교육 운영", "프로그램 / 일정 / 출석 / 결과 / 진도", GREEN),
    ("🔄 후속 관리", "재교육 / 이수증 / 효과성 평가", ORANGE),
    ("👤 인력 & 역량", "직원 / 강사 / 역량 매트릭스 / 스킬갭", TEAL),
    ("🎓 신입 TQC", "(기존과 동일)", ACCENT),
    ("📋 프로젝트", "(기존과 동일)", INDIGO),
    ("🔍 품질 개선", "CAPA / 5PRS / 교육 추천", RED),
    ("📈 분석 & 리포트", "리포트 / 연간 계획 / 알림", BLUE),
    ("⚙️ 시스템 관리", "감사 / 감사 로그 / 교육 자료 / 동기화", EMERALD),
]

for i, (name, items, color) in enumerate(proposed):
    y = Inches(2.85) + i * Inches(0.45)
    badge = add_shape(slide, Inches(6.5), y, Inches(0.3), Inches(0.3), color, radius=0.08)
    add_text(slide, Inches(6.9), y, Inches(2.2), Inches(0.3),
             name, size=10, color=DARK, bold=True)
    add_text(slide, Inches(9.1), y, Inches(3.5), Inches(0.3),
             items, size=9, color=GRAY)


# Second improvement slide: Detailed proposal
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "사이드바 개선 제안 — 상세 비교", "각 메뉴 항목의 이동 경로")

comparison_data = [
    ["현재 위치", "메뉴 항목", "제안 위치", "이동 사유"],
    ["메인", "대시보드", "📊 대시보드", "대시보드 그룹으로 통합"],
    ["관리자", "경영진 대시보드", "📊 대시보드", "대시보드 그룹으로 통합"],
    ["관리", "부서별 대시보드", "📊 대시보드", "대시보드 그룹으로 통합"],
    ["메인", "프로그램/일정/출석/결과", "📚 교육 운영", "교육 운영 워크플로우로 통합"],
    ["메인", "진도 현황", "📚 교육 운영", "교육 진도 관련"],
    ["관리", "재교육 대상", "🔄 후속 관리", "교육 후 후속 조치"],
    ["관리", "이수증 발급", "🔄 후속 관리", "교육 후 후속 조치"],
    ["관리", "효과성 평가", "🔄 후속 관리", "교육 후 후속 조치"],
    ["메인", "직원 관리", "👤 인력 & 역량", "인력 관련 통합"],
    ["관리자", "강사 관리", "👤 인력 & 역량", "인력 관련 통합"],
    ["관리", "역량 매트릭스", "👤 인력 & 역량", "핵심 기능으로 승격"],
    ["관리", "스킬 갭 분석", "👤 인력 & 역량", "핵심 기능으로 승격"],
    ["관리", "리포트", "📈 분석 & 리포트", "분석 그룹으로 이동"],
    ["관리자", "연간 계획", "📈 분석 & 리포트", "계획/분석 그룹"],
    ["관리", "알림", "📈 분석 & 리포트", "정보 확인용"],
    ["관리자", "감사 대응/로그", "⚙️ 시스템 관리", "관리자 도구 통합"],
    ["관리", "교육 자료", "⚙️ 시스템 관리", "리소스 관리"],
    ["관리자", "데이터 동기화", "⚙️ 시스템 관리", "관리자 도구"],
]
add_table_slide(slide, Inches(0.5), Inches(2.2), Inches(12.3), comparison_data,
                [Inches(1.5), Inches(2.8), Inches(2.5), Inches(5.5)])


# ═══════════════════════════════════════════════════════════════
# SLIDE: 업무별 체크리스트
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "QIP 부서장 업무 체크리스트", "일/주/월별 정기 업무 가이드")

checklist_data = [
    ["주기", "업무", "시스템 위치", "핵심 확인 사항"],
    ["매일", "대시보드 KPI 확인", "/dashboard", "이수율, 재교육 대상, 만료 예정"],
    ["매일", "알림 확인 & 처리", "/notifications", "긴급 알림 즉시 대응"],
    ["매일", "교육 결과 확인", "/results", "당일 교육 결과 검토"],
    ["매주", "재교육 대상 처리", "/retraining", "FAIL/ABSENT 세션 배정"],
    ["매주", "TQC 신입 면담 확인", "/new-tqc/meetings", "면담 일정 미이행 확인"],
    ["매주", "5PRS 불량 데이터 확인", "/five-prs", "기준 초과 항목 교육 추천 확인"],
    ["매월", "경영진 보고 준비", "/executive", "KPI 리포트 생성"],
    ["매월", "만료 예정 자격 관리", "/retraining", "30일 내 만료 → 재교육 계획"],
    ["매월", "CAPA 현황 점검", "/capa", "진행 중 건 추적, 미완료 건 독촉"],
    ["매월", "역량 갭 분석", "/skill-gap", "부서별 갭 확인, 학습 경로 업데이트"],
    ["분기", "연간 계획 리뷰", "/training-plan", "계획 대비 실적 분석"],
    ["분기", "감사 대응 준비", "/audit", "규정 준수 현황 확인"],
]
add_table_slide(slide, Inches(0.5), Inches(2.3), Inches(12.3), checklist_data,
                [Inches(0.9), Inches(2.5), Inches(2.8), Inches(6.1)])


# ═══════════════════════════════════════════════════════════════
# SLIDE: 핵심 규칙 정리
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Q-TRAIN 핵심 규칙 & 원칙", "반드시 숙지해야 할 시스템 운영 원칙")

# Data integrity rules
add_shape(slide, Inches(0.5), Inches(2.3), Inches(6), Inches(2.0), RED_BG, radius=0.04)
add_text(slide, Inches(0.7), Inches(2.4), Inches(5.5), Inches(0.3),
         "데이터 무결성 규칙 (절대 위반 불가)", size=14, color=RED, bold=True)

integrity_items = [
    "교육 결과 삭제 불가 (No Delete Policy)",
    "결과 수정 시 반드시 사유 입력 → 감사 로그 기록",
    "프로그램 삭제 = 비활성화 (Soft Delete)",
    "감사 로그는 읽기 전용 (수정/삭제 불가)",
]
txBox = add_text(slide, Inches(0.7), Inches(2.8), Inches(5.5), Inches(1.3), "", size=12, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in integrity_items:
    add_para(tf, "  " + item, size=12, color=RED, bold=True, space_before=Pt(6))

# Auto features
add_shape(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(2.0), GREEN_BG, radius=0.04)
add_text(slide, Inches(7), Inches(2.4), Inches(5.5), Inches(0.3),
         "자동화 기능", size=14, color=EMERALD, bold=True)

auto_items = [
    "등급 자동 계산 (AA≥95 / A≥85 / B≥70)",
    "불합격/결석 → 자동 재교육 대상 식별",
    "자격 만료 30일 전 → 자동 알림",
    "주간 5PRS 분석 → 자동 교육 추천 (월요일 06:00)",
]
txBox = add_text(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(1.3), "", size=12, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in auto_items:
    add_para(tf, "  " + item, size=12, color=EMERALD, space_before=Pt(6))

# Multi-language
add_shape(slide, Inches(0.5), Inches(4.5), Inches(6), Inches(1.5), BLUE_BG, radius=0.04)
add_text(slide, Inches(0.7), Inches(4.6), Inches(5.5), Inches(0.3),
         "다국어 지원", size=14, color=BLUE, bold=True)

lang_items = [
    "UI: 베트남어 / 한국어 / 영어 전환 가능",
    "프로그램명: 3개 언어로 등록 가능",
    "우측 상단 언어 전환 버튼으로 즉시 변경",
]
txBox = add_text(slide, Inches(0.7), Inches(5.0), Inches(5.5), Inches(0.8), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in lang_items:
    add_para(tf, "  " + item, size=11, color=DARK, space_before=Pt(5))

# Shortcuts
add_shape(slide, Inches(6.8), Inches(4.5), Inches(6), Inches(1.5), PURPLE_BG, radius=0.04)
add_text(slide, Inches(7), Inches(4.6), Inches(5.5), Inches(0.3),
         "유용한 단축키 & 기능", size=14, color=ACCENT, bold=True)

shortcut_items = [
    "Ctrl + K → 글로벌 검색 (직원/프로그램 빠른 검색)",
    "브레드크럼 → 현재 위치 확인 + 상위 메뉴 이동",
    "PWA → 모바일에서도 앱처럼 사용 가능",
]
txBox = add_text(slide, Inches(7), Inches(5.0), Inches(5.5), Inches(0.8), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in shortcut_items:
    add_para(tf, "  " + item, size=11, color=DARK, space_before=Pt(5))

# Performance targets
add_shape(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(1.0), LIGHT_GRAY, radius=0.04)
add_text(slide, Inches(0.7), Inches(6.3), Inches(11.5), Inches(0.7),
         "성능 기준:  LCP < 2.5초  |  FID < 100ms  |  CLS < 0.1  |  3개 언어 즉시 전환  |  PWA 오프라인 지원\n"
         "보안:  Google 인증 기반  |  역할별 접근 제어  |  Firestore 보안 규칙 적용  |  모든 변경 감사 기록",
         size=10, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE: FAQ
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "자주 묻는 질문 (FAQ)", "QIP 부서장이 자주 하는 질문과 답변")

faqs = [
    ("Q: 교육 결과를 잘못 입력했으면?",
     "A: 수정 가능하지만 반드시 사유를 입력해야 합니다. 모든 수정 이력은 영구 보존됩니다."),
    ("Q: 교육 결과를 삭제할 수 있나요?",
     "A: 아니요. No Delete Policy에 따라 교육 결과는 절대 삭제 불가합니다."),
    ("Q: 자격증 유효기간이 만료되면?",
     "A: 시스템이 자동으로 30일 전 경고 → 만료 시 EXPIRED 처리 → 재교육 대상에 추가됩니다."),
    ("Q: 여러 명의 결과를 한번에 입력하려면?",
     "A: 세션별 일괄 입력 가능합니다. 엑셀 대량 업로드도 지원합니다."),
    ("Q: 5PRS 추천이 매주 자동으로 되나요?",
     "A: 네. 매주 월요일 06:00에 자동 분석 후, IMMEDIATE 대상자는 자동 등록됩니다."),
    ("Q: 부서별로 교육 현황을 비교하려면?",
     "A: /department 부서별 대시보드 또는 /reports에서 부서별 리포트를 생성하세요."),
    ("Q: CAPA와 교육은 어떻게 연계되나요?",
     "A: CAPA 조치 계획에서 '교육 연계' 항목으로 해당 교육 프로그램을 링크할 수 있습니다."),
    ("Q: 모바일에서도 사용할 수 있나요?",
     "A: 네. PWA 지원으로 모바일 브라우저에서 사용하거나 홈 화면에 추가하여 앱처럼 사용 가능합니다."),
]

add_shape(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(5.0), WHITE, radius=0.04)
y_pos = Inches(2.35)
for q, a in faqs:
    add_text(slide, Inches(0.7), y_pos, Inches(11.5), Inches(0.25),
             q, size=11, color=NAVY, bold=True)
    add_text(slide, Inches(0.7), y_pos + Inches(0.25), Inches(11.5), Inches(0.25),
             a, size=10, color=DARK)
    y_pos += Inches(0.57)


# ═══════════════════════════════════════════════════════════════
# SLIDE: Thank You
# ═══════════════════════════════════════════════════════════════
thank_you_slide(prs, "https://q-train-web.web.app", "ksmoon@hsvina.com")


# ═══════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "Q-TRAIN_QIP_Head_Guide_KO.pptx")
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
