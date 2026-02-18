#!/usr/bin/env python3
"""
Q-TRAIN System Guide - PowerPoint Presentation Generator
Educational slide deck explaining how to use Q-TRAIN system
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors ──
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
BLUE = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
GREEN = RGBColor(0x05, 0x96, 0x69)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT = RGBColor(0x7C, 0x3A, 0xED)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ════════════════════════════════════════
# Helper Functions
# ════════════════════════════════════════

def add_bg(slide, color):
    """Set solid background color"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, radius=None):
    """Add a colored rectangle shape"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if hasattr(shape, 'adjustments') and radius:
        try:
            shape.adjustments[0] = radius
        except:
            pass
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font_name='Arial'):
    """Add a text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_para(text_frame, text, size=16, color=DARK, bold=False, space_before=Pt(4), space_after=Pt(4), align=PP_ALIGN.LEFT, bullet=False):
    """Add a paragraph to existing text frame"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Arial'
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    if bullet:
        p.level = 0
    return p


def add_card(slide, left, top, width, height, title, value, subtitle="", bg_color=WHITE, title_color=GRAY, value_color=DARK, icon=""):
    """Add a KPI card"""
    card = add_shape(slide, left, top, width, height, bg_color, radius=0.05)
    card.shadow.inherit = False

    if icon:
        add_text(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.5), Inches(0.4),
                 icon, size=22, color=title_color)

    add_text(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.3),
             title, size=11, color=title_color)
    add_text(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), Inches(0.5),
             value, size=28, color=value_color, bold=True)
    if subtitle:
        add_text(slide, left + Inches(0.2), top + height - Inches(0.35), width - Inches(0.4), Inches(0.3),
                 subtitle, size=10, color=GRAY)
    return card


def slide_header(slide, title, subtitle="", bg_color=WHITE):
    """Add consistent header bar"""
    add_bg(slide, bg_color)
    # Top bar
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), BLUE)
    # Logo area
    add_text(slide, Inches(0.5), Inches(0.25), Inches(3), Inches(0.4),
             "Q-TRAIN", size=14, color=BLUE, bold=True)
    add_text(slide, Inches(0.5), Inches(0.55), Inches(5), Inches(0.3),
             "HWK Vietnam QIP Training Management System", size=9, color=GRAY)
    # Title
    add_text(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
             title, size=28, color=NAVY, bold=True)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.65), Inches(12), Inches(0.4),
                 subtitle, size=14, color=GRAY)
    # Divider
    add_shape(slide, Inches(0.5), Inches(2.05), Inches(1.5), Inches(0.04), BLUE)


def add_table_slide(slide, left, top, width, rows_data, col_widths, header_color=NAVY):
    """Add a styled table"""
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.4 * rows))
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for r, row_data in enumerate(rows_data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = 'Arial'
                if r == 0:
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.bold = True
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    return table_shape


# ════════════════════════════════════════
# SLIDE 1: Title Slide
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, NAVY)

# Accent shapes
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), BLUE)
add_shape(slide, Inches(0), Inches(7.0), W, Inches(0.5), RGBColor(0x15, 0x2A, 0x4A))

# Main title
add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
         "Q-TRAIN", size=60, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6),
         "QIP Training Management System", size=28, color=RGBColor(0x93, 0xC5, 0xFD))

# Divider
add_shape(slide, Inches(1), Inches(3.3), Inches(2), Inches(0.05), BLUE)

# Subtitle
add_text(slide, Inches(1), Inches(3.6), Inches(10), Inches(0.5),
         "시스템 사용 가이드 — 시뮬레이션으로 배우는 Q-TRAIN", size=20, color=RGBColor(0xBF, 0xDB, 0xFE))

# Info boxes
add_text(slide, Inches(1), Inches(4.5), Inches(5), Inches(0.3),
         "HWK Vietnam (화승비나) | QIP 교육 관리", size=14, color=RGBColor(0x93, 0xC5, 0xFD))
add_text(slide, Inches(1), Inches(4.9), Inches(5), Inches(0.3),
         "https://q-train-web.web.app", size=13, color=RGBColor(0x60, 0xA5, 0xFA))

# Bottom
add_text(slide, Inches(1), Inches(6.5), Inches(5), Inches(0.3),
         "2026.02 | Version 1.0", size=11, color=RGBColor(0x64, 0x78, 0x96))


# ════════════════════════════════════════
# SLIDE 2: Table of Contents
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "목차", "Table of Contents")

chapters = [
    ("01", "등장인물 소개", "시뮬레이션 캐릭터와 시스템 역할"),
    ("02", "마스터 데이터 등록", "시스템 운영을 위한 기초 데이터 세팅"),
    ("03", "교육 프로그램 생성", "프로그램 설정과 평가 기준"),
    ("04", "세션 스케줄링", "교육 일정 계획과 참석자 배정"),
    ("05", "출석 & 결과 입력", "교육 당일 출결 관리와 성적 입력"),
    ("06", "자격증 발급", "합격자 인증서 자동 발급"),
    ("07", "재교육 관리", "불합격·결석·만료자 재교육 추적"),
    ("08", "대시보드 & 분석", "KPI, ROI, 리포트 활용"),
    ("09", "신입 TQC 교육", "신입사원 4단계 교육 과정"),
    ("10", "CAPA 프로세스", "문제 발생 시 시정/예방 조치"),
    ("11", "역량 & 스킬 갭", "역량 평가와 학습 경로 관리"),
]

y = Inches(2.3)
for i, (num, title, desc) in enumerate(chapters):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + col * Inches(6.2)
    cy = y + row * Inches(0.75)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, cy, Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False

    add_text(slide, x + Inches(0.55), cy - Inches(0.02), Inches(5), Inches(0.25),
             title, size=14, color=NAVY, bold=True)
    add_text(slide, x + Inches(0.55), cy + Inches(0.22), Inches(5), Inches(0.2),
             desc, size=10, color=GRAY)


# ════════════════════════════════════════
# SLIDE 3: Characters
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Chapter 1: 등장인물 소개", "시뮬레이션에 등장하는 4명의 캐릭터")

characters = [
    ("문경서 차장", "QIP 관리자", "ADMIN", BLUE,
     "시스템 전체 관리 권한\n프로그램/직원/설정 관리\n리포트 및 감사 접근"),
    ("Nguyen Thi Mai", "QA 팀장 / 사내 강사", "TRAINER", GREEN,
     "교육 세션 생성 및 진행\n출석 체크 및 결과 입력\n자격증 발급 요청"),
    ("Tran Van Duc", "생산 1동 라인장", "VIEWER", ORANGE,
     "대시보드 조회\n팀원 교육 현황 확인\n일정 및 알림 수신"),
    ("Le Thi Hoa", "신입 작업자 (입사 3일차)", "교육 대상자", ACCENT,
     "교육 수강 및 시험 응시\n알림 수신\nTQC 신입 교육 대상"),
]

for i, (name, role, badge, color, desc) in enumerate(characters):
    x = Inches(0.5) + i * Inches(3.15)
    y = Inches(2.4)

    card = add_shape(slide, x, y, Inches(2.95), Inches(4.5), WHITE, radius=0.04)
    card.shadow.inherit = False

    # Color bar
    add_shape(slide, x, y, Inches(2.95), Inches(0.5), color, radius=0.04)
    add_shape(slide, x, y + Inches(0.3), Inches(2.95), Inches(0.2), color)

    # Badge
    badge_shape = add_shape(slide, x + Inches(0.15), y + Inches(0.1), Inches(1.2), Inches(0.3), WHITE, radius=0.1)
    tf = badge_shape.text_frame
    tf.paragraphs[0].text = badge
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Name and role
    add_text(slide, x + Inches(0.15), y + Inches(0.7), Inches(2.65), Inches(0.35),
             name, size=16, color=NAVY, bold=True)
    add_text(slide, x + Inches(0.15), y + Inches(1.05), Inches(2.65), Inches(0.25),
             role, size=11, color=GRAY)

    # Divider
    add_shape(slide, x + Inches(0.15), y + Inches(1.45), Inches(2.65), Inches(0.02), LIGHT_BLUE)

    # Description
    txBox = add_text(slide, x + Inches(0.15), y + Inches(1.6), Inches(2.65), Inches(2.5),
                     "", size=11, color=DARK)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for line in desc.split('\n'):
        add_para(tf, "  " + line, size=11, color=DARK, space_before=Pt(6), space_after=Pt(2))


# ════════════════════════════════════════
# SLIDE 4: Master Data Overview
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Chapter 2: 마스터 데이터 등록", "시스템 운영 전 등록해야 하는 기초 데이터")

# Data registration flow
items = [
    ("1", "직원 (Employees)", "인사시스템 연동 완료", GREEN, "1,200명 자동 동기화\n사번, 이름, 부서, 직급, 동/라인"),
    ("2", "교육 프로그램 (Programs)", "관리자가 직접 등록", BLUE, "프로그램 코드, 3개국어 명칭\n카테고리, 합격기준, 유효기간"),
    ("3", "강사 (Trainers)", "엑셀 템플릿으로 수집", ORANGE, "사내/사외 강사 정보\n전문분야, 자격증"),
    ("4", "역량 (Competencies)", "엑셀 템플릿으로 수집", ORANGE, "스킬 코드, 3개국어 명칭\n카테고리, 핵심역량 여부"),
    ("5", "TQC 팀", "엑셀 템플릿으로 수집", ORANGE, "신입교육 팀 구성\n공장/라인 배정"),
    ("6", "교육비용 (선택)", "엑셀 템플릿으로 수집", GRAY, "프로그램별 월간 비용\nROI 분석용"),
]

for i, (num, title, status, color, desc) in enumerate(items):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + col * Inches(4.1)
    y = Inches(2.4) + row * Inches(2.45)

    card = add_shape(slide, x, y, Inches(3.9), Inches(2.2), WHITE, radius=0.04)
    card.shadow.inherit = False

    # Number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15), Inches(0.35), Inches(0.35))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, x + Inches(0.6), y + Inches(0.12), Inches(3), Inches(0.3),
             title, size=13, color=NAVY, bold=True)
    add_text(slide, x + Inches(0.6), y + Inches(0.38), Inches(3), Inches(0.22),
             status, size=9, color=color, bold=True)

    # Desc
    txBox = add_text(slide, x + Inches(0.15), y + Inches(0.75), Inches(3.6), Inches(1.3),
                     "", size=10, color=DARK)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for line in desc.split('\n'):
        add_para(tf, "  " + line, size=10, color=DARK, space_before=Pt(3))


# ════════════════════════════════════════
# SLIDE 5: Program Creation
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Chapter 3: 교육 프로그램 생성", "/programs 페이지에서 프로그램 등록")

# Example program
fields = [
    ["항목", "설정 값", "설명"],
    ["프로그램 코드", "QIP-001", "고유 식별 코드"],
    ["프로그램명 (EN)", "Visual Inspection Basic", "영문 교육명"],
    ["프로그램명 (VI)", "Kiểm tra trực quan cơ bản", "베트남어 교육명"],
    ["프로그램명 (KR)", "외관검사 기초", "한국어 교육명"],
    ["카테고리", "QIP", "QIP / Production / Retraining / Newcomer / Promotion"],
    ["대상 직급", "Worker, Line Leader", "이 교육이 필요한 직급"],
    ["평가 방식", "SCORE (점수제)", "점수제 또는 합격/불합격"],
    ["합격 점수", "70점", "이 점수 이상이면 PASS"],
    ["등급 기준", "AA≥95 / A≥85 / B≥70", "점수 구간별 자동 등급 부여"],
    ["교육 시간", "4시간", "1회 교육 소요 시간"],
    ["유효기간", "12개월", "자격 만료까지 기간 (재교육 주기)"],
]

add_table_slide(slide, Inches(0.5), Inches(2.3), Inches(12.3), fields,
                [Inches(2.2), Inches(3.5), Inches(6.6)])

# Note
add_text(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
         "* 프로그램 생성/수정/삭제 시 자동으로 program_change_logs에 변경 이력이 기록됩니다.", size=10, color=GRAY)


# ════════════════════════════════════════
# SLIDE 6: Session Scheduling
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Chapter 4: 세션 스케줄링", "/schedule 캘린더에서 교육 일정 생성")

# Left: Session creation
add_shape(slide, Inches(0.5), Inches(2.3), Inches(5.8), Inches(4.8), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(2.45), Inches(5), Inches(0.3),
         "새 세션 생성", size=16, color=NAVY, bold=True)

session_fields = [
    "프로그램: QIP-001 (외관검사 기초)",
    "날짜: 2026-02-20 (목)",
    "시간: 09:00 - 13:00",
    "강사: Nguyen Thi Mai",
    "장소: Training Room A - Building A",
    "최대 인원: 30명",
    "참석 대상: Building A Worker 28명 선택",
    "상태: PLANNED",
]

txBox = add_text(slide, Inches(0.7), Inches(2.9), Inches(5.2), Inches(3.8), "", size=12, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for f in session_fields:
    add_para(tf, "  " + f, size=12, color=DARK, space_before=Pt(8))

# Right: Notification
add_shape(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(2.2), LIGHT_BLUE, radius=0.04)
add_text(slide, Inches(7), Inches(2.45), Inches(5.5), Inches(0.3),
         "자동 알림 발송", size=16, color=NAVY, bold=True)

txBox = add_text(slide, Inches(7), Inches(2.9), Inches(5.5), Inches(1.4), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
add_para(tf, "  세션 생성 시 참석 대상자에게 자동 알림", size=11, color=DARK, space_before=Pt(6))
add_para(tf, "  알림 유형: TRAINING_REMINDER", size=11, color=BLUE, bold=True, space_before=Pt(6))
add_para(tf, "  만료 예정 시: 30/14/7/1일 전 경고", size=11, color=DARK, space_before=Pt(6))
add_para(tf, "  사용자별 알림 설정 커스터마이징 가능", size=11, color=DARK, space_before=Pt(6))

# Right: Status flow
add_shape(slide, Inches(6.8), Inches(4.8), Inches(6), Inches(2.3), WHITE, radius=0.04)
add_text(slide, Inches(7), Inches(4.95), Inches(5.5), Inches(0.3),
         "세션 상태 흐름", size=16, color=NAVY, bold=True)

# Status boxes
statuses = [("PLANNED", BLUE), ("COMPLETED", GREEN), ("CANCELLED", RED)]
for i, (st, col) in enumerate(statuses):
    sx = Inches(7.2) + i * Inches(1.8)
    box = add_shape(slide, sx, Inches(5.5), Inches(1.5), Inches(0.4), col, radius=0.08)
    tf = box.text_frame
    tf.paragraphs[0].text = st
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    if i < 2:
        add_text(slide, sx + Inches(1.5), Inches(5.5), Inches(0.3), Inches(0.4),
                 "→", size=16, color=GRAY, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(7), Inches(6.1), Inches(5.5), Inches(0.5),
         "교육 완료 후 결과 입력 시 자동으로 COMPLETED로 전환\nCANCELLED는 세션 취소 시 (참석자에게 알림 발송)",
         size=10, color=GRAY)


# ════════════════════════════════════════
# SLIDE 7: Attendance & Results
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Chapter 5: 출석 & 결과 입력", "교육 당일: 출석 체크 → 시험 → 결과 기록")

# Attendance table
add_text(slide, Inches(0.5), Inches(2.2), Inches(5), Inches(0.3),
         "출석 기록 (/attendance)", size=14, color=NAVY, bold=True)

att_data = [
    ["사번", "이름", "출석", "체크인", "체크아웃"],
    ["HV-2024-0847", "Le Thi Hoa", "PRESENT", "08:50", "13:05"],
    ["HV-2023-0312", "Pham Van Thanh", "PRESENT", "08:55", "13:00"],
    ["HV-2022-0156", "Vo Thi Lan", "LATE", "09:20", "13:00"],
    ["HV-2024-0901", "Bui Quang Minh", "ABSENT", "-", "-"],
]
add_table_slide(slide, Inches(0.5), Inches(2.55), Inches(6), att_data,
                [Inches(1.4), Inches(1.5), Inches(1.1), Inches(1), Inches(1)])

# Result table
add_text(slide, Inches(6.8), Inches(2.2), Inches(6), Inches(0.3),
         "결과 입력 (/results)", size=14, color=NAVY, bold=True)

res_data = [
    ["이름", "점수", "등급", "결과", "재교육"],
    ["Le Thi Hoa", "92", "A", "PASS", "NO"],
    ["Pham Van Thanh", "78", "B", "PASS", "NO"],
    ["Vo Thi Lan", "55", "C", "FAIL", "YES"],
    ["Bui Quang Minh", "-", "-", "ABSENT", "YES"],
]
add_table_slide(slide, Inches(6.8), Inches(2.55), Inches(6), res_data,
                [Inches(1.5), Inches(0.8), Inches(0.8), Inches(1.1), Inches(0.8)])

# Key rules
add_shape(slide, Inches(0.5), Inches(5), Inches(12.3), Inches(2.2), RGBColor(0xFE, 0xF3, 0xC7), radius=0.04)
add_text(slide, Inches(0.7), Inches(5.1), Inches(12), Inches(0.3),
         "핵심 규칙", size=14, color=RGBColor(0x92, 0x40, 0x0E), bold=True)

rules = [
    "점수 입력 시 등급 자동 계산:  AA (≥95) → A (≥85) → B (≥70) → C (<70)",
    "합격 점수(70) 미만 = FAIL → needs_retraining 자동 활성화",
    "ABSENT(결석)도 자동으로 재교육 대상으로 분류",
    "결과 수정 시 반드시 '수정 사유' 입력 필수 → result_edit_logs에 영구 보존",
    "결과 삭제는 절대 불가 (No Delete Policy) — 감사 추적을 위한 핵심 원칙",
]

txBox = add_text(slide, Inches(0.7), Inches(5.45), Inches(11.5), Inches(1.6), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for r in rules:
    add_para(tf, "  " + r, size=11, color=RGBColor(0x78, 0x35, 0x0F), space_before=Pt(5))


# ════════════════════════════════════════
# SLIDE 8: Certificate & Retraining
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Chapter 6-7: 자격증 발급 & 재교육 관리", "합격자 인증 → 불합격·만료자 재교육 추적")

# Left: Certificate
add_shape(slide, Inches(0.5), Inches(2.3), Inches(6), Inches(4.8), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(2.45), Inches(5), Inches(0.3),
         "자격증 발급 (/certificates)", size=16, color=NAVY, bold=True)

cert_items = [
    "자동 채번: CERT-2026-000147",
    "합격자 개별 / 세션별 일괄 발급 가능",
    "커스텀 템플릿: 로고, 서명, 배경 이미지",
    "PDF 다운로드 및 인쇄",
    "유효기간 자동 계산 (프로그램 설정 기반)",
    "자격증 취소(Revoke): 사유 기록 후 무효화",
    "",
    "자격증 예시:",
    "  직원: Le Thi Hoa",
    "  프로그램: QIP-001 외관검사 기초",
    "  점수: 92점 / A등급",
    "  발급일: 2026-02-20",
    "  만료일: 2027-02-20",
]

txBox = add_text(slide, Inches(0.7), Inches(2.9), Inches(5.5), Inches(3.8), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in cert_items:
    add_para(tf, "  " + item, size=11, color=DARK, space_before=Pt(4))

# Right: Retraining
add_shape(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(4.8), WHITE, radius=0.04)
add_text(slide, Inches(7), Inches(2.45), Inches(5.5), Inches(0.3),
         "재교육 관리 (/retraining)", size=16, color=NAVY, bold=True)

retrain_items = [
    "자동 식별되는 재교육 대상:",
    "",
    "  FAILED — 시험 불합격자",
    "  ABSENT — 교육 결석자",
    "  EXPIRED — 자격 유효기간 만료",
    "  EXPIRING — 30일 내 만료 예정",
    "",
    "재교육 프로세스:",
    "  1. 시스템이 대상자 자동 식별",
    "  2. 관리자가 재교육 세션 배정",
    "  3. 재교육 완료 시 새 결과 기록",
    "  4. 합격 시 새 자격증 발급",
    "  5. 재교육 이력 전체 추적 가능",
]

txBox = add_text(slide, Inches(7), Inches(2.9), Inches(5.5), Inches(3.8), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in retrain_items:
    color = RED if "FAILED" in item or "ABSENT" in item else (ORANGE if "EXPIR" in item else DARK)
    add_para(tf, "  " + item, size=11, color=color, space_before=Pt(4),
             bold=("FAILED" in item or "ABSENT" in item or "EXPIRED" in item or "EXPIRING" in item))


# ════════════════════════════════════════
# SLIDE 9: Dashboard & Analytics
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Chapter 8: 대시보드 & 분석", "실시간 KPI 모니터링과 리포트")

# KPI Cards
add_card(slide, Inches(0.5), Inches(2.3), Inches(2.9), Inches(1.3),
         "전체 직원", "1,200명", "Active employees", LIGHT_BLUE, BLUE, NAVY)
add_card(slide, Inches(3.6), Inches(2.3), Inches(2.9), Inches(1.3),
         "이번달 완료", "3건", "Monthly sessions", LIGHT_BLUE, BLUE, NAVY)
add_card(slide, Inches(6.7), Inches(2.3), Inches(2.9), Inches(1.3),
         "전체 이수율", "67.2%", "Target: 90%", LIGHT_BLUE, BLUE, NAVY)
add_card(slide, Inches(9.8), Inches(2.3), Inches(2.9), Inches(1.3),
         "재교육 대상", "45명", "Needs retraining", RGBColor(0xFE, 0xF2, 0xF2), RED, RED)

# Dashboard features
add_shape(slide, Inches(0.5), Inches(3.9), Inches(6), Inches(3.3), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(4.05), Inches(5.5), Inches(0.3),
         "메인 대시보드 (/dashboard)", size=14, color=NAVY, bold=True)

dash_items = [
    "등급 분포 파이차트 (AA/A/B/C)",
    "월별 교육 추세 (계획 vs 완료) 막대 그래프",
    "동별 재교육/만료 현황 차트",
    "30일 내 만료 예정 자격 목록",
    "다가오는 교육 세션 일정",
    "최근 재교육 대상자 목록",
]
txBox = add_text(slide, Inches(0.7), Inches(4.45), Inches(5.5), Inches(2.5), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in dash_items:
    add_para(tf, "  " + item, size=11, color=DARK, space_before=Pt(5))

# Executive dashboard
add_shape(slide, Inches(6.8), Inches(3.9), Inches(6), Inches(3.3), WHITE, radius=0.04)
add_text(slide, Inches(7), Inches(4.05), Inches(5.5), Inches(0.3),
         "경영진 대시보드 (/executive)", size=14, color=NAVY, bold=True)

exec_data = [
    ["KPI", "현재", "목표", "달성률"],
    ["교육 이수율", "67.2%", "90%", "74.7%"],
    ["자격 보유율", "58.3%", "85%", "68.6%"],
    ["신입 이직률", "8.2%", "5%", "-"],
    ["교육 ROI", "187%", "200%", "93.5%"],
]
add_table_slide(slide, Inches(7), Inches(4.5), Inches(5.5), exec_data,
                [Inches(1.5), Inches(1), Inches(1), Inches(1)])

add_text(slide, Inches(7), Inches(6.7), Inches(5.5), Inches(0.4),
         "리포트 (/reports): 부서별, 프로그램별, 직원별\n기간 필터 (1M/3M/6M/1Y) | CSV/PDF 내보내기",
         size=10, color=GRAY)


# ════════════════════════════════════════
# SLIDE 10: New TQC Training
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Chapter 9: 신입 TQC 교육", "신입사원 4단계 교육 과정 관리 (/new-tqc)")

# 4 Stages
stages = [
    ("Stage 1", "오리엔테이션\nOrientation", "공장 안전 교육\n사내 규정 안내\n기본 품질 개념", GREEN),
    ("Stage 2", "기초 교육\nBasic Training", "QIP 교육 수강\n시험 응시\n합격 기준 충족", BLUE),
    ("Stage 3", "라인 배치\nLine Assignment", "실제 라인에서 OJT\n선임자 지도하에\n실무 능력 습득", ORANGE),
    ("Stage 4", "현장 평가\nField Evaluation", "독립 작업 능력 평가\n최종 합격 판정\n자격증 발급", ACCENT),
]

for i, (label, title, desc, color) in enumerate(stages):
    x = Inches(0.5) + i * Inches(3.15)
    y = Inches(2.3)

    card = add_shape(slide, x, y, Inches(2.95), Inches(3), WHITE, radius=0.04)
    # Top bar
    add_shape(slide, x, y, Inches(2.95), Inches(0.7), color, radius=0.04)
    add_shape(slide, x, y + Inches(0.5), Inches(2.95), Inches(0.2), color)

    add_text(slide, x + Inches(0.1), y + Inches(0.05), Inches(2.75), Inches(0.25),
             label, size=10, color=RGBColor(0xBF, 0xDB, 0xFE), bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(0.25), Inches(2.75), Inches(0.45),
             title, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    txBox = add_text(slide, x + Inches(0.15), y + Inches(0.85), Inches(2.65), Inches(2), "", size=11, color=DARK)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for line in desc.split('\n'):
        add_para(tf, "  " + line, size=11, color=DARK, space_before=Pt(6))

    # Arrow between stages
    if i < 3:
        add_text(slide, x + Inches(2.95), y + Inches(1.2), Inches(0.2), Inches(0.3),
                 "→", size=18, color=GRAY, bold=True)

# Additional TQC features
add_shape(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.7), LIGHT_BLUE, radius=0.04)
add_text(slide, Inches(0.7), Inches(5.6), Inches(11.5), Inches(0.3),
         "TQC 추가 기능", size=14, color=NAVY, bold=True)

tqc_features = [
    "정기 면담: 1주 / 1개월 / 3개월 면담 자동 스케줄링 (/new-tqc/meetings)",
    "색맹 검사: 입사 시 PASS/FAIL 기록 (품질 검사 직무 배정 기준)",
    "퇴사 추적: 퇴사 사유 8가지 카테고리 분류, 교육 기간 중 이직률 분석 (/new-tqc/resignations)",
    "TQC 대시보드: 신입 인원, 완료율, 이직률, 평균 온보딩 기간 종합 모니터링",
]

txBox = add_text(slide, Inches(0.7), Inches(5.95), Inches(11.5), Inches(1.1), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for f in tqc_features:
    add_para(tf, "  " + f, size=10, color=DARK, space_before=Pt(4))


# ════════════════════════════════════════
# SLIDE 11: CAPA Process
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Chapter 10: CAPA 프로세스", "문제 발생 시 5단계 시정/예방 조치 (/capa)")

# 5 Stage flow
capa_stages = [
    ("1", "발견\nDiscovery", "문제 기술\n출처 분류\n즉시 조치", RGBColor(0xEF, 0x44, 0x44)),
    ("2", "조사\nInvestigation", "근본 원인 분석\n5-Why / Fishbone\n영향 평가", RGBColor(0xF5, 0x9E, 0x0B)),
    ("3", "조치\nAction", "시정 조치 계획\n예방 조치 계획\n담당자/기한 배정", RGBColor(0x25, 0x63, 0xEB)),
    ("4", "검증\nVerification", "효과성 검증\n모니터링 기간\n재발 여부 확인", RGBColor(0x7C, 0x3A, 0xED)),
    ("5", "종료\nClosure", "교훈 기록\n지식 공유\n문서화 완료", RGBColor(0x05, 0x96, 0x69)),
]

for i, (num, title, desc, color) in enumerate(capa_stages):
    x = Inches(0.3) + i * Inches(2.55)
    y = Inches(2.3)

    card = add_shape(slide, x, y, Inches(2.35), Inches(2.8), WHITE, radius=0.04)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85), y + Inches(0.15), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, x + Inches(0.1), y + Inches(0.8), Inches(2.15), Inches(0.5),
             title, size=11, color=color, bold=True, align=PP_ALIGN.CENTER)

    txBox = add_text(slide, x + Inches(0.1), y + Inches(1.4), Inches(2.15), Inches(1.2), "", size=10, color=DARK)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for line in desc.split('\n'):
        add_para(tf, line, size=10, color=DARK, space_before=Pt(4), align=PP_ALIGN.CENTER)

    # Arrow
    if i < 4:
        add_text(slide, x + Inches(2.35), y + Inches(1.1), Inches(0.2), Inches(0.3),
                 "→", size=16, color=GRAY, bold=True)

# CAPA Example
add_shape(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.9), RGBColor(0xFE, 0xF2, 0xF2), radius=0.04)
add_text(slide, Inches(0.7), Inches(5.4), Inches(11.5), Inches(0.3),
         "시뮬레이션 예시: CAPA-2026-003", size=14, color=RED, bold=True)

txBox = add_text(slide, Inches(0.7), Inches(5.75), Inches(11.5), Inches(1.3), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
add_para(tf, "  문제: Building A L3 외관검사 불량률 3.2% (기준 1.5% 초과)", size=10, color=DARK, space_before=Pt(3))
add_para(tf, "  원인: QIP-001 교육이 이론 위주 → 실무 적용 능력 부족", size=10, color=DARK, space_before=Pt(3))
add_para(tf, "  조치: 보충교육 실시 + 미세 스크래치 판별 가이드 제작 + 커리큘럼에 실습 2시간 추가", size=10, color=DARK, space_before=Pt(3))
add_para(tf, "  결과: 불량률 3.2% → 1.1% (기준 이내) | 효과성 85/100 | CAPA → 프로그램 개선으로 연결", size=10, color=GREEN, bold=True, space_before=Pt(3))


# ════════════════════════════════════════
# SLIDE 12: Competency & Skill Gap
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Chapter 11: 역량 & 스킬 갭 분석", "역량 평가 → 갭 분석 → 학습 경로 추천")

# Competency assessment
add_text(slide, Inches(0.5), Inches(2.2), Inches(6), Inches(0.3),
         "직원별 역량 평가 (/competency)", size=14, color=NAVY, bold=True)

comp_data = [
    ["역량", "현재 수준", "목표 수준", "갭"],
    ["외관 검사", "COMPETENT (3)", "PROFICIENT (4)", "-1"],
    ["봉제 품질", "BEGINNER (2)", "COMPETENT (3)", "-1"],
    ["안전 절차", "COMPETENT (3)", "COMPETENT (3)", "0"],
    ["팀 리더십", "NOVICE (1)", "NOVICE (1)", "0"],
]
add_table_slide(slide, Inches(0.5), Inches(2.55), Inches(6), comp_data,
                [Inches(1.5), Inches(1.5), Inches(1.5), Inches(0.8)])

# Skill gap
add_text(slide, Inches(6.8), Inches(2.2), Inches(6), Inches(0.3),
         "팀/부서 스킬 갭 분석 (/skill-gap)", size=14, color=NAVY, bold=True)

gap_data = [
    ["역량", "팀 평균", "목표", "갭", "미달 인원"],
    ["외관 검사", "2.8", "4.0", "-1.2", "15명"],
    ["봉제 품질", "2.5", "3.0", "-0.5", "8명"],
    ["안전 절차", "3.2", "3.0", "+0.2", "0명"],
]
add_table_slide(slide, Inches(6.8), Inches(2.55), Inches(6), gap_data,
                [Inches(1.3), Inches(1), Inches(1), Inches(0.8), Inches(1)])

# Learning path
add_shape(slide, Inches(0.5), Inches(5), Inches(12.3), Inches(2.2), LIGHT_BLUE, radius=0.04)
add_text(slide, Inches(0.7), Inches(5.1), Inches(11.5), Inches(0.3),
         "학습 경로 & 개발 계획", size=14, color=NAVY, bold=True)

lp_items = [
    "5가지 역량 수준: NOVICE (1) → BEGINNER (2) → COMPETENT (3) → PROFICIENT (4) → EXPERT (5)",
    "6가지 역량 카테고리: TECHNICAL / QUALITY / SAFETY / LEADERSHIP / COMMUNICATION / PROCESS",
    "학습 경로 유형: ONBOARDING (입사) / POSITION (직무) / PROMOTION (승진) / SPECIALIZATION (전문) / REMEDIAL (보충)",
    "개인별 개발 계획(IDP): 목표 역량, 할당된 학습 경로, 관리자/본인 코멘트, 진행률 추적",
    "역량-프로그램 매핑: 각 역량에 필요한 교육 프로그램 자동 추천",
]

txBox = add_text(slide, Inches(0.7), Inches(5.5), Inches(11.5), Inches(1.5), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
for item in lp_items:
    add_para(tf, "  " + item, size=10, color=DARK, space_before=Pt(5))


# ════════════════════════════════════════
# SLIDE 13: Complete Workflow
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "전체 워크플로우 요약", "Q-TRAIN 시스템 운영 흐름 한눈에 보기")

# Flow steps
flow_steps = [
    ("1", "기초 데이터\n등록", "직원/강사/역량\nTQC팀/비용", NAVY),
    ("2", "프로그램\n생성", "교육과정 설정\n평가기준 정의", BLUE),
    ("3", "세션\n스케줄링", "일정/강사/장소\n참석자 배정", RGBColor(0x06, 0x69, 0xB4)),
    ("4", "교육 실시\n& 출석", "출석 체크\n교육 진행", RGBColor(0x05, 0x96, 0x69)),
    ("5", "결과\n입력", "점수/등급\n합격 판정", RGBColor(0x7C, 0x3A, 0xED)),
    ("6", "자격증\n발급", "합격자 인증\nPDF 발급", RGBColor(0xDB, 0x27, 0x77)),
    ("7", "재교육\n관리", "불합격/만료\n추적 & 배정", ORANGE),
    ("8", "대시보드\n& 분석", "KPI/ROI\n리포트", RGBColor(0x05, 0x96, 0x69)),
]

y = Inches(2.4)
for i, (num, title, desc, color) in enumerate(flow_steps):
    x = Inches(0.3) + i * Inches(1.6)

    # Circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.35), y, Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, x, y + Inches(0.9), Inches(1.5), Inches(0.5),
             title, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(1.45), Inches(1.5), Inches(0.45),
             desc, size=8, color=GRAY, align=PP_ALIGN.CENTER)

    if i < 7:
        add_text(slide, x + Inches(1.25), y + Inches(0.2), Inches(0.35), Inches(0.4),
                 "→", size=18, color=LIGHT_GRAY, bold=True, align=PP_ALIGN.CENTER)

# CAPA & TQC side flows
add_shape(slide, Inches(0.5), Inches(4.6), Inches(6), Inches(2.5), WHITE, radius=0.04)
add_text(slide, Inches(0.7), Inches(4.7), Inches(5.5), Inches(0.3),
         "문제 발생 시 → CAPA 프로세스", size=13, color=RED, bold=True)
txBox = add_text(slide, Inches(0.7), Inches(5.05), Inches(5.5), Inches(1.8), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
add_para(tf, "  발견 → 조사(5-Why) → 시정/예방 조치 → 검증 → 종료", size=10, color=DARK, space_before=Pt(4))
add_para(tf, "  CAPA 결과가 프로그램 개선으로 피드백", size=10, color=DARK, space_before=Pt(4))
add_para(tf, "  교훈(Lessons Learned) 전사 공유", size=10, color=DARK, space_before=Pt(4))
add_para(tf, "  Adidas 감사 대응: 6개 카테고리 컴플라이언스 추적", size=10, color=DARK, space_before=Pt(4))

add_shape(slide, Inches(6.8), Inches(4.6), Inches(6), Inches(2.5), WHITE, radius=0.04)
add_text(slide, Inches(7), Inches(4.7), Inches(5.5), Inches(0.3),
         "신입사원 → TQC 교육 과정", size=13, color=ACCENT, bold=True)
txBox = add_text(slide, Inches(7), Inches(5.05), Inches(5.5), Inches(1.8), "", size=10, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
add_para(tf, "  4단계: 오리엔테이션 → 기초교육 → 라인배치 → 현장평가", size=10, color=DARK, space_before=Pt(4))
add_para(tf, "  정기 면담: 1주 / 1개월 / 3개월", size=10, color=DARK, space_before=Pt(4))
add_para(tf, "  색맹검사 / 퇴사 추적 / 온보딩 기간 분석", size=10, color=DARK, space_before=Pt(4))
add_para(tf, "  TQC 대시보드에서 종합 모니터링", size=10, color=DARK, space_before=Pt(4))


# ════════════════════════════════════════
# SLIDE 14: Data Registration Summary
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "데이터 등록 순서 요약", "시스템 운영을 위해 필요한 데이터와 현재 상태")

reg_data = [
    ["순서", "데이터", "상태", "등록 방법", "설명"],
    ["①", "직원 (Employees)", "완료", "인사시스템 연동", "1,200명 자동 동기화"],
    ["②", "교육 프로그램 (Programs)", "완료", "관리자 직접 등록", "프로그램 코드, 평가기준, 유효기간"],
    ["③", "강사 (Trainers)", "요청중", "엑셀 → 시스템 입력", "사내/사외 강사, 전문분야, 자격"],
    ["④", "역량 (Competencies)", "요청중", "엑셀 → 시스템 입력", "스킬 코드, 3개국어명, 카테고리"],
    ["⑤", "TQC 팀 (Teams)", "요청중", "엑셀 → 시스템 입력", "신입교육 팀, 공장/라인 배정"],
    ["⑥", "교육비용 (Costs)", "요청중", "엑셀 → 시스템 입력", "프로그램별 월간 비용 (ROI용)"],
]

add_table_slide(slide, Inches(0.5), Inches(2.3), Inches(12.3), reg_data,
                [Inches(0.7), Inches(2.5), Inches(1), Inches(2), Inches(6.1)])

# After registration
add_shape(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.8), RGBColor(0xEC, 0xFD, 0xF5), radius=0.04)
add_text(slide, Inches(0.7), Inches(5.5), Inches(11.5), Inches(0.3),
         "③~⑥ 데이터 등록 후 운영 가능한 기능:", size=13, color=GREEN, bold=True)

txBox = add_text(slide, Inches(0.7), Inches(5.85), Inches(11.5), Inches(1.2), "", size=11, color=DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = ""
add_para(tf, "  세션 생성 → 출석 관리 → 결과 입력 → 자격증 발급 → 재교육 관리 → CAPA → 대시보드/리포트/역량분석", size=11, color=DARK, space_before=Pt(4))
add_para(tf, "  현재 ①②는 완료, ③~⑥은 담당자에게 엑셀 템플릿으로 데이터를 요청한 상태입니다.", size=11, color=GRAY, space_before=Pt(6))


# ════════════════════════════════════════
# SLIDE 15: Thank You
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), BLUE)

add_text(slide, Inches(1), Inches(2), Inches(11), Inches(0.8),
         "Q-TRAIN", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.5),
         "QIP Training Management System", size=22, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.04), BLUE)

add_text(slide, Inches(1), Inches(3.9), Inches(11), Inches(0.4),
         "Thank You", size=28, color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.3),
         "https://q-train-web.web.app", size=14, color=RGBColor(0x60, 0xA5, 0xFA), align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(5.4), Inches(11), Inches(0.5),
         "HWK Vietnam (화승비나) | QIP Training Management\n문의: ksmoon@hsvina.com",
         size=12, color=RGBColor(0x64, 0x78, 0x96), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# Save
# ════════════════════════════════════════
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "Q-TRAIN_System_Guide.pptx")
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
